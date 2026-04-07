"""
Tests for PPTX diff engine.

Validates slide-level comparison, text/shape detection, and image hashing.
"""
import pytest
from pathlib import Path
from io import BytesIO

# Skip all tests if python-pptx not available
pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches, Pt

from ikam.diff.pptx_diff import (
    diff_pptx,
    diff_pptx_files,
    _extract_slide_text,
    _extract_slide_fingerprint,
    _compare_slides,
)


# Fixtures for creating test presentations

def _create_presentation_bytes(slide_texts: list[str]) -> bytes:
    """Create a simple PPTX with text slides."""
    prs = Presentation()
    
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = text
    
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _create_presentation_with_shapes(shapes_per_slide: list[int]) -> bytes:
    """Create PPTX with varying shape counts per slide."""
    prs = Presentation()
    
    for shape_count in shapes_per_slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        for i in range(shape_count):
            left = Inches(1 + i * 0.5)
            top = Inches(1)
            width = Inches(1)
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = f"Shape {i + 1}"
    
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# Test: Identical presentations

def test_diff_pptx_identical():
    """Identical presentations should return no changes."""
    pptx_bytes = _create_presentation_bytes(["Slide 1", "Slide 2", "Slide 3"])
    
    changes = diff_pptx(pptx_bytes, pptx_bytes)
    
    assert len(changes) == 0, "Identical presentations should have no differences"


# Test: Slide additions

def test_diff_pptx_slides_added():
    """Detect slides added at the end."""
    old_pptx = _create_presentation_bytes(["Slide 1", "Slide 2"])
    new_pptx = _create_presentation_bytes(["Slide 1", "Slide 2", "Slide 3", "Slide 4"])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 2, "Should detect 2 added slides"
    
    # Check added slides
    added_changes = [c for c in changes if c.change_type == "added"]
    assert len(added_changes) == 2
    assert added_changes[0].path == "Slide 3"
    assert added_changes[1].path == "Slide 4"


# Test: Slide removals

def test_diff_pptx_slides_removed():
    """Detect slides removed from the end."""
    old_pptx = _create_presentation_bytes(["Slide 1", "Slide 2", "Slide 3", "Slide 4"])
    new_pptx = _create_presentation_bytes(["Slide 1", "Slide 2"])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 2, "Should detect 2 removed slides"
    
    # Check removed slides
    removed_changes = [c for c in changes if c.change_type == "removed"]
    assert len(removed_changes) == 2
    assert removed_changes[0].path == "Slide 3"
    assert removed_changes[1].path == "Slide 4"


# Test: Text modifications

def test_diff_pptx_text_modified():
    """Detect text changes within slides."""
    old_pptx = _create_presentation_bytes(["Original Title", "Slide 2"])
    new_pptx = _create_presentation_bytes(["Modified Title", "Slide 2"])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 1, "Should detect 1 modified slide"
    
    change = changes[0]
    assert change.path == "Slide 1"
    assert change.change_type == "modified"
    assert "text_modified" in change.new_value["changes"]


# Test: Multiple text modifications

def test_diff_pptx_multiple_slides_modified():
    """Detect modifications in multiple slides."""
    old_pptx = _create_presentation_bytes(["Title A", "Title B", "Title C"])
    new_pptx = _create_presentation_bytes(["Title A", "Modified B", "Modified C"])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 2, "Should detect 2 modified slides"
    
    modified_changes = [c for c in changes if c.change_type == "modified"]
    assert len(modified_changes) == 2
    assert modified_changes[0].path == "Slide 2"
    assert modified_changes[1].path == "Slide 3"


# Test: Shape count changes

def test_diff_pptx_shapes_added():
    """Detect shapes added to slides."""
    old_pptx = _create_presentation_with_shapes([2, 3])
    new_pptx = _create_presentation_with_shapes([2, 5])  # Added 2 shapes to slide 2
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 1, "Should detect 1 modified slide"
    
    change = changes[0]
    assert change.path == "Slide 2"
    assert any("shapes_added" in c for c in change.new_value["changes"])


def test_diff_pptx_shapes_removed():
    """Detect shapes removed from slides."""
    old_pptx = _create_presentation_with_shapes([4, 3])
    new_pptx = _create_presentation_with_shapes([2, 3])  # Removed 2 shapes from slide 1
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 1, "Should detect 1 modified slide"
    
    change = changes[0]
    assert change.path == "Slide 1"
    assert any("shapes_removed" in c for c in change.new_value["changes"])


# Test: Empty presentations

def test_diff_pptx_empty_to_content():
    """Detect adding slides to empty presentation."""
    old_pptx = _create_presentation_bytes([])
    new_pptx = _create_presentation_bytes(["Slide 1"])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 1
    assert changes[0].change_type == "added"


def test_diff_pptx_content_to_empty():
    """Detect removing all slides."""
    old_pptx = _create_presentation_bytes(["Slide 1", "Slide 2"])
    new_pptx = _create_presentation_bytes([])
    
    changes = diff_pptx(old_pptx, new_pptx)
    
    assert len(changes) == 2
    assert all(c.change_type == "removed" for c in changes)


# Test: Slide text extraction

def test_extract_slide_text():
    """Verify text extraction from slide shapes."""
    pptx_bytes = _create_presentation_bytes(["Test Title"])
    prs = Presentation(BytesIO(pptx_bytes))
    slide = prs.slides[0]
    
    text = _extract_slide_text(slide)
    
    assert "Test Title" in text


# Test: Slide fingerprint

def test_extract_slide_fingerprint():
    """Verify fingerprint extraction captures key slide properties."""
    pptx_bytes = _create_presentation_bytes(["Fingerprint Test"])
    prs = Presentation(BytesIO(pptx_bytes))
    slide = prs.slides[0]
    
    fingerprint = _extract_slide_fingerprint(slide, 0)
    
    assert "index" in fingerprint
    assert "text" in fingerprint
    assert "shape_count" in fingerprint
    assert "image_hashes" in fingerprint
    assert fingerprint["index"] == 0
    assert "Fingerprint Test" in fingerprint["text"]


# Test: Compare slides helper

def test_compare_slides_identical():
    """Identical slide fingerprints should return no changes."""
    fp = {
        "index": 0,
        "text": "Same text",
        "shape_count": 2,
        "image_hashes": ["abc123"],
    }
    
    changes = _compare_slides(fp, fp)
    
    assert len(changes) == 0


def test_compare_slides_text_modified():
    """Detect text changes between slides."""
    old_fp = {"index": 0, "text": "Old", "shape_count": 1, "image_hashes": []}
    new_fp = {"index": 0, "text": "New", "shape_count": 1, "image_hashes": []}
    
    changes = _compare_slides(old_fp, new_fp)
    
    assert "text_modified" in changes


# Test: Error handling

def test_diff_pptx_invalid_bytes():
    """Invalid PPTX bytes should raise ValueError."""
    with pytest.raises(ValueError, match="Failed to parse"):
        diff_pptx(b"invalid", b"also invalid")


# Test: File convenience function

def test_diff_pptx_files(tmp_path):
    """Test file-based diff convenience function."""
    old_pptx = _create_presentation_bytes(["Slide 1"])
    new_pptx = _create_presentation_bytes(["Slide 1", "Slide 2"])
    
    old_path = tmp_path / "old.pptx"
    new_path = tmp_path / "new.pptx"
    
    old_path.write_bytes(old_pptx)
    new_path.write_bytes(new_pptx)
    
    changes = diff_pptx_files(old_path, new_path)
    
    assert len(changes) == 1
    assert changes[0].change_type == "added"
