"""Performance benchmarks for advanced diff features (PPTX, arrays, formulas)."""

import time
from io import BytesIO

import pytest
from openpyxl import Workbook
from pptx import Presentation

from ikam.diff import diff_pptx, diff_arrays, compute_xlsx_diff


# ============================================================================
# PPTX Performance Benchmarks
# ============================================================================

def create_presentation_with_slides(num_slides: int) -> bytes:
    """Create a PPTX presentation with N slides for benchmarking."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    
    for i in range(num_slides):
        slide = prs.slides.add_slide(blank_slide_layout)
        # Add title text box
        left = top = width = height = 1000000  # 1 inch in EMUs
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Slide {i+1} Title"
        
        # Add content text box
        left = top + height + 100000
        contentBox = slide.shapes.add_textbox(left, top, width, height)
        cf = contentBox.text_frame
        cf.text = f"This is content for slide {i+1}."
    
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_pptx_diff_5_slides_performance():
    """PPTX diff: 5 slides should complete in <200ms."""
    old_pptx = create_presentation_with_slides(5)
    new_pptx = create_presentation_with_slides(5)
    
    start = time.perf_counter()
    result = diff_pptx(old_pptx, new_pptx)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 200, f"5-slide PPTX diff took {duration:.2f}ms (target: <200ms)"
    assert isinstance(result, list)


def test_pptx_diff_20_slides_performance():
    """PPTX diff: 20 slides should complete in <500ms."""
    old_pptx = create_presentation_with_slides(20)
    new_pptx = create_presentation_with_slides(20)
    
    start = time.perf_counter()
    result = diff_pptx(old_pptx, new_pptx)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 500, f"20-slide PPTX diff took {duration:.2f}ms (target: <500ms)"
    assert isinstance(result, list)


def test_pptx_diff_50_slides_performance():
    """PPTX diff: 50 slides should complete in <1500ms."""
    old_pptx = create_presentation_with_slides(50)
    new_pptx = create_presentation_with_slides(50)
    
    start = time.perf_counter()
    result = diff_pptx(old_pptx, new_pptx)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 1500, f"50-slide PPTX diff took {duration:.2f}ms (target: <1500ms)"
    assert isinstance(result, list)


# ============================================================================
# Array Diff Performance Benchmarks
# ============================================================================

def test_array_diff_100_elements_performance():
    """Array diff: 100 elements should complete in <20ms."""
    old_array = list(range(100))
    new_array = list(range(100))
    
    start = time.perf_counter()
    result = diff_arrays(old_array, new_array)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 20, f"100-element array diff took {duration:.2f}ms (target: <20ms)"
    assert isinstance(result, list)


def test_array_diff_500_elements_performance():
    """Array diff: 500 elements should complete in <100ms."""
    old_array = list(range(500))
    new_array = list(range(500))
    
    start = time.perf_counter()
    result = diff_arrays(old_array, new_array)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 100, f"500-element array diff took {duration:.2f}ms (target: <100ms)"
    assert isinstance(result, list)


def test_array_diff_2000_elements_performance():
    """Array diff: 2000 elements should complete in <600ms."""
    old_array = list(range(2000))
    new_array = list(range(2000))
    
    start = time.perf_counter()
    result = diff_arrays(old_array, new_array)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 600, f"2000-element array diff took {duration:.2f}ms (target: <600ms)"
    assert isinstance(result, list)


# ============================================================================
# XLSX Formula Performance Benchmarks
# ============================================================================

def create_workbook_with_formulas(num_cells: int) -> bytes:
    """Create a workbook with N cells containing formulas."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    
    # Create formulas in a grid (10 columns wide)
    cols_per_row = 10
    for i in range(num_cells):
        row = (i // cols_per_row) + 1
        col = (i % cols_per_row) + 1
        
        # Create a simple SUM formula
        ws.cell(row=row, column=col).value = f"=SUM(A{row+10}:A{row+20})"
    
    buffer = BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def test_xlsx_formula_diff_100_cells_performance():
    """XLSX formula diff: 100 cells should complete in <200ms."""
    old_xlsx = create_workbook_with_formulas(100)
    new_xlsx = create_workbook_with_formulas(100)
    
    start = time.perf_counter()
    result = compute_xlsx_diff(old_xlsx, new_xlsx, include_formulas=True)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 200, f"100-cell formula diff took {duration:.2f}ms (target: <200ms)"
    assert hasattr(result, 'changes')


def test_xlsx_formula_diff_500_cells_performance():
    """XLSX formula diff: 500 cells should complete in <800ms."""
    old_xlsx = create_workbook_with_formulas(500)
    new_xlsx = create_workbook_with_formulas(500)
    
    start = time.perf_counter()
    result = compute_xlsx_diff(old_xlsx, new_xlsx, include_formulas=True)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 800, f"500-cell formula diff took {duration:.2f}ms (target: <800ms)"
    assert hasattr(result, 'changes')


# ============================================================================
# Combined Realistic Scenarios
# ============================================================================

def test_pptx_diff_with_modifications_performance():
    """PPTX diff with actual changes: 10 slides with 5 modified."""
    old_pptx = create_presentation_with_slides(10)
    
    # Create modified version
    prs = Presentation(BytesIO(old_pptx))
    # Modify some slides
    for i in [2, 4, 6, 8]:
        if i < len(prs.slides):
            slide = prs.slides[i]
            if slide.shapes:
                slide.shapes[0].text_frame.text = f"Modified slide {i+1}"
    
    buffer = BytesIO()
    prs.save(buffer)
    new_pptx = buffer.getvalue()
    
    start = time.perf_counter()
    result = diff_pptx(old_pptx, new_pptx)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 300, f"10-slide PPTX diff with changes took {duration:.2f}ms (target: <300ms)"
    # Should detect modifications
    assert any(c.change_type == "modified" for c in result)


def test_array_diff_with_moves_performance():
    """Array diff with move detection: 200 elements with 50 moved."""
    old_array = list(range(200))
    new_array = old_array.copy()
    
    # Move some elements
    for i in range(50):
        old_idx = i * 3
        new_idx = (i * 3 + 50) % 200
        if old_idx < len(new_array) and new_idx < len(new_array):
            new_array[old_idx], new_array[new_idx] = new_array[new_idx], new_array[old_idx]
    
    start = time.perf_counter()
    result = diff_arrays(old_array, new_array, detect_moves=True)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 150, f"200-element array diff with moves took {duration:.2f}ms (target: <150ms)"
    # Should detect some moves
    assert any(op.operation == "move" for op in result)


def test_xlsx_formula_modification_performance():
    """XLSX formula diff with actual changes: 100 cells with 20 modified."""
    wb1 = Workbook()
    ws1 = wb1.active
    ws1.title = "Sheet1"
    
    for i in range(100):
        row = i + 1
        ws1.cell(row=row, column=1).value = f"=SUM(B1:B{i+5})"
    
    buffer1 = BytesIO()
    wb1.save(buffer1)
    old_xlsx = buffer1.getvalue()
    
    # Create modified version
    wb2 = Workbook()
    ws2 = wb2.active
    ws2.title = "Sheet1"
    
    for i in range(100):
        row = i + 1
        # Change formulas for rows divisible by 5
        if i % 5 == 0:
            ws2.cell(row=row, column=1).value = f"=AVERAGE(B1:B{i+5})"
        else:
            ws2.cell(row=row, column=1).value = f"=SUM(B1:B{i+5})"
    
    buffer2 = BytesIO()
    wb2.save(buffer2)
    new_xlsx = buffer2.getvalue()
    
    start = time.perf_counter()
    result = compute_xlsx_diff(old_xlsx, new_xlsx, include_formulas=True)
    duration = (time.perf_counter() - start) * 1000
    
    assert duration < 250, f"100-cell formula diff with changes took {duration:.2f}ms (target: <250ms)"
    assert result.change_count > 0


# ============================================================================
# Performance Summary Test
# ============================================================================

@pytest.mark.parametrize("scenario,func,target_ms", [
    ("PPTX 5 slides", lambda: diff_pptx(create_presentation_with_slides(5), create_presentation_with_slides(5)), 200),
    ("PPTX 20 slides", lambda: diff_pptx(create_presentation_with_slides(20), create_presentation_with_slides(20)), 500),
    ("Array 100 elements", lambda: diff_arrays(list(range(100)), list(range(100))), 20),
    ("Array 500 elements", lambda: diff_arrays(list(range(500)), list(range(500))), 100),
    ("XLSX 100 formulas", lambda: compute_xlsx_diff(create_workbook_with_formulas(100), create_workbook_with_formulas(100), include_formulas=True), 200),
])
def test_performance_summary(scenario: str, func, target_ms: float):
    """Run performance scenario and report actual vs target time."""
    start = time.perf_counter()
    func()
    duration = (time.perf_counter() - start) * 1000
    
    # Always pass but report timing
    print(f"\n{scenario}: {duration:.2f}ms (target: {target_ms}ms) - {'✓ PASS' if duration < target_ms else '✗ SLOW'}")
    
    # Assert for CI
    assert duration < target_ms * 2, f"{scenario} took {duration:.2f}ms (2x target exceeded)"
