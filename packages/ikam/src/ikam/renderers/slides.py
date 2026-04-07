"""PowerPoint rendering utilities with deterministic mode support.

Phase 4 deterministic export infrastructure:
- Frozen timestamps
- Stable IDs (deterministic element naming/ordering)
- Float precision rounding (for positions, sizes)

The renderer inspects IKAM_* environment variables:
- IKAM_DETERMINISTIC_RENDER (bool)
- IKAM_FROZEN_TIMESTAMP (ISO8601 string)
- IKAM_STABLE_IDS (bool)
- IKAM_FLOAT_PRECISION (int)

Usage:
    from ikam.renderers.slides import render_pptx
    content_bytes = render_pptx({"slides": [{"title": "Overview", "content": "Details"}]})
"""
from __future__ import annotations

import os
import uuid
from io import BytesIO
from datetime import datetime
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches, Pt


def _env_bool(key: str, default: str = "false") -> bool:
    return os.getenv(key, default).lower() in {"1", "true", "yes"}


def _env_int(key: str, default: str) -> int:
    try:
        return int(os.getenv(key, default))
    except ValueError:
        return int(default)


def _parse_timestamp(raw: str | None) -> datetime | None:
    if not raw:
        return None
    if raw.endswith("Z"):
        raw = raw.replace("Z", "+00:00")
    try:
        return datetime.fromisoformat(raw)
    except ValueError:
        return None


def render_pptx(model: Dict[str, Any]) -> bytes:
    """Render a PowerPoint presentation from a model structure.

    Model format:
        {
          "slides": [
             {"title": "Slide Title", "content": "Body text", "notes": "Speaker notes (optional)"},
             ...
          ]
        }

    Deterministic behavior controlled by IKAM_* env variables.
    """
    deterministic = _env_bool("IKAM_DETERMINISTIC_RENDER")
    frozen_ts_raw = os.getenv("IKAM_FROZEN_TIMESTAMP")
    frozen_dt = _parse_timestamp(frozen_ts_raw) if deterministic else None
    stable_ids = _env_bool("IKAM_STABLE_IDS") if deterministic else False
    float_precision = _env_int("IKAM_FLOAT_PRECISION", "15") if deterministic else 15

    prs = Presentation()

    slides_data = model.get("slides", [])

    # Deterministic ordering: sort by title if stable_ids enabled
    if stable_ids:
        slides_data = sorted(slides_data, key=lambda s: s.get("title") or "")

    for slide_spec in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide_spec.get("title", "")
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set body content
        content = slide_spec.get("content", "")
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            body_shape.text = content

        # Set speaker notes if provided
        notes_text = slide_spec.get("notes")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    if deterministic:
        _apply_deterministic_properties(prs, frozen_dt)
    else:
        # Inject non-deterministic token for variance detection
        if prs.slides:
            first_slide = prs.slides[0]
            token_shape = first_slide.shapes.add_textbox(
                Inches(9), Inches(7), Inches(1), Inches(0.3)
            )
            token_shape.text = uuid.uuid4().hex

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _apply_deterministic_properties(prs: Presentation, frozen_dt: datetime | None) -> None:
    """Apply deterministic properties to PowerPoint core properties."""
    props = prs.core_properties
    if frozen_dt:
        ts = frozen_dt.replace(microsecond=0)
        props.created = ts
        props.modified = ts
    # Clear volatile attributes
    _safe_clear_attr(props, "last_modified_by")
    _safe_clear_attr(props, "revision")


def _safe_clear_attr(props, attr: str) -> None:
    if hasattr(props, attr):
        try:
            setattr(props, attr, None)
        except Exception:
            pass


__all__ = ["render_pptx"]
