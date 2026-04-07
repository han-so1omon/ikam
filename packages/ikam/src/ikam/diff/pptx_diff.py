"""
PPTX diff engine for PowerPoint presentations.

Provides slide-level comparison with text/shape change detection and image hash-based diffing.
"""
from typing import Any, Dict, List, Optional
from pathlib import Path
import hashlib

try:
    from pptx import Presentation
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .types import DiffChange


class SlideChange:
    """Represents a change to a slide."""
    
    def __init__(
        self,
        operation: str,  # "added", "removed", "modified", "reordered"
        slide_index: int,
        changes: Optional[List[str]] = None,
        old_index: Optional[int] = None,
        new_index: Optional[int] = None,
    ):
        self.operation = operation
        self.slide_index = slide_index
        self.changes = changes or []
        self.old_index = old_index
        self.new_index = new_index


def _extract_slide_text(slide: "Slide") -> str:
    """Extract all text content from a slide."""
    if not PPTX_AVAILABLE:
        return ""
    
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            texts.append(shape.text)
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texts.append(run.text)
    
    return "\n".join(texts)


def _hash_image(shape: "BaseShape") -> Optional[str]:
    """Generate hash for image shape."""
    if not PPTX_AVAILABLE:
        return None
    
    if not hasattr(shape, "image"):
        return None
    
    try:
        image_bytes = shape.image.blob
        return hashlib.blake2b(image_bytes, digest_size=16).hexdigest()
    except Exception:
        return None


def _extract_slide_fingerprint(slide: "Slide", slide_index: int) -> Dict[str, Any]:
    """
    Extract a fingerprint for a slide including text, shapes, and image hashes.
    
    Returns a dict with:
    - index: slide position
    - text: extracted text content
    - shape_count: number of shapes
    - image_hashes: list of image content hashes
    """
    if not PPTX_AVAILABLE:
        return {}
    
    text = _extract_slide_text(slide)
    shape_count = len(slide.shapes)
    
    image_hashes = []
    for shape in slide.shapes:
        img_hash = _hash_image(shape)
        if img_hash:
            image_hashes.append(img_hash)
    
    return {
        "index": slide_index,
        "text": text,
        "shape_count": shape_count,
        "image_hashes": image_hashes,
    }


def _compare_slides(
    old_fingerprint: Dict[str, Any],
    new_fingerprint: Dict[str, Any],
) -> List[str]:
    """
    Compare two slide fingerprints and return list of changes.
    
    Returns descriptions like:
    - "text_modified"
    - "shapes_added" or "shapes_removed"
    - "image_modified"
    """
    changes = []
    
    if old_fingerprint["text"] != new_fingerprint["text"]:
        changes.append("text_modified")
    
    old_shapes = old_fingerprint["shape_count"]
    new_shapes = new_fingerprint["shape_count"]
    if old_shapes != new_shapes:
        if new_shapes > old_shapes:
            changes.append(f"shapes_added ({new_shapes - old_shapes})")
        else:
            changes.append(f"shapes_removed ({old_shapes - new_shapes})")
    
    old_images = set(old_fingerprint["image_hashes"])
    new_images = set(new_fingerprint["image_hashes"])
    
    if old_images != new_images:
        added = new_images - old_images
        removed = old_images - new_images
        if added:
            changes.append(f"images_added ({len(added)})")
        if removed:
            changes.append(f"images_removed ({len(removed)})")
    
    return changes


def diff_pptx(old_bytes: bytes, new_bytes: bytes) -> List[DiffChange]:
    """
    Compute structural diff between two PPTX files.
    
    Returns a list of DiffChange objects with:
    - path: "Slide {index}"
    - change_type: "added", "removed", "modified"
    - old_value/new_value: slide fingerprints (for "modified")
    
    Raises:
        ImportError: if python-pptx is not installed
        ValueError: if files cannot be parsed
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx library required for PPTX diff. "
            "Install with: pip install python-pptx"
        )
    
    try:
        # Load presentations from bytes
        from io import BytesIO
        old_prs = Presentation(BytesIO(old_bytes))
        new_prs = Presentation(BytesIO(new_bytes))
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX files: {e}")
    
    # Extract fingerprints for all slides
    old_fingerprints = [
        _extract_slide_fingerprint(slide, i)
        for i, slide in enumerate(old_prs.slides)
    ]
    new_fingerprints = [
        _extract_slide_fingerprint(slide, i)
        for i, slide in enumerate(new_prs.slides)
    ]
    
    changes: List[DiffChange] = []
    
    # Simple approach: compare by index first, then detect adds/removes
    old_count = len(old_fingerprints)
    new_count = len(new_fingerprints)
    
    # Compare overlapping slides
    for i in range(min(old_count, new_count)):
        old_fp = old_fingerprints[i]
        new_fp = new_fingerprints[i]
        
        slide_changes = _compare_slides(old_fp, new_fp)
        
        if slide_changes:
            # Store changes as a dict in the values for modified slides
            old_val = {**old_fp, "changes": slide_changes}
            new_val = {**new_fp, "changes": slide_changes}
            
            changes.append(
                DiffChange(
                    path=f"Slide {i + 1}",
                    change_type="modified",
                    old_value=old_val,
                    new_value=new_val,
                )
            )
    
    # Detect added slides
    if new_count > old_count:
        for i in range(old_count, new_count):
            new_fp = new_fingerprints[i]
            changes.append(
                DiffChange(
                    path=f"Slide {i + 1}",
                    change_type="added",
                    old_value=None,
                    new_value=new_fp,
                )
            )
    
    # Detect removed slides
    if old_count > new_count:
        for i in range(new_count, old_count):
            old_fp = old_fingerprints[i]
            changes.append(
                DiffChange(
                    path=f"Slide {i + 1}",
                    change_type="removed",
                    old_value=old_fp,
                    new_value=None,
                )
            )
    
    return changes


def diff_pptx_files(old_path: Path, new_path: Path) -> List[DiffChange]:
    """
    Convenience function to diff PPTX files from disk.
    
    Args:
        old_path: Path to old PPTX file
        new_path: Path to new PPTX file
    
    Returns:
        List of DiffChange objects describing differences
    """
    with open(old_path, "rb") as f:
        old_bytes = f.read()
    with open(new_path, "rb") as f:
        new_bytes = f.read()
    
    return diff_pptx(old_bytes, new_bytes)
