"""TDD tests for artifact download endpoint.

The download endpoint reconstructs artifact content from V3 fragments
stored in GraphSnapshot, using IKAM algebra (root relation DAG traversal).
"""

from fastapi.testclient import TestClient
from uuid import uuid4

from ikam.fragments import Fragment
from ikam_perf_report.benchmarks import runner
from ikam_perf_report.benchmarks.store import GraphSnapshot, STORE
from ikam_perf_report.main import app


def _run_case_and_await_async_pipeline(client: TestClient, case_id: str) -> tuple[str, str]:
    """Run benchmark and advance async debug pipeline to completion."""
    run_resp = client.post(
        "/benchmarks/run",
        params={"case_ids": case_id, "include_evaluation": "false"},
    )
    assert run_resp.status_code == 200
    run = run_resp.json()["runs"][0]
    run_id = run["run_id"]
    graph_id = run["graph_id"]

    state_resp = client.get(f"/benchmarks/runs/{run_id}/debug-state")
    assert state_resp.status_code == 200
    state = state_resp.json()
    assert state.get("status") == "ok"

    control_resp = client.post(
        f"/benchmarks/runs/{run_id}/control",
        json={
            "command_id": str(uuid4()),
            "action": "resume",
            "pipeline_id": state["pipeline_id"],
            "pipeline_run_id": state["pipeline_run_id"],
        },
    )
    assert control_resp.status_code == 200
    control = control_resp.json()
    assert control.get("status") == "ok"
    assert control.get("state", {}).get("execution_state") in {"completed", "paused"}

    return run_id, graph_id


def test_graph_snapshot_stores_fragments():
    """GraphSnapshot should accept a 'fragments' field holding V3 Fragment objects."""
    frag = Fragment(cas_id="abc123", value="hello", mime_type="text/plain")
    snap = GraphSnapshot(graph_id="test-graph", fragments=[frag])
    assert snap.fragments == [frag]
    assert snap.fragments[0].cas_id == "abc123"


def test_graph_snapshot_fragments_defaults_empty():
    """GraphSnapshot.fragments should default to empty list for backward compat."""
    snap = GraphSnapshot(graph_id="test-graph")
    assert snap.fragments == []


def test_run_benchmark_populates_fragments(case_fixtures_root, monkeypatch):
    """run_benchmark() should store V3 fragments in GraphSnapshot for reconstruction."""
    monkeypatch.setattr(
        runner,
        "run_semantic_pipeline",
        lambda text: {"entities": [{"id": "e1"}], "relations": []},
    )
    STORE.reset()
    client = TestClient(app)
    run_id, _ = _run_case_and_await_async_pipeline(client, "s-construction-v01")
    record = STORE.get_run(run_id)
    assert record is not None
    assert record.graph.fragments, "GraphSnapshot must retain V3 fragments for reconstruction"
    # Fragments should include at least one Fragment object with cas_id
    assert any(hasattr(f, "cas_id") and f.cas_id for f in record.graph.fragments)


def test_download_artifact_returns_reconstructed_content(case_fixtures_root, monkeypatch):
    """GET /artifacts/{artifact_id}/download reconstructs content via IKAM algebra."""
    monkeypatch.setattr(
        runner,
        "run_semantic_pipeline",
        lambda text: {"entities": [{"id": "e1"}], "relations": []},
    )
    STORE.reset()
    client = TestClient(app)

    # Run and await async benchmark completion to populate graph/fragments
    _, graph_id = _run_case_and_await_async_pipeline(client, "s-construction-v01")

    # Find a text-type artifact node ID from the graph
    nodes_resp = client.get("/graph/nodes", params={"graph_id": graph_id})
    assert nodes_resp.status_code == 200
    nodes = nodes_resp.json()
    artifact_nodes = [n for n in nodes if n.get("kind") == "artifact" or n.get("type") == "artifact"]
    assert artifact_nodes, "Should have at least one artifact node"

    # Pick a text artifact (not the image one, since binary reconstruction differs)
    text_artifact = None
    for anode in artifact_nodes:
        if anode["label"] not in ("image",):
            text_artifact = anode
            break
    assert text_artifact is not None, "Should have a non-image artifact"
    
    import urllib.parse
    artifact_id = text_artifact["id"]
    encoded_id = urllib.parse.quote(artifact_id, safe="")
    resp = client.get(f"/artifacts/{encoded_id}/download")
    assert resp.status_code == 200
    assert resp.headers.get("content-disposition"), "Should have Content-Disposition header"
    assert "attachment" in resp.headers["content-disposition"]
    # The reconstructed content should be non-empty
    assert len(resp.content) > 0


def test_download_artifact_404_when_not_found(case_fixtures_root, monkeypatch):
    """GET /artifacts/{artifact_id}/download returns 404 for unknown artifact."""
    monkeypatch.setattr(
        runner,
        "run_semantic_pipeline",
        lambda text: {"entities": [{"id": "e1"}], "relations": []},
    )
    STORE.reset()
    client = TestClient(app)
    resp = client.get("/artifacts/artifact:nonexistent123/download")
    assert resp.status_code == 404


def test_download_preserves_original_filename(case_fixtures_root, monkeypatch):
    """Download should use original filename (e.g. 'revenue_plan.md'), not a heuristic."""
    monkeypatch.setattr(
        runner,
        "run_semantic_pipeline",
        lambda text: {"entities": [{"id": "e1"}], "relations": []},
    )
    STORE.reset()
    client = TestClient(app)

    _, graph_id = _run_case_and_await_async_pipeline(client, "s-construction-v01")

    nodes_resp = client.get("/graph/nodes", params={"graph_id": graph_id})
    nodes = nodes_resp.json()
    artifact_nodes = [n for n in nodes if n.get("type") == "artifact"]

    # Find the 'metrics' artifact — its original file is metrics.json.
    # The old heuristic would produce 'metrics.md' (wrong extension).
    target = None
    for anode in artifact_nodes:
        label = anode.get("label", "")
        if "metrics" in label.lower():
            target = anode
            break

    assert target is not None, "Should find the metrics artifact"
    import urllib.parse
    encoded_id = urllib.parse.quote(target['id'], safe="")
    resp = client.get(f"/artifacts/{encoded_id}/download")
    assert resp.status_code == 200

    disposition = resp.headers.get("content-disposition", "")
    # Must contain the original filename with correct extension
    assert "metrics.json" in disposition, (
        f"Expected original filename 'metrics.json' in Content-Disposition, got: {disposition}"
    )


def test_artifact_node_stores_file_name():
    """fragments_to_graph() should store file_name on artifact nodes when provided."""
    from ikam_perf_report.benchmarks.ikam_flow import fragments_to_graph

    frag = Fragment(cas_id="abc123", value="hello world", mime_type="text/plain")
    nodes, edges, _ = fragments_to_graph(
        [frag],
        labels_by_fragment_id={"abc123": "revenue plan"},
        filenames_by_label={"revenue plan": "revenue_plan.md"},
    )
    artifact_nodes = [n for n in nodes if n["type"] == "artifact"]
    assert artifact_nodes, "Should have an artifact node"
    assert artifact_nodes[0].get("file_name") == "revenue_plan.md", (
        f"Artifact node should have file_name='revenue_plan.md', got: {artifact_nodes[0].get('file_name')}"
    )


def test_download_binary_artifact_uses_binary_reconstruction(case_fixtures_root, monkeypatch):
    """Binary artifacts with 'plan' in label should not be treated as text."""
    monkeypatch.setattr(
        runner,
        "run_semantic_pipeline",
        lambda text: {"entities": [{"id": "e1"}], "relations": []},
    )
    # Create a real pptx fixture whose label includes "plan".
    from io import BytesIO
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Delivery Plan"
    slide.placeholders[1].text = "Binary artifact for reconstruction test"
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    (case_fixtures_root / "s-construction-v01" / "deliverable-plan.pptx").write_bytes(pptx_bytes.getvalue())
    STORE.reset()
    client = TestClient(app)

    _, graph_id = _run_case_and_await_async_pipeline(client, "s-construction-v01")

    nodes_resp = client.get("/graph/nodes", params={"graph_id": graph_id})
    assert nodes_resp.status_code == 200
    nodes = nodes_resp.json()

    pptx_node = None
    for node in nodes:
        if node.get("type") != "artifact":
            continue
        file_name = (
            (node.get("meta") or {}).get("file_name")
            or node.get("file_name")
            or node.get("id")
        )
        if isinstance(file_name, str) and file_name.endswith(".pptx") and "plan" in file_name:
            pptx_node = node
            break

    assert pptx_node is not None, "Should find a pptx artifact with plan in name"

    import urllib.parse
    encoded_id = urllib.parse.quote(pptx_node['id'], safe="")
    resp = client.get(f"/artifacts/{encoded_id}/download")
    assert resp.status_code == 200
    disposition = resp.headers.get("content-disposition", "")
    assert ".pptx" in disposition, f"Expected .pptx filename, got: {disposition}"
