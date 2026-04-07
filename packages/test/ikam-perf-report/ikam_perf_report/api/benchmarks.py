from __future__ import annotations

import asyncio
import base64
from copy import deepcopy
from dataclasses import dataclass
from datetime import UTC, datetime
from functools import lru_cache
import io
import json
import os
from pathlib import Path
import sys
from threading import get_ident
import zipfile
from time import perf_counter
from typing import Any
from uuid import uuid4
from xml.etree import ElementTree

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field
from ikam.fragments import Relation, is_relation_fragment
from ikam.inspection import InspectionEdge, InspectionNode, InspectionRef, InspectionSubgraph, ResolveInspectionRequest, edge_id_for, node_id_for
from interacciones.schemas import WorkflowDefinition
from ikam.forja.boundary_policy import build_boundary_diagnostic
from ikam.forja.debug_execution import StepExecutionState, execute_step, next_step_name

from ikam_perf_report.benchmarks.case_fixtures import available_cases, case_fixture_dir, load_case_fixture
from ikam_perf_report.benchmarks.debug_models import DebugRunState, DebugStepEvent
from modelado.core.execution_scope import DefaultExecutionScope
from modelado.environment_scope import EnvironmentScope
from modelado.inspection_runtime import HotInspectionResolver, PersistentInspectionResolver
from modelado.executors.transition_validation import build_runtime_transition_validation, build_runtime_transition_validation_for_direction
from modelado.oraculo.persistent_graph_state import PersistentGraphState
from ikam_perf_report.benchmarks.ikam_flow import fragments_to_graph
from ikam_perf_report.benchmarks.runner import (
    run_benchmark,
    run_merge_benchmark,
    run_staging_normalize_promote_enrich_poc,
)
from ikam_perf_report.benchmarks.store import BenchmarkRunRecord, GraphSnapshot, STORE

router = APIRouter(prefix="/benchmarks", tags=["benchmarks"])
_RUN_STEP_TASKS: dict[str, asyncio.Task[None]] = {}
_RUNNING_LOG_HEARTBEAT_INTERVAL_SECONDS = 5.0
_EXECUTOR_STREAM_REDIRECT_LOCK = asyncio.Lock()


class _StreamingLogBuffer:
    def __init__(self, sink: list[str], publish: callable, record_line: callable | None = None) -> None:
        self._sink = sink
        self._publish = publish
        self._record_line = record_line
        self._pending = ""

    def write(self, text: str) -> int:
        if not text:
            return 0
        self._pending += text
        while "\n" in self._pending:
            line, self._pending = self._pending.split("\n", 1)
            if line.strip():
                self._sink.append(line)
                if self._record_line is not None:
                    self._record_line(line)
                self._publish()
        return len(text)

    def flush(self) -> None:
        if self._pending.strip():
            self._sink.append(self._pending)
            if self._record_line is not None:
                self._record_line(self._pending)
            self._pending = ""
            self._publish()


class _ThreadScopedStreamProxy:
    def __init__(self, *, original: Any, target_thread_id: list[int | None], buffer: _StreamingLogBuffer) -> None:
        self._original = original
        self._target_thread_id = target_thread_id
        self._buffer = buffer

    def write(self, text: str) -> int:
        if self._target_thread_id[0] == get_ident():
            return self._buffer.write(text)
        return self._original.write(text)

    def flush(self) -> None:
        if self._target_thread_id[0] == get_ident():
            self._buffer.flush()
            return
        self._original.flush()

    def __getattr__(self, name: str) -> Any:
        return getattr(self._original, name)


def _append_log_event(
    log_events: list[dict[str, Any]],
    *,
    at: str,
    source: str,
    stream: str,
    message: str,
) -> None:
    normalized_message = str(message).strip()
    if not normalized_message:
        return
    log_events.append(
        {
            "seq": len(log_events) + 1,
            "at": at,
            "source": source,
            "stream": stream,
            "message": normalized_message,
        }
    )


def _merge_log_blocks(
    *,
    executor_stdout_lines: list[str],
    executor_stderr_lines: list[str],
    system_stdout_lines: list[str],
    system_stderr_lines: list[str],
) -> dict[str, list[str]]:
    return {
        "stdout_lines": list(executor_stdout_lines) + list(system_stdout_lines),
        "stderr_lines": list(executor_stderr_lines) + list(system_stderr_lines),
    }


def _legacy_logs_from_log_events(log_events: list[dict[str, Any]]) -> dict[str, list[str]]:
    return {
        "stdout_lines": [
            str(item.get("message"))
            for item in log_events
            if isinstance(item, dict) and item.get("stream") == "stdout" and str(item.get("message") or "").strip()
        ],
        "stderr_lines": [
            str(item.get("message"))
            for item in log_events
            if isinstance(item, dict) and item.get("stream") == "stderr" and str(item.get("message") or "").strip()
        ],
    }


def _named_logs_from_log_events(log_events: list[dict[str, Any]], *, source: str) -> dict[str, list[str]]:
    return {
        "stdout_lines": [
            str(item.get("message"))
            for item in log_events
            if isinstance(item, dict)
            and item.get("source") == source
            and item.get("stream") == "stdout"
            and str(item.get("message") or "").strip()
        ],
        "stderr_lines": [
            str(item.get("message"))
            for item in log_events
            if isinstance(item, dict)
            and item.get("source") == source
            and item.get("stream") == "stderr"
            and str(item.get("message") or "").strip()
        ],
    }


def _publish_running_event_logs(
    *,
    run_id: str,
    step_id: str,
    executor_stdout_lines: list[str],
    executor_stderr_lines: list[str],
    system_stdout_lines: list[str],
    system_stderr_lines: list[str],
    log_events: list[dict[str, Any]],
) -> None:
    try:
        STORE.update_debug_event(
            run_id=run_id,
            step_id=step_id,
            metrics={
                "executor_logs": _named_logs_from_log_events(log_events, source="executor"),
                "system_logs": _named_logs_from_log_events(log_events, source="system"),
                "logs": _legacy_logs_from_log_events(log_events),
                "log_events": list(log_events),
            },
        )
    except KeyError:
        return


async def _run_execute_step_with_thread_scoped_capture(
    *,
    implementation_next_name: str,
    execution_state: StepExecutionState,
    debug_scope: DefaultExecutionScope,
    stdout_buffer: _StreamingLogBuffer,
    stderr_buffer: _StreamingLogBuffer,
) -> Any:
    stdout_thread_id: list[int | None] = [None]
    stderr_thread_id: list[int | None] = stdout_thread_id
    stdout_proxy = _ThreadScopedStreamProxy(original=sys.stdout, target_thread_id=stdout_thread_id, buffer=stdout_buffer)
    stderr_proxy = _ThreadScopedStreamProxy(original=sys.stderr, target_thread_id=stderr_thread_id, buffer=stderr_buffer)

    def _run_in_worker_thread() -> Any:
        stdout_thread_id[0] = get_ident()
        try:
            return asyncio.run(
                execute_step(
                    implementation_next_name,
                    execution_state,
                    scope=debug_scope,
                )
            )
        finally:
            stdout_thread_id[0] = None

    async with _EXECUTOR_STREAM_REDIRECT_LOCK:
        original_stdout = sys.stdout
        original_stderr = sys.stderr
        sys.stdout = stdout_proxy
        sys.stderr = stderr_proxy
        try:
            return await asyncio.to_thread(_run_in_worker_thread)
        finally:
            sys.stdout = original_stdout
            sys.stderr = original_stderr
            stdout_buffer.flush()
            stderr_buffer.flush()


@dataclass
class _PreparedPipelineStep:
    debug_scope: DefaultExecutionScope
    implementation_next_name: str
    runtime_context: dict[str, Any]
    execution_state: StepExecutionState
    running_event: DebugStepEvent
    retry_parent_step_id: str | None
    start_iso: str
    start_perf: float
    executor_stdout_lines: list[str]
    executor_stderr_lines: list[str]
    system_stdout_lines: list[str]
    system_stderr_lines: list[str]
    log_events: list[dict[str, Any]]

_CANONICAL_DEBUG_PIPELINE_STEPS = [
    "init.initialize",
    "map.conceptual.lift.surface_fragments",
    "map.conceptual.lift.entities_and_relationships",
    "map.conceptual.lift.claims",
    "map.conceptual.lift.summarize",
    "map.conceptual.embed.discovery_index",
    "map.conceptual.normalize.discovery",
    "map.reconstructable.embed",
    "map.reconstructable.search.dependency_resolution",
    "map.reconstructable.normalize",
    "map.reconstructable.compose.reconstruction_programs",
    "map.conceptual.verify.discovery_gate",
    "map.conceptual.commit.semantic_only",
    "map.reconstructable.build_subgraph.reconstruction",
]


def _canonical_debug_pipeline_steps(scope: DefaultExecutionScope | None = None) -> list[str]:
    base_scope = scope or DefaultExecutionScope()
    raw_steps = base_scope.get_dynamic_execution_steps()
    if raw_steps == ["init.initialize", "parse_artifacts", "lift_fragments"]:
        return list(_CANONICAL_DEBUG_PIPELINE_STEPS)
    return list(raw_steps)


def _node_id_to_step_name(node_id: str) -> str:
    head, sep, tail = node_id.partition("-")
    if not sep:
        return node_id.replace("-", "_")
    return f"{head}.{tail.replace('-', '_')}"


@lru_cache(maxsize=32)
def _declared_workflow_steps(pipeline_id: str) -> tuple[str, ...]:
    workflow_path = Path(__file__).resolve().parents[2] / "preseed" / "workflows" / f"{pipeline_id}.yaml"
    if not workflow_path.exists():
        return ()

    try:
        import yaml

        payload = yaml.safe_load(workflow_path.read_text(encoding="utf-8")) or {}
    except Exception:
        return ()

    if payload.get("workflow_id") != pipeline_id:
        return ()

    nodes = payload.get("nodes")
    if not isinstance(nodes, list):
        return ()

    steps: list[str] = []
    for node in nodes:
        if not isinstance(node, dict):
            continue
        node_id = node.get("node_id")
        if isinstance(node_id, str) and node_id:
            steps.append(_node_id_to_step_name(node_id))
    return tuple(steps)


def _debug_pipeline_steps(pipeline_id: str | None, scope: DefaultExecutionScope | None = None) -> list[str]:
    if pipeline_id:
        declared = list(_declared_workflow_steps(pipeline_id))
        if declared:
            return declared
    return _canonical_debug_pipeline_steps(scope)


@lru_cache(maxsize=32)
def _workflow_definition(pipeline_id: str) -> WorkflowDefinition | None:
    workflow_path = Path(__file__).resolve().parents[2] / "preseed" / "workflows" / f"{pipeline_id}.yaml"
    if not workflow_path.exists():
        return None

    try:
        import yaml

        payload = yaml.safe_load(workflow_path.read_text(encoding="utf-8")) or {}
        return WorkflowDefinition.model_validate(payload)
    except Exception:
        return None


def _workflow_node_id(step_name: str) -> str:
    presentation = _presentation_step_name("ingestion-early-parse", step_name)
    return presentation.replace(".", "-").replace("_", "-")


def _compiled_transition_step_name(pipeline_id: str, step_name: str) -> str:
    if pipeline_id != "ingestion-early-parse":
        return step_name
    aliases = {
        "load_documents": "load.documents",
        "parse_chunk": "parse.chunk",
        "parse_entities_and_relationships": "parse.entities_and_relationships",
        "parse_claims": "parse.claims",
        "parse_artifacts": "parse_artifacts",
        "lift_fragments": "lift_fragments",
    }
    return aliases.get(step_name, step_name)


@lru_cache(maxsize=32)
def _executor_declarations_by_capability(pipeline_id: str) -> dict[str, dict[str, str]]:
    declarations_dir = Path(__file__).resolve().parents[2] / "preseed" / "declarations"
    if not declarations_dir.exists():
        return {}

    import yaml

    capability_map: dict[str, dict[str, str]] = {}
    for declaration_path in declarations_dir.glob("*.yaml"):
        try:
            payload = yaml.safe_load(declaration_path.read_text(encoding="utf-8")) or {}
        except Exception:
            continue
        executor_id = str(payload.get("executor_id") or "").strip()
        executor_kind = str(payload.get("executor_kind") or "").strip()
        capabilities = payload.get("capabilities") if isinstance(payload.get("capabilities"), list) else []
        if not executor_id or not executor_kind:
            continue
        for capability in capabilities:
            if isinstance(capability, str) and capability.strip():
                capability_map[capability.strip()] = {
                    "executor_id": executor_id,
                    "executor_kind": executor_kind,
                }
    return capability_map


def _workflow_validators(pipeline_id: str | None, step_name: str) -> list[dict[str, Any]]:
    if not pipeline_id:
        return []
    workflow = _workflow_definition(pipeline_id)
    if workflow is None:
        return []
    node_id = _workflow_node_id(step_name)
    for node in workflow.nodes:
        if node.node_id != node_id:
            continue
        return [validator.model_dump(mode="json") for validator in node.validators]
    return []


@lru_cache(maxsize=32)
def _compiled_operator_metadata(pipeline_id: str) -> dict[str, dict[str, Any]]:
    if pipeline_id != "ingestion-early-parse":
        return {}

    import yaml

    compiled_dir = Path(__file__).resolve().parents[2] / "preseed" / "compiled"
    net_path = compiled_dir / "ingestion_net_v2.yaml"
    if not net_path.exists():
        return {}

    try:
        net_payload = yaml.safe_load(net_path.read_text(encoding="utf-8")) or {}
    except Exception:
        return {}

    fragments = net_payload.get("fragments") if isinstance(net_payload.get("fragments"), list) else []
    transition_operator_fragment_ids: dict[str, str] = {}
    transition_ids_to_step_names: dict[str, str] = {}

    for fragment in fragments:
        if not isinstance(fragment, dict):
            continue
        value = fragment.get("value") if isinstance(fragment.get("value"), dict) else {}
        if value.get("predicate") == "executed_by_operator":
            subject = str(value.get("subject") or "")
            operator_fragment_id = str(value.get("object") or "")
            if subject and operator_fragment_id:
                transition_operator_fragment_ids[subject] = operator_fragment_id
        transitions = value.get("transitions") if isinstance(value.get("transitions"), dict) else None
        if transitions:
            for step_name, transition_id in transitions.items():
                if isinstance(step_name, str) and step_name and isinstance(transition_id, str) and transition_id:
                    transition_ids_to_step_names[transition_id] = step_name

    operator_file_by_fragment_id = {
        "24f19013-3759-5991-a148-39076da4284e": compiled_dir / "load_documents_operator.yaml",
        "adf6d2e6-c660-5947-901d-2726b057681d": compiled_dir / "parse_chunk_operator.yaml",
        "22250859-c67e-5bf1-9370-3d59b7cf4c0b": compiled_dir / "entities_and_relationships_operator.yaml",
        "773da82d-e6d0-5c12-8657-7fe44f9720e7": compiled_dir / "claims_operator.yaml",
    }

    operator_metadata: dict[str, dict[str, Any]] = {}
    workflow = _workflow_definition(pipeline_id)
    workflow_nodes = {node.node_id: node for node in workflow.nodes} if workflow is not None else {}

    for transition_id, compiled_step_name in transition_ids_to_step_names.items():
        step_name = _compiled_transition_step_name(pipeline_id, compiled_step_name)
        operator_fragment_id = transition_operator_fragment_ids.get(transition_id)
        if not operator_fragment_id:
            continue
        operator_file = operator_file_by_fragment_id.get(operator_fragment_id)
        if operator_file is None or not operator_file.exists():
            continue
        try:
            operator_payload = yaml.safe_load(operator_file.read_text(encoding="utf-8")) or {}
        except Exception:
            continue
        operator_fragments = operator_payload.get("fragments") if isinstance(operator_payload.get("fragments"), list) else []
        operator_id = None
        executor_fragment_id = None
        for operator_fragment in operator_fragments:
            if not isinstance(operator_fragment, dict):
                continue
            value = operator_fragment.get("value") if isinstance(operator_fragment.get("value"), dict) else {}
            if value.get("operator_id"):
                operator_id = str(value.get("operator_id"))
            if value.get("predicate") == "executed_by":
                executor_fragment_id = str(value.get("object") or "") or None
        node = workflow_nodes.get(_workflow_node_id(step_name))
        policy = dict(node.policy) if node is not None else {}
        capability = str(getattr(node, "capability", "") or "") if node is not None else ""
        metadata: dict[str, Any] = {
            "step_name": step_name,
            "implementation_step_name": _implementation_step_name(pipeline_id, step_name),
            "operator_id": operator_id,
            "executor_fragment_id": executor_fragment_id,
            "capability": capability,
            "policy": policy,
        }
        executor_identity = _executor_declarations_by_capability(pipeline_id).get(capability) if capability else None
        if isinstance(executor_identity, dict):
            metadata.update(executor_identity)
        if node is not None and node.operator_selection:
            metadata["operator_selection"] = dict(node.operator_selection)
        if node is not None and node.executor_selection:
            metadata["executor_selection"] = dict(node.executor_selection)
        operator_metadata[step_name] = metadata
        implementation_step_name = metadata["implementation_step_name"]
        if implementation_step_name != step_name:
            operator_metadata[implementation_step_name] = metadata
    return operator_metadata


def _document_set_payload(*, artifact_id: str, fragment_refs: list[str], subgraph_ref: str) -> dict[str, Any]:
    return {
        "kind": "document_set",
        "artifact_head_ref": artifact_id,
        "subgraph_ref": subgraph_ref,
        "document_refs": fragment_refs,
    }


def _chunk_extraction_set_payload(*, source_subgraph_ref: str, subgraph_ref: str, extraction_refs: list[str]) -> dict[str, Any]:
    return {
        "kind": "chunk_extraction_set",
        "source_subgraph_ref": source_subgraph_ref,
        "subgraph_ref": subgraph_ref,
        "extraction_refs": extraction_refs,
    }


def _entity_relationship_set_payload(*, source_subgraph_ref: str, subgraph_ref: str, entity_relationship_refs: list[str]) -> dict[str, Any]:
    return {
        "kind": "entity_relationship_set",
        "source_subgraph_ref": source_subgraph_ref,
        "subgraph_ref": subgraph_ref,
        "entity_relationship_refs": entity_relationship_refs,
    }


def _claim_set_payload(*, source_subgraph_ref: str, subgraph_ref: str, claim_refs: list[str]) -> dict[str, Any]:
    return {
        "kind": "claim_set",
        "source_subgraph_ref": source_subgraph_ref,
        "subgraph_ref": subgraph_ref,
        "claim_refs": claim_refs,
    }


def _persist_runtime_hot_document_set(
    *,
    implementation_next_name: str,
    run_id: str,
    step_id: str,
    runtime_context: dict[str, Any],
    step_outputs: dict[str, Any],
) -> None:
    if implementation_next_name != "parse_artifacts":
        return
    document_fragment_refs = [
        str(item)
        for item in (step_outputs.get("document_fragment_refs") or [])
        if isinstance(item, str) and item
    ]
    if not document_fragment_refs:
        return
    artifact_id = str(runtime_context.get("artifact_id") or step_outputs.get("artifact_id") or "")
    document_set_payload = _document_set_payload(
        artifact_id=artifact_id,
        fragment_refs=document_fragment_refs,
        subgraph_ref="",
    )
    documents = [item for item in (step_outputs.get("documents") or []) if isinstance(item, dict)]
    hydrated_documents: list[dict[str, Any]] = []
    for index, fragment_ref in enumerate(document_fragment_refs):
        document_value = dict(documents[index]) if index < len(documents) else {}
        hydrated_documents.append(
            {
                "cas_id": fragment_ref,
                "mime_type": "application/vnd.ikam.loaded-document+json",
                "value": document_value,
            }
        )
    if hydrated_documents:
        document_set_payload["documents"] = hydrated_documents
    document_set_ref = STORE.put_hot_subgraph(
        run_id=run_id,
        step_id=step_id,
        contract_type="document_set",
        payload=document_set_payload,
    )
    runtime_context["document_set_ref"] = document_set_ref


def _persist_runtime_hot_chunk_extraction_set(
    *,
    run_id: str,
    step_id: str,
    runtime_context: dict[str, Any],
    step_outputs: dict[str, Any],
) -> None:
    extraction_refs = _extract_runtime_output_fragment_ids(step_outputs)
    if not extraction_refs:
        return
    source_subgraph_ref = runtime_context.get("document_set_ref")
    if not isinstance(source_subgraph_ref, str) or not source_subgraph_ref:
        return
    chunk_rows = [item for item in (step_outputs.get("chunks") or []) if isinstance(item, dict)]
    extractions: list[dict[str, Any]] = []
    document_chunk_sets: list[dict[str, Any]] = []
    edges: list[dict[str, str]] = []
    for index, chunk in enumerate(chunk_rows):
        fragment_ref = str(chunk.get("fragment_id") or chunk.get("chunk_id") or "")
        if not fragment_ref and index < len(extraction_refs):
            fragment_ref = extraction_refs[index]
        if not fragment_ref:
            continue
        chunk_value = dict(chunk)
        chunk_value["fragment_id"] = fragment_ref
        source_document_fragment_id = str(chunk.get("source_document_fragment_id") or "")
        if not source_document_fragment_id:
            chunk_document_id = str(chunk.get("document_id") or "")
            document_set_payload = STORE.get_hot_subgraph(source_subgraph_ref)
            if isinstance(document_set_payload, dict):
                documents = document_set_payload.get("documents") if isinstance(document_set_payload.get("documents"), list) else []
                matched_document = next(
                    (
                        item
                        for item in documents
                        if isinstance(item, dict)
                        and isinstance(item.get("value"), dict)
                        and str((item.get("value") or {}).get("document_id") or "") == chunk_document_id
                    ),
                    None,
                )
                if isinstance(matched_document, dict):
                    source_document_fragment_id = str(matched_document.get("cas_id") or matched_document.get("id") or "")
        if source_document_fragment_id:
            chunk_value["source_document_fragment_id"] = source_document_fragment_id
        extractions.append(
            {
                "cas_id": fragment_ref,
                "mime_type": "application/vnd.ikam.chunk-extraction+json",
                "value": chunk_value,
            }
        )
        if source_document_fragment_id:
            edges.append(
                {
                    "from": f"fragment:{fragment_ref}",
                    "to": f"fragment:{source_document_fragment_id}",
                    "edge_label": "knowledge:derives",
                }
            )
    chunk_extraction_set_payload = _chunk_extraction_set_payload(
        source_subgraph_ref=source_subgraph_ref,
        subgraph_ref="",
        extraction_refs=extraction_refs,
    )
    if extractions:
        chunk_extraction_set_payload["extractions"] = extractions
        chunk_extraction_set_payload["chunk_extractions"] = extractions
    chunk_extraction_output = step_outputs.get("chunk_extraction_set") if isinstance(step_outputs.get("chunk_extraction_set"), dict) else {}
    raw_document_chunk_sets = step_outputs.get("document_chunk_sets")
    if isinstance(raw_document_chunk_sets, list) and raw_document_chunk_sets:
        for document_chunk_set in raw_document_chunk_sets:
            if not isinstance(document_chunk_set, dict):
                continue
            fragment_ref = str(document_chunk_set.get("cas_id") or document_chunk_set.get("id") or "")
            document_chunk_set_value = document_chunk_set.get("value") if isinstance(document_chunk_set.get("value"), dict) else {}
            if not fragment_ref:
                fragment_ref = str(document_chunk_set_value.get("fragment_id") or "")
            if not fragment_ref:
                document_id = str(document_chunk_set_value.get("document_id") or len(document_chunk_sets))
                fragment_ref = f"document-chunk-set:{document_id}"
            document_chunk_set_value = dict(document_chunk_set_value)
            document_chunk_set_value["fragment_id"] = fragment_ref
            document_chunk_sets.append(
                {
                    "cas_id": fragment_ref,
                    "mime_type": str(document_chunk_set.get("mime_type") or "application/vnd.ikam.document-chunk-set+json"),
                    "value": document_chunk_set_value,
                }
            )
    else:
        for document_chunk_set in chunk_extraction_output.get("document_chunk_sets") or []:
            if not isinstance(document_chunk_set, dict):
                continue
            document_chunk_set_value = dict(document_chunk_set)
            fragment_ref = str(document_chunk_set_value.get("fragment_id") or "")
            if not fragment_ref:
                document_id = str(document_chunk_set_value.get("document_id") or len(document_chunk_sets))
                fragment_ref = f"document-chunk-set:{document_id}"
            document_chunk_set_value["fragment_id"] = fragment_ref
            document_chunk_sets.append(
                {
                    "cas_id": fragment_ref,
                    "mime_type": "application/vnd.ikam.document-chunk-set+json",
                    "value": document_chunk_set_value,
                }
            )
    if document_chunk_sets:
        chunk_extraction_set_payload["document_chunk_sets"] = document_chunk_sets
    document_set_payload = STORE.get_hot_subgraph(source_subgraph_ref)
    if isinstance(document_set_payload, dict) and isinstance(document_set_payload.get("documents"), list):
        chunk_extraction_set_payload["documents"] = [item for item in document_set_payload["documents"] if isinstance(item, dict)]
    if edges:
        chunk_extraction_set_payload["edges"] = edges
    chunk_extraction_set_ref = STORE.put_hot_subgraph(
        run_id=run_id,
        step_id=step_id,
        contract_type="chunk_extraction_set",
        payload=chunk_extraction_set_payload,
    )
    runtime_context["chunk_extraction_set_ref"] = chunk_extraction_set_ref


def _persist_runtime_hot_entity_relationship_set(
    *,
    implementation_next_name: str,
    run_id: str,
    step_id: str,
    runtime_context: dict[str, Any],
    step_outputs: dict[str, Any],
) -> None:
    if implementation_next_name != "map.conceptual.lift.entities_and_relationships":
        return
    entity_relationship_refs = _extract_runtime_output_fragment_ids(step_outputs)
    if not entity_relationship_refs:
        return
    source_subgraph_ref = runtime_context.get("chunk_extraction_set_ref")
    if not isinstance(source_subgraph_ref, str) or not source_subgraph_ref:
        return
    entity_relationship_set_payload = _entity_relationship_set_payload(
        source_subgraph_ref=source_subgraph_ref,
        subgraph_ref="",
        entity_relationship_refs=entity_relationship_refs,
    )
    entity_relationship_set_ref = STORE.put_hot_subgraph(
        run_id=run_id,
        step_id=step_id,
        contract_type="entity_relationship_set",
        payload=entity_relationship_set_payload,
    )
    runtime_context["entity_relationship_set_ref"] = entity_relationship_set_ref


def _persist_runtime_hot_claim_set(
    *,
    implementation_next_name: str,
    run_id: str,
    step_id: str,
    runtime_context: dict[str, Any],
    step_outputs: dict[str, Any],
) -> None:
    if implementation_next_name != "map.conceptual.lift.claims":
        return
    claim_refs = _extract_runtime_output_fragment_ids(step_outputs)
    if not claim_refs:
        return
    source_subgraph_ref = runtime_context.get("entity_relationship_set_ref")
    if not isinstance(source_subgraph_ref, str) or not source_subgraph_ref:
        return
    claim_set_payload = _claim_set_payload(
        source_subgraph_ref=source_subgraph_ref,
        subgraph_ref="",
        claim_refs=claim_refs,
    )
    claim_set_ref = STORE.put_hot_subgraph(
        run_id=run_id,
        step_id=step_id,
        contract_type="claim_set",
        payload=claim_set_payload,
    )
    runtime_context["claim_set_ref"] = claim_set_ref


def _hydrate_runtime_hot_inputs(
    *,
    implementation_next_name: str,
    runtime_context: dict[str, Any],
    outputs: dict[str, Any],
) -> dict[str, Any]:
    if implementation_next_name == "map.conceptual.lift.surface_fragments":
        document_set_ref = runtime_context.get("document_set_ref")
        if not isinstance(document_set_ref, str) or not document_set_ref:
            return outputs

        document_set_payload = STORE.get_hot_subgraph(document_set_ref)
        if not isinstance(document_set_payload, dict):
            return outputs

        document_fragment_refs = [
            str(item)
            for item in (document_set_payload.get("document_refs") or [])
            if isinstance(item, str) and item
        ]
        if not document_fragment_refs:
            return outputs

        hydrated_outputs = dict(outputs)
        hydrated_outputs["document_fragment_refs"] = document_fragment_refs
        existing_inputs = hydrated_outputs.get("inputs") if isinstance(hydrated_outputs.get("inputs"), dict) else {}
        hydrated_outputs["inputs"] = {
            **existing_inputs,
            "document_set_ref": document_set_ref,
            "document_fragment_refs": document_fragment_refs,
        }
        return hydrated_outputs

    if implementation_next_name == "map.conceptual.lift.entities_and_relationships":
        chunk_extraction_set_ref = runtime_context.get("chunk_extraction_set_ref")
        if not isinstance(chunk_extraction_set_ref, str) or not chunk_extraction_set_ref:
            return outputs

        chunk_extraction_set_payload = STORE.get_hot_subgraph(chunk_extraction_set_ref)
        if not isinstance(chunk_extraction_set_payload, dict):
            return outputs

        extraction_refs = [
            str(item)
            for item in (chunk_extraction_set_payload.get("extraction_refs") or [])
            if isinstance(item, str) and item
        ]
        if not extraction_refs:
            return outputs

        hydrated_outputs = dict(outputs)
        hydrated_outputs["fragment_ids"] = extraction_refs
        hydrated_outputs["chunk_extraction_set"] = dict(chunk_extraction_set_payload)
        existing_inputs = hydrated_outputs.get("inputs") if isinstance(hydrated_outputs.get("inputs"), dict) else {}
        hydrated_outputs["inputs"] = {
            **existing_inputs,
            "chunk_extraction_set_ref": chunk_extraction_set_ref,
            "fragment_ids": extraction_refs,
        }
        return hydrated_outputs

    if implementation_next_name == "map.conceptual.lift.claims":
        entity_relationship_set_ref = runtime_context.get("entity_relationship_set_ref")
        if not isinstance(entity_relationship_set_ref, str) or not entity_relationship_set_ref:
            return outputs

        entity_relationship_set_payload = STORE.get_hot_subgraph(entity_relationship_set_ref)
        if not isinstance(entity_relationship_set_payload, dict):
            return outputs

        entity_relationship_refs = [
            str(item)
            for item in (entity_relationship_set_payload.get("entity_relationship_refs") or [])
            if isinstance(item, str) and item
        ]
        if not entity_relationship_refs:
            return outputs

        hydrated_outputs = dict(outputs)
        hydrated_outputs["fragment_ids"] = entity_relationship_refs
        hydrated_outputs["entity_relationship_set"] = dict(entity_relationship_set_payload)
        existing_inputs = hydrated_outputs.get("inputs") if isinstance(hydrated_outputs.get("inputs"), dict) else {}
        hydrated_outputs["inputs"] = {
            **existing_inputs,
            "entity_relationship_set_ref": entity_relationship_set_ref,
            "fragment_ids": entity_relationship_refs,
        }
        return hydrated_outputs

    return outputs


def _validation_has_failures(validation: dict[str, Any] | None) -> bool:
    if not isinstance(validation, dict):
        return False
    results = validation.get("results")
    if not isinstance(results, list):
        return False
    return any(isinstance(item, dict) and item.get("status") != "passed" for item in results)


def _combine_phase_validations(
    input_validation: dict[str, Any] | None,
    output_validation: dict[str, Any] | None,
) -> dict[str, Any] | None:
    if input_validation is None and output_validation is None:
        return None
    return {
        "specs": [
            *(input_validation.get("specs", []) if isinstance(input_validation, dict) else []),
            *(output_validation.get("specs", []) if isinstance(output_validation, dict) else []),
        ],
        "resolved_inputs": input_validation.get("resolved_inputs", {}) if isinstance(input_validation, dict) else {},
        "resolved_outputs": output_validation.get("resolved_outputs", {}) if isinstance(output_validation, dict) else {},
        "results": [
            *(input_validation.get("results", []) if isinstance(input_validation, dict) else []),
            *(output_validation.get("results", []) if isinstance(output_validation, dict) else []),
        ],
    }


def _pending_validation_for_direction(validators: list[dict[str, Any]], direction: str) -> dict[str, Any] | None:
    specs = [validator for validator in validators if validator.get("direction") == direction]
    if not specs:
        return None
    return {
        "specs": specs,
        "resolved_inputs": {} if direction == "output" else {},
        "resolved_outputs": {} if direction == "input" else {},
        "results": [],
    }


def _split_transition_validation_by_direction(
    transition_validation: dict[str, Any] | None,
) -> tuple[dict[str, Any] | None, dict[str, Any] | None]:
    if not isinstance(transition_validation, dict):
        return None, None
    specs = transition_validation.get("specs") if isinstance(transition_validation.get("specs"), list) else []
    results = transition_validation.get("results") if isinstance(transition_validation.get("results"), list) else []
    resolved_inputs = transition_validation.get("resolved_inputs") if isinstance(transition_validation.get("resolved_inputs"), dict) else {}
    resolved_outputs = transition_validation.get("resolved_outputs") if isinstance(transition_validation.get("resolved_outputs"), dict) else {}
    input_specs = [item for item in specs if isinstance(item, dict) and item.get("direction") == "input"]
    output_specs = [item for item in specs if isinstance(item, dict) and item.get("direction") == "output"]
    input_names = {str(item.get("name")) for item in input_specs}
    output_names = {str(item.get("name")) for item in output_specs}
    input_results = [item for item in results if isinstance(item, dict) and str(item.get("name")) in input_names]
    output_results = [item for item in results if isinstance(item, dict) and str(item.get("name")) in output_names]
    input_validation = {
        "specs": input_specs,
        "resolved_inputs": resolved_inputs,
        "resolved_outputs": {},
        "results": input_results,
    } if input_specs or input_results or resolved_inputs else None
    output_validation = {
        "specs": output_specs,
        "resolved_inputs": {},
        "resolved_outputs": resolved_outputs,
        "results": output_results,
    } if output_specs or output_results or resolved_outputs else None
    return input_validation, output_validation


def _validation_is_effectively_empty(validation: dict[str, Any] | None, *, direction: str) -> bool:
    if not isinstance(validation, dict):
        return True
    key = "resolved_inputs" if direction == "input" else "resolved_outputs"
    resolved = validation.get(key)
    if not isinstance(resolved, dict):
        return True
    return all(not isinstance(value, list) or len(value) == 0 for value in resolved.values())


def _validation_has_unresolved_refs(validation: dict[str, Any] | None, *, direction: str) -> bool:
    if not isinstance(validation, dict):
        return False
    key = "resolved_inputs" if direction == "input" else "resolved_outputs"
    resolved = validation.get(key)
    if not isinstance(resolved, dict):
        return False
    for values in resolved.values():
        if not isinstance(values, list):
            continue
        for item in values:
            if not isinstance(item, dict):
                continue
            inspection = item.get("inspection")
            if not isinstance(inspection, dict):
                continue
            refs = inspection.get("refs") if isinstance(inspection.get("refs"), list) else []
            resolved_refs = inspection.get("resolved_refs") if isinstance(inspection.get("resolved_refs"), list) else []
            if refs and not resolved_refs:
                return True
    return False


def _normalize_step_outputs_for_output_validation(
    *,
    step_outputs: dict[str, Any],
    pre_execution_outputs: dict[str, Any],
) -> dict[str, Any]:
    normalized = deepcopy(step_outputs)
    current_fragment_ids = _extract_runtime_output_fragment_ids(normalized)
    prior_fragment_ids = _extract_runtime_output_fragment_ids(pre_execution_outputs)
    if current_fragment_ids and current_fragment_ids == prior_fragment_ids:
        normalized.pop("fragment_ids", None)
    current_document_refs = [str(item) for item in (normalized.get("document_fragment_refs") or []) if isinstance(item, str) and item]
    prior_document_refs = [str(item) for item in (pre_execution_outputs.get("document_fragment_refs") or []) if isinstance(item, str) and item]
    if current_document_refs and current_document_refs == prior_document_refs:
        normalized.pop("document_fragment_refs", None)
    return normalized


def _extract_runtime_output_fragment_ids(step_outputs: dict[str, Any]) -> list[str]:
    direct = [str(item) for item in (step_outputs.get("fragment_ids") or []) if isinstance(item, str) and item]
    if direct:
        return direct

    decomposition = step_outputs.get("decomposition")
    structural = None
    if decomposition is not None:
        structural = getattr(decomposition, "structural", None)
        if structural is None and isinstance(decomposition, dict):
            structural = decomposition.get("structural")
    derived = [
        str(getattr(item, "cas_id", None) or getattr(item, "id", None) or "")
        for item in (structural or [])
    ]
    derived = [item for item in derived if item]
    if derived:
        return derived

    candidates = step_outputs.get("map_segment_candidates") if isinstance(step_outputs.get("map_segment_candidates"), list) else []
    derived = [str(item.get("segment_id") or "") for item in candidates if isinstance(item, dict)]
    return [item for item in derived if item]


def _implementation_step_name(pipeline_id: str | None, step_name: str) -> str:
    if pipeline_id != "ingestion-early-parse":
        return step_name

    aliases = {
        "load.documents": "parse_artifacts",
        "parse.chunk": "map.conceptual.lift.surface_fragments",
        "parse.entities_and_relationships": "map.conceptual.lift.entities_and_relationships",
        "parse.claims": "map.conceptual.lift.claims",
        "complete": "map.conceptual.lift.summarize",
    }
    return aliases.get(step_name, step_name)


def _presentation_step_name(pipeline_id: str | None, step_name: str) -> str:
    if pipeline_id != "ingestion-early-parse":
        return step_name

    aliases = {
        "parse_artifacts": "load.documents",
        "map.conceptual.lift.surface_fragments": "parse.chunk",
        "map.conceptual.lift.entities_and_relationships": "parse.entities_and_relationships",
        "map.conceptual.lift.claims": "parse.claims",
        "map.conceptual.lift.summarize": "complete",
    }
    return aliases.get(step_name, step_name)


def _executor_identity_for_step(pipeline_id: str | None, step_name: str) -> dict[str, str]:
    if pipeline_id:
        compiled = _compiled_operator_metadata(pipeline_id).get(step_name) or {}
        executor_id = str(compiled.get("executor_id") or "").strip()
        executor_kind = str(compiled.get("executor_kind") or "").strip()
        if executor_id or executor_kind:
            return {
                "executor_id": executor_id,
                "executor_kind": executor_kind,
            }

    implementation_name = _implementation_step_name(pipeline_id, step_name)
    aliases = {
        "parse_artifacts": {
            "executor_id": "executor://python-primary",
            "executor_kind": "python-executor",
        },
        "map.conceptual.lift.surface_fragments": {
            "executor_id": "executor://python-primary",
            "executor_kind": "python-executor",
        },
        "map.conceptual.lift.entities_and_relationships": {
            "executor_id": "executor://ml-primary",
            "executor_kind": "ml-executor",
        },
        "map.conceptual.lift.claims": {
            "executor_id": "executor://ml-primary",
            "executor_kind": "ml-executor",
        },
        "init.initialize": {
            "executor_id": "executor://debug-local",
            "executor_kind": "debug-executor",
        },
        "map.conceptual.lift.summarize": {
            "executor_id": "executor://debug-local",
            "executor_kind": "debug-executor",
        },
    }
    return dict(aliases.get(implementation_name, {}))


def _event_executor_identity(event: DebugStepEvent) -> dict[str, str]:
    metrics = event.metrics if isinstance(event.metrics, dict) else {}
    trace = metrics.get("trace") if isinstance(metrics.get("trace"), dict) else {}
    executor_id = str(trace.get("executor_id") or "").strip()
    executor_kind = str(trace.get("executor_kind") or "").strip()
    if executor_id or executor_kind:
        return {
            "executor_id": executor_id,
            "executor_kind": executor_kind,
        }
    return _executor_identity_for_step(event.pipeline_id, event.step_name)


def _capture_step_logs(step_name: str, metrics: dict[str, Any], *, started_at: str, ended_at: str) -> None:
    raw_log_events = metrics.get("log_events") if isinstance(metrics.get("log_events"), list) else []
    log_events = [dict(item) for item in raw_log_events if isinstance(item, dict)]

    trace = metrics.get("trace") if isinstance(metrics.get("trace"), dict) else {}
    timeline = _serialize_trace_timeline(trace.get("timeline"))
    raw_events = _serialize_raw_events(trace.get("raw_events"))
    if not timeline and trace:
        topic_sequence = _serialize_trace_topic_sequence(trace.get("topic_sequence"))
        timeline = [
            {
                "topic": item["topic"],
                "event_type": item["event_type"],
                "status": item["status"],
                "occurred_at": started_at if index == 0 else ended_at,
                "payload": {},
            }
            for index, item in enumerate(topic_sequence)
        ]
    if not raw_events and timeline:
        raw_events = [{"topic": item["topic"], "payload": item.get("payload", {})} for item in timeline]

    metrics["executor_logs"] = _named_logs_from_log_events(log_events, source="executor")
    metrics["system_logs"] = _named_logs_from_log_events(log_events, source="system")
    metrics["logs"] = _legacy_logs_from_log_events(log_events)
    metrics["log_events"] = log_events
    metrics["trace"] = {
        **trace,
        "timeline": timeline,
        "raw_events": raw_events,
    }


class _PerfReportDebugExecutionScope:
    def __init__(self, base_scope: DefaultExecutionScope | None = None, pipeline_id: str | None = None) -> None:
        self._base_scope = base_scope or DefaultExecutionScope()
        self._pipeline_id = pipeline_id

    def get_dynamic_execution_steps(self) -> list[str]:
        return _debug_pipeline_steps(self._pipeline_id, self._base_scope)

    def get_available_tools(self) -> list[dict[str, Any]]:
        return self._base_scope.get_available_tools()

    def get_step_execution_metadata(self, step_name: str) -> dict[str, Any]:
        if not self._pipeline_id:
            return {}
        return dict(_compiled_operator_metadata(self._pipeline_id).get(step_name) or {})


def _serialize_debug_event(event: DebugStepEvent) -> dict[str, Any]:
    payload = dict(event.__dict__)
    payload.pop("env_type", None)
    payload.pop("env_id", None)
    payload["step_name"] = str(event.step_name)
    payload["ref"] = _legacy_scope_ref(event.env_type, event.env_id)
    payload.update(_event_executor_identity(event))
    return payload


def _serialize_debug_state(state: DebugRunState) -> dict[str, Any]:
    payload = dict(state.__dict__)
    payload.pop("env_type", None)
    payload.pop("env_id", None)
    payload["current_step_name"] = str(state.current_step_name)
    payload["ref"] = _legacy_scope_ref(state.env_type, state.env_id)
    return payload


def _run_scope_ref_from_state(state: DebugRunState) -> str:
    return _legacy_scope_ref(state.env_type, state.env_id)


def _legacy_scope_ref(env_type: str, env_id: str) -> str:
    normalized_type = env_type.strip()
    normalized_id = env_id.strip()
    if normalized_type == "committed":
        return EnvironmentScope(ref="refs/heads/main").ref
    if normalized_type == "staging":
        return EnvironmentScope(ref=f"refs/heads/staging/{normalized_id}").ref
    if normalized_type == "dev":
        return EnvironmentScope(ref=f"refs/heads/run/{normalized_id}").ref
    raise HTTPException(status_code=422, detail={"status": "schema_invalid", "missing": ["ref"]})


def _resolve_scope_query(*, ref: str | None = None, env_type: str | None = None, env_id: str | None = None) -> EnvironmentScope:
    if isinstance(ref, str) and ref.strip():
        return EnvironmentScope(ref=ref.strip())
    if isinstance(env_type, str) and env_type.strip() and isinstance(env_id, str) and env_id.strip():
        return EnvironmentScope(ref=_legacy_scope_ref(env_type, env_id))
    raise HTTPException(status_code=422, detail={"status": "schema_invalid", "missing": ["ref"]})


def _inspection_request(ref: str, max_depth: int) -> tuple[InspectionRef, ResolveInspectionRequest]:
    try:
        parsed_ref = InspectionRef.parse(ref)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail={"status": "schema_invalid", "reason": str(exc)}) from exc

    return parsed_ref, ResolveInspectionRequest(
        inspection_ref=parsed_ref,
        max_depth=max_depth,
        include_payload=True,
        include_edges=True,
        include_provenance=True,
    )


def _resolve_hot_inspection(run_id: str, parsed_ref: InspectionRef, request: ResolveInspectionRequest):
    subgraph_ref = parsed_ref.locator.get("subgraph_ref")
    if not isinstance(subgraph_ref, str) or not subgraph_ref:
        return None
    refs_by_name = getattr(STORE, "_hot_subgraph_refs_by_run", {}).get(run_id, {})
    hot_ref = refs_by_name.get(subgraph_ref)
    if hot_ref is None:
        return None
    resolver = HotInspectionResolver(getattr(STORE, "_hot_subgraph_store"))
    return resolver.resolve(
        request.model_copy(
            update={
                "inspection_ref": InspectionRef(
                    backend="hot",
                    locator={
                        "subgraph_ref": subgraph_ref,
                        "head_fragment_id": hot_ref["head_fragment_id"],
                    },
                )
            }
        )
    )


def _resolve_persistent_inspection(run_id: str, parsed_ref: InspectionRef, request: ResolveInspectionRequest):
    runtime_context = STORE.get_debug_runtime_context(run_id) or {}
    persistent_state = runtime_context.get("persistent_graph_state")
    if not isinstance(persistent_state, PersistentGraphState):
        return None
    subgraph_ref = parsed_ref.locator.get("subgraph_ref")
    if not isinstance(subgraph_ref, str) or not subgraph_ref:
        return None
    resolver = PersistentInspectionResolver(persistent_state)
    return resolver.resolve(
        request.model_copy(
            update={
                "inspection_ref": InspectionRef(
                    backend="persistent",
                    locator={"subgraph_ref": subgraph_ref},
                )
            }
        )
    )


def _resolve_fragment_inspection(run_id: str, parsed_ref: InspectionRef, request: ResolveInspectionRequest):
    fragment_id = parsed_ref.locator.get("cas_id")
    if not isinstance(fragment_id, str) or not fragment_id:
        return None
    state = STORE.get_debug_run_state(run_id)
    if state is None:
        return None
    environment_fragments = STORE.list_environment_fragments(
        run_id=run_id,
        env_type=state.env_type,
        env_id=state.env_id,
    )
    payload = next(
        (
            item for item in environment_fragments
            if isinstance(item, dict) and str(item.get("id") or item.get("cas_id") or "") == fragment_id
        ),
        None,
    )
    hot_fragment_payload = None
    hot_container_payloads: list[dict[str, Any]] = []
    refs_by_name = getattr(STORE, "_hot_subgraph_refs_by_run", {}).get(run_id, {})
    for subgraph_ref in refs_by_name:
        hot_payload = STORE.get_hot_subgraph(subgraph_ref)
        if not isinstance(hot_payload, dict):
            continue
        matched_in_payload = False
        documents = hot_payload.get("documents") if isinstance(hot_payload.get("documents"), list) else []
        hot_fragment_payload = next(
            (
                {
                    "id": fragment_id,
                    "cas_id": item.get("cas_id") or fragment_id,
                    "mime_type": item.get("mime_type"),
                    "value": item.get("value"),
                    "_hot_subgraph_payload": hot_payload,
                }
                for item in documents
                if isinstance(item, dict) and str(item.get("cas_id") or item.get("id") or "") == fragment_id
            ),
            hot_fragment_payload,
        )
        matched_in_payload = hot_fragment_payload is not None and isinstance(hot_fragment_payload.get("_hot_subgraph_payload"), dict) and hot_fragment_payload.get("_hot_subgraph_payload") is hot_payload
        extractions = hot_payload.get("extractions") if isinstance(hot_payload.get("extractions"), list) else []
        if not matched_in_payload:
            hot_fragment_payload = next(
                (
                    {
                        "id": fragment_id,
                        "cas_id": item.get("cas_id") or fragment_id,
                        "mime_type": item.get("mime_type"),
                        "value": item.get("value"),
                        "_hot_subgraph_payload": hot_payload,
                    }
                    for item in extractions
                    if isinstance(item, dict) and str(item.get("cas_id") or item.get("id") or "") == fragment_id
                ),
                hot_fragment_payload,
            )
            matched_in_payload = hot_fragment_payload is not None and isinstance(hot_fragment_payload.get("_hot_subgraph_payload"), dict) and hot_fragment_payload.get("_hot_subgraph_payload") is hot_payload
        document_chunk_sets = hot_payload.get("document_chunk_sets") if isinstance(hot_payload.get("document_chunk_sets"), list) else []
        if not matched_in_payload:
            hot_fragment_payload = next(
                (
                    {
                        "id": fragment_id,
                        "cas_id": item.get("cas_id") or fragment_id,
                        "mime_type": item.get("mime_type"),
                        "value": item.get("value"),
                        "_hot_subgraph_payload": hot_payload,
                    }
                    for item in document_chunk_sets
                    if isinstance(item, dict) and str(item.get("cas_id") or item.get("id") or "") == fragment_id
                ),
                hot_fragment_payload,
            )
            matched_in_payload = hot_fragment_payload is not None and isinstance(hot_fragment_payload.get("_hot_subgraph_payload"), dict) and hot_fragment_payload.get("_hot_subgraph_payload") is hot_payload
        if matched_in_payload:
            container_payload = dict(hot_payload)
            if not str(container_payload.get("subgraph_ref") or ""):
                container_payload["subgraph_ref"] = subgraph_ref
            hot_container_payloads.append(container_payload)
    if payload is not None and isinstance(payload, dict) and isinstance(hot_fragment_payload, dict):
        merged_payload = dict(payload)
        if isinstance(payload.get("value"), dict) and isinstance(hot_fragment_payload.get("value"), dict):
            merged_payload["value"] = {**hot_fragment_payload["value"], **payload["value"]}
        elif payload.get("value") is None and hot_fragment_payload.get("value") is not None:
            merged_payload["value"] = hot_fragment_payload.get("value")
        merged_payload["_hot_subgraph_payload"] = hot_fragment_payload.get("_hot_subgraph_payload")
        payload = merged_payload
    if payload is None:
        payload = hot_fragment_payload
    if payload is None:
        return None
    value = payload.get("value") if isinstance(payload.get("value"), dict) else payload.get("value")
    label = fragment_id
    if isinstance(value, dict):
        for candidate in (value.get("name"), value.get("filename"), value.get("file_name"), value.get("document_id")):
            if isinstance(candidate, str) and candidate:
                label = candidate
                break
    node_id = node_id_for("fragment", {"cas_id": fragment_id})
    nodes = [
        InspectionNode(
            id=node_id,
            kind="fragment",
            label=label,
            payload={
                "cas_id": payload.get("cas_id") or fragment_id,
                "mime_type": payload.get("mime_type"),
                "value": value,
            } if request.include_payload else {},
            refs={"self": InspectionRef(backend="hot", locator={"cas_id": fragment_id})},
            provenance={"source_backend": "hot", "run_id": run_id},
        )
    ]
    edges: list[InspectionEdge] = []
    source_document_fragment_id = str(value.get("source_document_fragment_id") or "") if isinstance(value, dict) else ""
    hot_subgraph_payload = payload.get("_hot_subgraph_payload") if isinstance(payload.get("_hot_subgraph_payload"), dict) else None
    if source_document_fragment_id and hot_subgraph_payload is not None:
        documents = hot_subgraph_payload.get("documents") if isinstance(hot_subgraph_payload.get("documents"), list) else []
        source_document_payload = next(
            (
                item
                for item in documents
                if isinstance(item, dict) and str(item.get("cas_id") or item.get("id") or "") == source_document_fragment_id
            ),
            None,
        )
        if isinstance(source_document_payload, dict):
            source_node_id = node_id_for("fragment", {"cas_id": source_document_fragment_id})
            source_label = source_document_fragment_id
            source_value = source_document_payload.get("value") if isinstance(source_document_payload.get("value"), dict) else source_document_payload.get("value")
            if isinstance(source_value, dict):
                for candidate in (source_value.get("name"), source_value.get("filename"), source_value.get("file_name"), source_value.get("document_id")):
                    if isinstance(candidate, str) and candidate:
                        source_label = candidate
                        break
            nodes.append(
                InspectionNode(
                    id=source_node_id,
                    kind="fragment",
                    label=source_label,
                    payload={
                        "cas_id": source_document_payload.get("cas_id") or source_document_fragment_id,
                        "mime_type": source_document_payload.get("mime_type"),
                        "value": source_value,
                    } if request.include_payload else {},
                    refs={"self": InspectionRef(backend="hot", locator={"cas_id": source_document_fragment_id})},
                    provenance={"source_backend": "hot", "run_id": run_id},
                )
            )
            if request.include_edges:
                edges.append(
                    InspectionEdge(
                        id=edge_id_for("derives", node_id, source_node_id),
                        **{"from": node_id, "to": source_node_id},
                        relation="derives",
                        provenance={"source_backend": "hot", "run_id": run_id},
                    )
                )
    fragment_document_id = str(value.get("document_id") or "") if isinstance(value, dict) else ""
    related_document_chunk_sets: list[tuple[dict[str, Any], dict[str, Any]]] = []
    for container_payload in hot_container_payloads:
        container_subgraph_ref = str(container_payload.get("subgraph_ref") or "")
        if container_subgraph_ref:
            container_node_id = node_id_for("subgraph", {"subgraph_ref": container_subgraph_ref})
            if not any(node.id == container_node_id for node in nodes):
                nodes.append(
                    InspectionNode(
                        id=container_node_id,
                        kind="subgraph",
                        label=container_subgraph_ref,
                        payload=container_payload if request.include_payload else {},
                        refs={"self": InspectionRef(backend="hot", locator={"subgraph_ref": container_subgraph_ref})},
                        provenance={"source_backend": "hot", "run_id": run_id},
                    )
                )
            if request.include_edges:
                edges.append(
                    InspectionEdge(
                        id=edge_id_for("contains", container_node_id, node_id),
                        **{"from": container_node_id, "to": node_id},
                        relation="contains",
                        provenance={"source_backend": "hot", "run_id": run_id},
                    )
                )
        if str(container_payload.get("kind") or "") != "chunk_extraction_set":
            continue
        for document_chunk_set in container_payload.get("document_chunk_sets") or []:
            if not isinstance(document_chunk_set, dict):
                continue
            document_chunk_set_id = str(document_chunk_set.get("cas_id") or document_chunk_set.get("id") or "")
            document_chunk_set_value = document_chunk_set.get("value") if isinstance(document_chunk_set.get("value"), dict) else {}
            if not document_chunk_set_id:
                continue
            if fragment_id in (document_chunk_set_value.get("chunk_refs") or []):
                related_document_chunk_sets.append((container_payload, document_chunk_set))
                continue
            if str(document_chunk_set_value.get("source_document_fragment_id") or "") == fragment_id:
                related_document_chunk_sets.append((container_payload, document_chunk_set))
                continue
            if fragment_document_id and str(document_chunk_set_value.get("document_id") or "") == fragment_document_id:
                related_document_chunk_sets.append((container_payload, document_chunk_set))
    for container_payload, document_chunk_set in related_document_chunk_sets:
        document_chunk_set_id = str(document_chunk_set.get("cas_id") or document_chunk_set.get("id") or "")
        document_chunk_set_value = document_chunk_set.get("value") if isinstance(document_chunk_set.get("value"), dict) else document_chunk_set.get("value")
        document_chunk_set_node_id = node_id_for("fragment", {"cas_id": document_chunk_set_id})
        document_chunk_set_label = document_chunk_set_id
        if isinstance(document_chunk_set_value, dict):
            for candidate in (document_chunk_set_value.get("document_id"), document_chunk_set_value.get("filename"), document_chunk_set_value.get("file_name")):
                if isinstance(candidate, str) and candidate:
                    document_chunk_set_label = candidate
                    break
        if not any(node.id == document_chunk_set_node_id for node in nodes):
            nodes.append(
                InspectionNode(
                    id=document_chunk_set_node_id,
                    kind="fragment",
                    label=document_chunk_set_label,
                    payload={
                        "cas_id": document_chunk_set_id,
                        "mime_type": document_chunk_set.get("mime_type"),
                        "value": document_chunk_set_value,
                    } if request.include_payload else {},
                    refs={"self": InspectionRef(backend="hot", locator={"cas_id": document_chunk_set_id})},
                    provenance={"source_backend": "hot", "run_id": run_id},
                )
            )
        if request.include_edges:
            container_subgraph_ref = str(container_payload.get("subgraph_ref") or "")
            if container_subgraph_ref:
                container_node_id = node_id_for("subgraph", {"subgraph_ref": container_subgraph_ref})
                edges.append(
                    InspectionEdge(
                        id=edge_id_for("contains", container_node_id, document_chunk_set_node_id),
                        **{"from": container_node_id, "to": document_chunk_set_node_id},
                        relation="contains",
                        provenance={"source_backend": "hot", "run_id": run_id},
                    )
                )
            relation = "contains" if fragment_id in ((document_chunk_set_value or {}).get("chunk_refs") or []) else "references"
            edges.append(
                InspectionEdge(
                    id=edge_id_for(relation, document_chunk_set_node_id, node_id),
                    **{"from": document_chunk_set_node_id, "to": node_id},
                    relation=relation,
                    provenance={"source_backend": "hot", "run_id": run_id},
                )
            )
    return InspectionSubgraph(
        schema_version="v1",
        root_node_id=node_id,
        nodes=nodes,
        edges=edges if request.include_edges else [],
        navigation={"focus": {"node_id": node_id}},
    )


def _resolve_inspection_subgraph(run_id: str, ref: str, max_depth: int):
    parsed_ref, request = _inspection_request(ref, max_depth)
    fragment_subgraph = _resolve_fragment_inspection(run_id, parsed_ref, request)
    if fragment_subgraph is not None:
        return fragment_subgraph
    hot_subgraph = _resolve_hot_inspection(run_id, parsed_ref, request)
    if hot_subgraph is not None:
        return hot_subgraph
    persistent_subgraph = _resolve_persistent_inspection(run_id, parsed_ref, request)
    if persistent_subgraph is not None:
        return persistent_subgraph
    raise HTTPException(status_code=404, detail={"status": "not_found", "ref": ref})


def _artifact_label(file_name: str) -> str:
    stem = file_name.rsplit("/", 1)[-1]
    stem = stem.rsplit(".", 1)[0] if "." in stem else stem
    return stem.replace("_", " ").replace("-", " ")


def _sync_graph_snapshot_from_step_outputs(*, run_id: str, step_name: str, step_outputs: dict[str, Any], runtime_context: dict[str, Any]) -> None:
    """Materialize graph snapshot state from async pipeline outputs."""
    if step_name not in ("map.conceptual.lift.surface_fragments", "map.conceptual.project.surface_fragments", "lift_fragments", "map.reconstructable.build_subgraph.reconstruction"):
        return

    run = STORE.get_run(run_id)
    if run is None:
        return

    if "graph_projection" in step_outputs:
        graph = step_outputs.get("graph_projection", {})
        print(f"DEBUG SYNC: synced graph_projection with {len(graph.get('nodes', []))} nodes and {len(graph.get('edges', []))} edges for step {step_name}.")
        if isinstance(graph, dict):
            run.graph.nodes = graph.get("nodes", [])
            run.graph.edges = graph.get("edges", [])
        return

    if step_name == "map.reconstructable.build_subgraph.reconstruction":
        return

    decomposition = step_outputs.get("decomposition")
    root_fragments = list(getattr(decomposition, "root_fragments", []) or [])
    if not root_fragments:
        return

    fragment_artifact_map = step_outputs.get("fragment_artifact_map")
    if not isinstance(fragment_artifact_map, dict):
        fragment_artifact_map = {}

    manifest = runtime_context.get("asset_manifest")
    if not isinstance(manifest, list):
        manifest = []

    artifact_label_by_id: dict[str, str] = {}
    filenames_by_label: dict[str, str] = {}
    for item in manifest:
        if not isinstance(item, dict):
            continue
        artifact_id = str(item.get("artifact_id") or "")
        file_name = str(item.get("filename") or "")
        if not artifact_id or not file_name:
            continue
        label = _artifact_label(file_name)
        artifact_label_by_id[artifact_id] = label
        filenames_by_label[label] = file_name

    labels_by_fragment_id: dict[str, str] = {}
    root_fragments_by_id: dict[str, Any] = {}
    for fragment in root_fragments:
        fragment_id = getattr(fragment, "cas_id", None) or getattr(fragment, "id", None)
        if not fragment_id:
            continue
        root_fragments_by_id[str(fragment_id)] = fragment
        artifact_id = fragment_artifact_map.get(str(fragment_id))
        label = artifact_label_by_id.get(str(artifact_id or ""))
        if label:
            labels_by_fragment_id[str(fragment_id)] = label

    resolved_artifact_by_fragment_id: dict[str, str] = {
        fragment_id: artifact_id
        for fragment_id, artifact_id in fragment_artifact_map.items()
        if isinstance(fragment_id, str) and isinstance(artifact_id, str)
    }

    changed = True
    while changed:
        changed = False
        for fragment_id, fragment in root_fragments_by_id.items():
            if not is_relation_fragment(fragment):
                continue
            value = getattr(fragment, "value", None)
            if not isinstance(value, dict):
                continue
            try:
                relation = Relation.model_validate(value)
            except Exception:
                continue

            referenced_ids = [slot.fragment_id for group in relation.binding_groups for slot in group.slots]
            known_ids = [fragment_id, *referenced_ids]
            artifact_id = next((resolved_artifact_by_fragment_id.get(fid) for fid in known_ids if resolved_artifact_by_fragment_id.get(fid)), None)
            if not artifact_id:
                continue
            for fid in known_ids:
                if fid and fid not in resolved_artifact_by_fragment_id:
                    resolved_artifact_by_fragment_id[fid] = artifact_id
                    changed = True

    for fragment_id, artifact_id in resolved_artifact_by_fragment_id.items():
        label = artifact_label_by_id.get(artifact_id)
        if label:
            labels_by_fragment_id[fragment_id] = label

    nodes, edges, manifests = fragments_to_graph(
        root_fragments,
        labels_by_fragment_id=labels_by_fragment_id,
        filenames_by_label=filenames_by_label,
    )
    run.graph.fragments = root_fragments
    run.graph.nodes = nodes
    run.graph.edges = edges
    run.graph.manifests = manifests


class ReviewPayload(BaseModel):
    query_id: str = Field(min_length=1)
    relevance: float = Field(ge=0.0, le=1.0)
    fidelity: float = Field(ge=0.0, le=1.0)
    clarity: float = Field(ge=0.0, le=1.0)
    note: str = ""


class RunControlPayload(BaseModel):
    command_id: str = Field(min_length=1)
    action: str = Field(pattern="^(set_mode|pause|resume|next_step|inject_verify_fail)$")
    pipeline_id: str = Field(min_length=1)
    pipeline_run_id: str = Field(min_length=1)
    mode: str | None = Field(default=None, pattern="^(autonomous|manual)$")
    drift_at: str | None = Field(default=None)


class SeedScenarioPayload(BaseModel):
    scenario_key: str = Field(min_length=1)
    overrides: dict[str, str] = Field(default_factory=dict)


def _test_gate_enabled() -> bool:
    return (
        os.getenv("IKAM_PERF_REPORT_TEST_MODE", "0").strip() == "1"
        and os.getenv("IKAM_ALLOW_DEBUG_INJECTION", "0").strip() == "1"
    )


def _require_test_gate() -> None:
    if not _test_gate_enabled():
        raise HTTPException(status_code=403, detail={"status": "forbidden", "reason": "test_mode_gate_disabled"})


def _async_next_step_enabled() -> bool:
    return _test_gate_enabled() and os.getenv("IKAM_ASYNC_NEXT_STEP", "0").strip() == "1"


def _control_availability(*, state: DebugRunState, events: list[DebugStepEvent], pipeline_steps: list[str]) -> dict[str, bool]:
    last_event = events[-1] if events else None
    previous_step_completed = isinstance(last_event, DebugStepEvent) and last_event.status == "succeeded"
    paused = state.execution_state == "paused"
    terminal_step = pipeline_steps[-1] if pipeline_steps else ""
    not_terminal = state.current_step_name != terminal_step
    can_resume = paused and previous_step_completed and not_terminal
    can_next_step = state.execution_mode == "manual" and paused and previous_step_completed and not_terminal
    return {
        "can_resume": can_resume,
        "can_next_step": can_next_step,
    }


async def _execute_next_pipeline_step(*, run_id: str, state: DebugRunState) -> DebugStepEvent:
    prepared_step, terminal_event = _prepare_next_pipeline_step(run_id=run_id, state=state)
    if terminal_event is not None:
        return terminal_event
    if prepared_step is None:
        raise RuntimeError("Expected prepared step or terminal event")
    return await _execute_prepared_pipeline_step(run_id=run_id, state=state, prepared_step=prepared_step)


def _prepare_next_pipeline_step(*, run_id: str, state: DebugRunState) -> tuple[_PreparedPipelineStep | None, DebugStepEvent | None]:
    events = STORE.list_debug_events(run_id)
    debug_scope = _PerfReportDebugExecutionScope(pipeline_id=state.pipeline_id)
    retry_parent_step_id = None

    latest_current_step_event = None
    for event in reversed(events):
        if event.step_name == state.current_step_name:
            latest_current_step_event = event
            break

    if latest_current_step_event and latest_current_step_event.status == "failed":
        retry_parent_step_id = latest_current_step_event.step_id
        state.retry_budget_remaining -= 1
        if state.retry_budget_remaining <= 0:
            state.execution_state = "budget_exhausted"
            terminal = DebugStepEvent(
                event_id=f"ev-{uuid4().hex}",
                run_id=state.run_id,
                pipeline_id=state.pipeline_id,
                pipeline_run_id=state.pipeline_run_id,
                project_id=state.project_id,
                operation_id=state.operation_id,
                env_type=state.env_type,
                env_id=state.env_id,
                step_name=state.current_step_name,
                step_id=f"step-{uuid4().hex[:10]}",
                status="failed",
                attempt_index=state.current_attempt_index,
                retry_parent_step_id=retry_parent_step_id,
                started_at=datetime.now(UTC).isoformat(),
                ended_at=datetime.now(UTC).isoformat(),
                duration_ms=0,
                metrics={"control": "next_step", "retry_boundary": True, "budget_exhausted": True},
                error={"reason": "retry budget exhausted"},
            )
            STORE.append_debug_event(terminal)
            return None, terminal
        state.current_attempt_index += 1
        drift_at = (latest_current_step_event.error or {}).get("drift_at")
        if drift_at:
            next_name = drift_at
        else:
            dynamic_steps = debug_scope.get_dynamic_execution_steps()
            next_name = dynamic_steps[1] if len(dynamic_steps) > 1 else dynamic_steps[0]
    else:
        next_name = next_step_name(state.current_step_name, scope=debug_scope)

    implementation_next_name = _implementation_step_name(state.pipeline_id, next_name)

    runtime_context = STORE.get_debug_runtime_context(run_id)
    if runtime_context is None:
        raise RuntimeError("Missing debug runtime context")
    source_bytes = runtime_context.get("source_bytes")
    mime_type = runtime_context.get("mime_type")
    artifact_id = runtime_context.get("artifact_id")
    if not isinstance(source_bytes, (bytes, bytearray)):
        raise RuntimeError("Runtime context source_bytes must be bytes")
    if not isinstance(mime_type, str) or not mime_type:
        raise RuntimeError("Runtime context missing mime_type")
    if not isinstance(artifact_id, str) or not artifact_id:
        raise RuntimeError("Runtime context missing artifact_id")

    execution_state = StepExecutionState(
        source_bytes=bytes(source_bytes),
        mime_type=mime_type,
        artifact_id=artifact_id,
        assets=list(runtime_context.get("asset_payloads") or []),
        outputs=_hydrate_runtime_hot_inputs(
            implementation_next_name=implementation_next_name,
            runtime_context=runtime_context,
            outputs=dict(runtime_context.get("step_outputs") or {}),
        ),
    )

    start_perf = perf_counter()
    start_iso = datetime.now(UTC).isoformat()
    running_event = DebugStepEvent(
        event_id=f"ev-{uuid4().hex}",
        run_id=state.run_id,
        pipeline_id=state.pipeline_id,
        pipeline_run_id=state.pipeline_run_id,
        project_id=state.project_id,
        operation_id=state.operation_id,
        env_type=state.env_type,
        env_id=state.env_id,
        step_name=next_name,
        step_id=f"step-{uuid4().hex[:10]}",
        status="running",
        attempt_index=state.current_attempt_index,
        retry_parent_step_id=retry_parent_step_id,
        started_at=start_iso,
        ended_at=None,
        duration_ms=None,
        metrics={
            "control": "next_step",
            "retry_boundary": retry_parent_step_id is not None,
            "executor_logs": {"stdout_lines": [], "stderr_lines": []},
            "system_logs": {"stdout_lines": [], "stderr_lines": []},
            "logs": {"stdout_lines": [], "stderr_lines": []},
            "log_events": [],
            "trace": {"timeline": [], "raw_events": []},
        },
        error=None,
    )
    STORE.append_debug_event(running_event)
    state.current_step_name = next_name
    executor_stdout_lines: list[str] = []
    executor_stderr_lines: list[str] = []
    system_stdout_lines: list[str] = [f"executing {next_name} operation"]
    system_stderr_lines: list[str] = []
    log_events: list[dict[str, Any]] = []
    _append_log_event(
        log_events,
        at=start_iso,
        source="system",
        stream="stdout",
        message=f"executing {next_name} operation",
    )
    started_message = f"[{next_name}] step started at {start_iso}"
    system_stdout_lines.append(started_message)
    _append_log_event(
        log_events,
        at=start_iso,
        source="system",
        stream="stdout",
        message=started_message,
    )
    _publish_running_event_logs(
        run_id=state.run_id,
        step_id=running_event.step_id,
        executor_stdout_lines=executor_stdout_lines,
        executor_stderr_lines=executor_stderr_lines,
        system_stdout_lines=system_stdout_lines,
        system_stderr_lines=system_stderr_lines,
        log_events=log_events,
    )

    return (
        _PreparedPipelineStep(
            debug_scope=debug_scope,
            implementation_next_name=implementation_next_name,
            runtime_context=runtime_context,
            execution_state=execution_state,
            running_event=running_event,
            retry_parent_step_id=retry_parent_step_id,
            start_iso=start_iso,
            start_perf=start_perf,
            executor_stdout_lines=executor_stdout_lines,
            executor_stderr_lines=executor_stderr_lines,
            system_stdout_lines=system_stdout_lines,
            system_stderr_lines=system_stderr_lines,
            log_events=log_events,
        ),
        None,
    )


async def _execute_prepared_pipeline_step(
    *,
    run_id: str,
    state: DebugRunState,
    prepared_step: _PreparedPipelineStep,
) -> DebugStepEvent:
    runtime_context = prepared_step.runtime_context
    execution_state = prepared_step.execution_state
    running_event = prepared_step.running_event
    retry_parent_step_id = prepared_step.retry_parent_step_id
    start_iso = prepared_step.start_iso
    start_perf = prepared_step.start_perf
    executor_stdout_lines = prepared_step.executor_stdout_lines
    executor_stderr_lines = prepared_step.executor_stderr_lines
    system_stdout_lines = prepared_step.system_stdout_lines
    system_stderr_lines = prepared_step.system_stderr_lines
    log_events = prepared_step.log_events
    next_name = running_event.step_name
    status = "succeeded"
    error = None
    metrics: dict[str, Any] = {}
    validators = _workflow_validators(state.pipeline_id, next_name)
    environment_fragments = STORE.list_environment_fragments(
        run_id=run_id,
        env_type=running_event.env_type,
        env_id=running_event.env_id,
    )
    input_validation = build_runtime_transition_validation_for_direction(
        validators=validators,
        artifact_id=str(runtime_context.get("artifact_id") or execution_state.outputs.get("artifact_id") or ""),
        mime_type=str(runtime_context.get("mime_type") or ""),
        fixture_path=str(runtime_context.get("fixture_path") or "") or None,
        run_id=run_id,
        step_id=running_event.step_id,
        step_outputs=execution_state.outputs,
        environment_fragments=environment_fragments,
        direction="input",
    )
    pre_execution_outputs = deepcopy(execution_state.outputs)
    if isinstance(input_validation, dict):
        persisted_inputs = runtime_context.get("step_input_validations")
        if not isinstance(persisted_inputs, dict):
            persisted_inputs = {}
        persisted_inputs[running_event.step_id] = deepcopy(input_validation)
        runtime_context["step_input_validations"] = persisted_inputs

        pending_output_validation = _pending_validation_for_direction(validators, "output")
        if isinstance(pending_output_validation, dict):
            persisted_outputs = runtime_context.get("step_output_validations")
            if not isinstance(persisted_outputs, dict):
                persisted_outputs = {}
            persisted_outputs[running_event.step_id] = deepcopy(pending_output_validation)
            runtime_context["step_output_validations"] = persisted_outputs

        persisted_transition = runtime_context.get("step_transition_validations")
        if not isinstance(persisted_transition, dict):
            persisted_transition = {}
        persisted_transition[running_event.step_id] = deepcopy(
            _combine_phase_validations(input_validation, pending_output_validation)
            or input_validation
        )
        runtime_context["step_transition_validations"] = persisted_transition

    snapshots = runtime_context.get("step_output_snapshots")
    if not isinstance(snapshots, dict):
        snapshots = {}
    snapshots[running_event.step_id] = deepcopy(execution_state.outputs)
    runtime_context["step_output_snapshots"] = snapshots
    STORE.set_debug_runtime_context(run_id, runtime_context)

    def publish_running_logs() -> None:
        _publish_running_event_logs(
            run_id=state.run_id,
            step_id=running_event.step_id,
            executor_stdout_lines=executor_stdout_lines,
            executor_stderr_lines=executor_stderr_lines,
            system_stdout_lines=system_stdout_lines,
            system_stderr_lines=system_stderr_lines,
            log_events=log_events,
        )

    heartbeat_task: asyncio.Task[None] | None = None

    async def emit_running_heartbeats() -> None:
        while True:
            await asyncio.sleep(_RUNNING_LOG_HEARTBEAT_INTERVAL_SECONDS)
            heartbeat_message = f"[{next_name}] still running"
            system_stdout_lines.append(heartbeat_message)
            _append_log_event(
                log_events,
                at=datetime.now(UTC).isoformat(),
                source="system",
                stream="stdout",
                message=heartbeat_message,
            )
            publish_running_logs()

    if _validation_has_failures(input_validation):
        status = "failed"
        error = {"reason": "input validation failed"}
        metrics = {
            "executor": "validation-gate",
            "control": "next_step",
            "retry_boundary": retry_parent_step_id is not None,
            "validation_phase": "input",
        }
    else:
        try:
            heartbeat_task = asyncio.create_task(emit_running_heartbeats())
            stdout_buffer = _StreamingLogBuffer(
                executor_stdout_lines,
                publish_running_logs,
                lambda line: _append_log_event(
                    log_events,
                    at=datetime.now(UTC).isoformat(),
                    source="executor",
                    stream="stdout",
                    message=line,
                ),
            )
            stderr_buffer = _StreamingLogBuffer(
                executor_stderr_lines,
                publish_running_logs,
                lambda line: _append_log_event(
                    log_events,
                    at=datetime.now(UTC).isoformat(),
                    source="executor",
                    stream="stderr",
                    message=line,
                ),
            )
            metrics = await _run_execute_step_with_thread_scoped_capture(
                implementation_next_name=prepared_step.implementation_next_name,
                execution_state=execution_state,
                debug_scope=prepared_step.debug_scope,
                stdout_buffer=stdout_buffer,
                stderr_buffer=stderr_buffer,
            )
        except Exception as exc:
            if 'stdout_buffer' in locals():
                stdout_buffer.flush()
            if 'stderr_buffer' in locals():
                stderr_buffer.flush()
            metrics = {
                "executor": "ikam.forja.debug_execution",
                "control": "next_step",
                "retry_boundary": retry_parent_step_id is not None,
                    "details": {
                        "error": str(exc),
                        "step": next_name,
                        "implementation_step": prepared_step.implementation_next_name,
                    },
                }
            if getattr(exc, "injection_used", False):
                metrics["injection_used"] = True
            status = "failed"
            error = {"reason": str(exc)}
            if hasattr(exc, "drift_at"):
                error["drift_at"] = exc.drift_at
        finally:
            if heartbeat_task is not None:
                heartbeat_task.cancel()
                try:
                    await heartbeat_task
                except asyncio.CancelledError:
                    pass

    output_validation = None
    if status == "succeeded":
        output_validation = build_runtime_transition_validation_for_direction(
            validators=validators,
            artifact_id=str(runtime_context.get("artifact_id") or execution_state.outputs.get("artifact_id") or ""),
            mime_type=str(runtime_context.get("mime_type") or ""),
            fixture_path=str(runtime_context.get("fixture_path") or "") or None,
            run_id=run_id,
            step_id=running_event.step_id,
            step_outputs=_normalize_step_outputs_for_output_validation(
                step_outputs=execution_state.outputs,
                pre_execution_outputs=pre_execution_outputs,
            ),
            environment_fragments=environment_fragments,
            direction="output",
        )
        if _validation_has_failures(output_validation):
            status = "failed"
            error = {"reason": "output validation failed"}
            metrics = {
                **(metrics if isinstance(metrics, dict) else {}),
                "validation_phase": "output",
            }
    duration_ms = max(1, int((perf_counter() - start_perf) * 1000))
    end_iso = datetime.now(UTC).isoformat()
    finished_message = f"[{next_name}] step finished at {end_iso}"
    system_stdout_lines.append(finished_message)
    _append_log_event(
        log_events,
        at=end_iso,
        source="system",
        stream="stdout",
        message=finished_message,
    )
    if not isinstance(metrics, dict):
        metrics = {"value": metrics}
    metrics["log_events"] = list(log_events)
    if executor_stdout_lines or executor_stderr_lines or system_stdout_lines or system_stderr_lines or status == "succeeded":
        metrics["executor_logs"] = {
            "stdout_lines": executor_stdout_lines,
            "stderr_lines": executor_stderr_lines,
        }
        metrics["system_logs"] = {
            "stdout_lines": system_stdout_lines,
            "stderr_lines": system_stderr_lines,
        }
        metrics["logs"] = _legacy_logs_from_log_events(log_events)
    _capture_step_logs(next_name, metrics, started_at=start_iso, ended_at=end_iso)

    runtime_context["step_outputs"] = execution_state.outputs
    if isinstance(input_validation, dict):
        persisted_inputs = runtime_context.get("step_input_validations")
        if not isinstance(persisted_inputs, dict):
            persisted_inputs = {}
        persisted_inputs[running_event.step_id] = deepcopy(input_validation)
        runtime_context["step_input_validations"] = persisted_inputs
    if isinstance(output_validation, dict):
        persisted_outputs = runtime_context.get("step_output_validations")
        if not isinstance(persisted_outputs, dict):
            persisted_outputs = {}
        persisted_outputs[running_event.step_id] = deepcopy(output_validation)
        runtime_context["step_output_validations"] = persisted_outputs
    combined_validation = _combine_phase_validations(input_validation, output_validation)
    if isinstance(combined_validation, dict):
        persisted = runtime_context.get("step_transition_validations")
        if not isinstance(persisted, dict):
            persisted = {}
        persisted[running_event.step_id] = deepcopy(combined_validation)
        runtime_context["step_transition_validations"] = persisted

    if status == "succeeded":
        _persist_runtime_hot_document_set(
            implementation_next_name=prepared_step.implementation_next_name,
            run_id=run_id,
            step_id=running_event.step_id,
            runtime_context=runtime_context,
            step_outputs=execution_state.outputs,
        )
        _persist_runtime_hot_chunk_extraction_set(
            run_id=run_id,
            step_id=running_event.step_id,
            runtime_context=runtime_context,
            step_outputs=execution_state.outputs,
        )
        _persist_runtime_hot_entity_relationship_set(
            implementation_next_name=prepared_step.implementation_next_name,
            run_id=run_id,
            step_id=running_event.step_id,
            runtime_context=runtime_context,
            step_outputs=execution_state.outputs,
        )
        _persist_runtime_hot_claim_set(
            implementation_next_name=prepared_step.implementation_next_name,
            run_id=run_id,
            step_id=running_event.step_id,
            runtime_context=runtime_context,
            step_outputs=execution_state.outputs,
        )
    snapshots = runtime_context.get("step_output_snapshots")
    if not isinstance(snapshots, dict):
        snapshots = {}
    snapshots[running_event.step_id] = deepcopy(execution_state.outputs)
    runtime_context["step_output_snapshots"] = snapshots
    STORE.set_debug_runtime_context(run_id, runtime_context)
    _sync_graph_snapshot_from_step_outputs(
        run_id=run_id,
        step_name=prepared_step.implementation_next_name,
        step_outputs=execution_state.outputs,
        runtime_context=runtime_context,
    )

    upcoming_step_name = next_step_name(next_name, scope=prepared_step.debug_scope)
    selected_next_transition_id = upcoming_step_name

    event = STORE.update_debug_event(
        run_id=state.run_id,
        step_id=running_event.step_id,
        status=status,
        ended_at=end_iso,
        duration_ms=duration_ms,
        metrics={
            **metrics,
            "control": "next_step",
            "retry_boundary": retry_parent_step_id is not None,
            "next_transition_id": selected_next_transition_id,
        },
        error=error,
    )
    _append_scoped_debug_fragments(run_id=run_id, event=event, step_outputs=execution_state.outputs)
    return event


def _append_scoped_debug_fragments(*, run_id: str, event: DebugStepEvent, step_outputs: dict[str, object]) -> None:
    """Persist step-scoped artifacts for debug detail and environment summaries."""
    run = STORE.get_run(run_id)
    if run is None:
        return

    implementation_step_name = _implementation_step_name(event.pipeline_id, event.step_name)

    base_meta = {
        "ref": _legacy_scope_ref(event.env_type, event.env_id),
        "env_type": event.env_type,
        "env_id": event.env_id,
        "step_id": event.step_id,
        "attempt_index": event.attempt_index,
    }

    def append_fragment(*, cas_id: str | None, mime_type: str, value: object, record_type: str) -> None:
        run.graph.fragments.append(
            {
                "id": cas_id,
                "cas_id": cas_id,
                "mime_type": mime_type,
                "value": value,
                "meta": {
                    **base_meta,
                    "record_type": record_type,
                },
            }
        )

    if implementation_step_name == "map.conceptual.lift.surface_fragments":
        decomposition = step_outputs.get("decomposition")
        fragment_artifact_map = step_outputs.get("fragment_artifact_map") if isinstance(step_outputs.get("fragment_artifact_map"), dict) else {}
        for fragment in getattr(decomposition, "structural", []):
            cas_id = getattr(fragment, "cas_id", None)
            append_fragment(
                cas_id=cas_id,
                mime_type=str(getattr(fragment, "mime_type", "application/octet-stream")),
                value=getattr(fragment, "value", None),
                record_type="surface_fragment",
            )
            if run.graph.fragments:
                latest = run.graph.fragments[-1]
                meta = latest.get("meta") if isinstance(latest, dict) and isinstance(latest.get("meta"), dict) else {}
                artifact_ref = fragment_artifact_map.get(str(cas_id)) if cas_id is not None else None
                if artifact_ref:
                    meta["artifact_id"] = artifact_ref
                    if isinstance(latest, dict):
                        latest["meta"] = meta

    if implementation_step_name == "parse_artifacts":
        document_fragment_artifact_map = step_outputs.get("document_fragment_artifact_map") if isinstance(step_outputs.get("document_fragment_artifact_map"), dict) else {}
        for fragment in step_outputs.get("document_fragments", []) if isinstance(step_outputs.get("document_fragments"), list) else []:
            cas_id = getattr(fragment, "cas_id", None)
            append_fragment(
                cas_id=cas_id,
                mime_type=str(getattr(fragment, "mime_type", "application/vnd.ikam.loaded-document+json")),
                value=getattr(fragment, "value", None),
                record_type="loaded_document",
            )
            if run.graph.fragments:
                latest = run.graph.fragments[-1]
                meta = latest.get("meta") if isinstance(latest, dict) and isinstance(latest.get("meta"), dict) else {}
                artifact_ref = document_fragment_artifact_map.get(str(cas_id)) if cas_id is not None else None
                if artifact_ref:
                    meta["artifact_id"] = artifact_ref
                if isinstance(latest, dict):
                    latest["meta"] = meta

    if implementation_step_name == "map.conceptual.normalize.discovery":
        lifted_from_map = step_outputs.get("lifted_from_map") if isinstance(step_outputs.get("lifted_from_map"), dict) else {}
        fragment_artifact_map = step_outputs.get("fragment_artifact_map") if isinstance(step_outputs.get("fragment_artifact_map"), dict) else {}
        for fragment in step_outputs.get("ir_fragments", []) if isinstance(step_outputs.get("ir_fragments"), list) else []:
            ir_id = getattr(fragment, "cas_id", None)
            append_fragment(
                cas_id=ir_id,
                mime_type=str(getattr(fragment, "mime_type", "application/octet-stream")),
                value=getattr(fragment, "value", None),
                record_type="ir_fragment",
            )
            if run.graph.fragments:
                latest = run.graph.fragments[-1]
                meta = latest.get("meta") if isinstance(latest, dict) and isinstance(latest.get("meta"), dict) else {}
                source_surface_id = lifted_from_map.get(str(ir_id)) if ir_id is not None else None
                if isinstance(source_surface_id, str):
                    meta["source_surface_fragment_id"] = source_surface_id
                    source_artifact_id = fragment_artifact_map.get(source_surface_id)
                    if source_artifact_id:
                        meta["artifact_id"] = source_artifact_id
                if isinstance(latest, dict):
                    latest["meta"] = meta

    if implementation_step_name == "map.conceptual.lift.claims":
        for fragment in step_outputs.get("claim_set", {}).get("claims", []) if isinstance(step_outputs.get("claim_set"), dict) else []:
            if not isinstance(fragment, dict):
                continue
            cas_id = str(fragment.get("cas_id") or "") or None
            append_fragment(
                cas_id=cas_id,
                mime_type=str(fragment.get("mime_type") or "application/vnd.ikam.claim+json"),
                value=fragment.get("value"),
                record_type="claim",
            )
            if run.graph.fragments:
                latest = run.graph.fragments[-1]
                meta = latest.get("meta") if isinstance(latest, dict) and isinstance(latest.get("meta"), dict) else {}
                value = fragment.get("value") if isinstance(fragment.get("value"), dict) else {}
                artifact_id = value.get("artifact_id")
                entity_relationship_fragment_id = value.get("entity_relationship_fragment_id")
                if isinstance(artifact_id, str) and artifact_id:
                    meta["artifact_id"] = artifact_id
                if isinstance(entity_relationship_fragment_id, str) and entity_relationship_fragment_id:
                    meta["entity_relationship_fragment_id"] = entity_relationship_fragment_id
                if isinstance(latest, dict):
                    latest["meta"] = meta

    if implementation_step_name == "map.reconstructable.normalize":
        normalized = step_outputs.get("normalized_fragments")
        if isinstance(normalized, list):
            for fragment in normalized:
                append_fragment(
                    cas_id=getattr(fragment, "cas_id", None),
                    mime_type=str(getattr(fragment, "mime_type", "application/octet-stream")),
                    value=getattr(fragment, "value", None),
                    record_type="normalized_fragment",
                )
        programs = step_outputs.get("reconstruction_programs")
        if isinstance(programs, list):
            for program in programs:
                append_fragment(
                    cas_id=getattr(program, "cas_id", None),
                    mime_type=str(getattr(program, "mime_type", "application/vnd.ikam.reconstruction-program+json")),
                    value=getattr(program, "value", None),
                    record_type="reconstruction_program",
                )

    if implementation_step_name == "map.conceptual.verify.discovery_gate":
        verification_fragment = step_outputs.get("verification_result_fragment")
        if verification_fragment is not None:
            append_fragment(
                cas_id=getattr(verification_fragment, "cas_id", None),
                mime_type=str(getattr(verification_fragment, "mime_type", "application/vnd.ikam.verification+json")),
                value=getattr(verification_fragment, "value", None),
                record_type="verification",
            )


async def _execute_next_pipeline_step_async(
    *,
    run_id: str,
    state: DebugRunState,
    prepared_step: _PreparedPipelineStep | None = None,
) -> None:
    """Execute a single pipeline step asynchronously for manual next_step controls."""
    try:
        terminal_event = None
        if prepared_step is None:
            prepared_step, terminal_event = _prepare_next_pipeline_step(run_id=run_id, state=state)
        if terminal_event is not None:
            event = terminal_event
        else:
            if prepared_step is None:
                raise RuntimeError("Expected prepared step or terminal event")
            event = await asyncio.to_thread(
                lambda: asyncio.run(
                    _execute_prepared_pipeline_step(run_id=run_id, state=state, prepared_step=prepared_step)
                )
            )
        pipeline_steps = _debug_pipeline_steps(state.pipeline_id)
        if event.status != "succeeded" and state.execution_state != "budget_exhausted":
            state.execution_state = "paused"
        elif event.status == "succeeded" and state.current_step_name == pipeline_steps[-1]:
            state.execution_state = "completed"
        else:
            state.execution_state = "paused"
        STORE.set_debug_run_state(run_id, state)
    finally:
        _RUN_STEP_TASKS.pop(run_id, None)


@router.get("/cases")
def list_cases():
    return {"cases": available_cases()}


@router.post("/run")
def run_benchmarks(case_ids: str | None = None, reset: bool = False, include_evaluation: bool = True, pipeline_id: str | None = None):
    return run_benchmark(case_ids=case_ids, reset=reset, include_evaluation=include_evaluation, pipeline_id=pipeline_id)


@router.post("/test/seed-scenario")
def seed_scenario(payload: SeedScenarioPayload):
    _require_test_gate()
    if payload.scenario_key != "core_stream_baseline":
        raise HTTPException(status_code=400, detail={"status": "invalid_scenario_key", "scenario_key": payload.scenario_key})

    run_id = payload.overrides.get("run_id") or f"seed-{uuid4().hex[:8]}"
    case_id = payload.overrides.get("case_id") or "s-local-retail-v01"
    project_id = payload.overrides.get("project_id") or f"{case_id}#seed"
    pipeline_id = payload.overrides.get("pipeline_id") or "compression-rerender/v1"
    pipeline_run_id = payload.overrides.get("pipeline_run_id") or f"pipe-{run_id}"

    asset_manifest: list[dict[str, Any]] = []
    asset_payloads: list[dict[str, Any]] = []
    try:
        fixture = load_case_fixture(case_id)
        for index, asset in enumerate(fixture.assets):
            file_name = str(getattr(asset, "file_name", f"asset-{index}"))
            mime_type = str(getattr(asset, "mime_type", "application/octet-stream"))
            payload_bytes = getattr(asset, "payload", b"")
            size_bytes = len(payload_bytes) if isinstance(payload_bytes, (bytes, bytearray)) else 0
            asset_manifest.append(
                {
                    "artifact_id": f"{project_id}:{file_name.replace('/', '_')}",
                    "filename": file_name,
                    "mime_type": mime_type,
                    "size_bytes": size_bytes,
                    "role": "source",
                    "included": True,
                    "exclusion_reason": None,
                }
            )
            asset_payloads.append(
                {
                    "artifact_id": f"{project_id}:{file_name.replace('/', '_')}",
                    "filename": file_name,
                    "mime_type": mime_type,
                    "payload": bytes(payload_bytes) if isinstance(payload_bytes, (bytes, bytearray)) else bytes(str(payload_bytes), encoding="utf-8"),
                }
            )
    except Exception:
        asset_manifest = []
        asset_payloads = []

    run = BenchmarkRunRecord(
        run_id=run_id,
        project_id=project_id,
        case_id=case_id,
        stages=[],
        decisions=[],
        project={},
        graph=GraphSnapshot(graph_id=project_id, fragments=[]),
    )
    STORE.add_run(run)
    STORE.create_debug_run_state(
        DebugRunState(
            run_id=run_id,
            pipeline_id=pipeline_id,
            pipeline_run_id=pipeline_run_id,
            project_id=project_id,
            operation_id=payload.overrides.get("operation_id") or f"op-{run_id}",
            env_type=payload.overrides.get("env_type") or "dev",
            env_id=payload.overrides.get("env_id") or "dev-seed",
            execution_mode=payload.overrides.get("execution_mode") or "manual",
            execution_state=payload.overrides.get("execution_state") or "paused",
            current_step_name=payload.overrides.get("current_step_name") or "map.conceptual.verify.discovery_gate",
            current_attempt_index=int(payload.overrides.get("current_attempt_index") or "1"),
        )
    )
    now = datetime.now(UTC).isoformat()
    STORE.append_debug_event(
        DebugStepEvent(
            event_id=f"ev-{uuid4().hex}",
            run_id=run_id,
            pipeline_id=pipeline_id,
            pipeline_run_id=pipeline_run_id,
            project_id=project_id,
            operation_id=payload.overrides.get("operation_id") or f"op-{run_id}",
            env_type=payload.overrides.get("env_type") or "dev",
            env_id=payload.overrides.get("env_id") or "dev-seed",
            step_name="init.initialize",
            step_id=f"step-{uuid4().hex[:10]}",
            status="succeeded",
            attempt_index=1,
            retry_parent_step_id=None,
            started_at=now,
            ended_at=now,
            duration_ms=1,
            metrics={
                "seeded": True,
                "immediate_stream": True,
                "details": {"case_id": case_id, "project_id": project_id, "asset_count": len(asset_manifest)},
            },
            error=None,
        )
    )
    STORE.set_debug_runtime_context(
        run_id,
        {
            "source_bytes": f"seed scenario {case_id}".encode("utf-8"),
            "mime_type": "text/markdown",
            "artifact_id": f"{project_id}:{case_id}",
            "fixture_path": str(case_fixture_dir(case_id).resolve()),
            "asset_manifest": asset_manifest,
            "asset_payloads": asset_payloads,
            "step_outputs": {},
        },
    )

    return {
        "status": "ok",
        "scenario_key": payload.scenario_key,
        "run": {
            "run_id": run_id,
            "project_id": project_id,
            "case_id": case_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
        },
    }


@router.post("/merge")
def merge_benchmarks(graph_ids: str, apply: bool = False):
    return run_merge_benchmark(graph_ids=graph_ids, apply=apply)


@router.post("/ingest-poc")
def run_ingest_poc():
    return run_staging_normalize_promote_enrich_poc()


@router.get("/runs")
def list_runs():
    result = []
    for run in STORE.list_runs():
        debug_state = STORE.get_debug_run_state(run.run_id)
        entry: dict = {
            "run_id": run.run_id,
            "project_id": run.project_id,
            "case_id": run.case_id,
            "stages": run.stages,
            "decisions": run.decisions,
            "semantic_entities": len((run.semantic or {}).get("entities", [])),
            "semantic_relations": len((run.semantic or {}).get("relations", [])),
            "answer_quality": run.answer_quality,
            "evaluation": run.evaluation,
        }
        if debug_state:
            entry["pipeline_id"] = debug_state.pipeline_id
            entry["pipeline_run_id"] = debug_state.pipeline_run_id
            entry["ref"] = _run_scope_ref_from_state(debug_state)
        result.append(entry)
    return result


@router.get("/runs/{run_id}")
def get_run(run_id: str):
    run = STORE.get_run(run_id)
    if not run:
        return {"run_id": run_id, "status": "missing"}
    debug_state = STORE.get_debug_run_state(run_id)
    return {
        "run_id": run.run_id,
        "project_id": run.project_id,
        "case_id": run.case_id,
        "stages": run.stages,
        "decisions": run.decisions,
        "project": run.project,
        "graph": run.graph,
        "answer_quality": run.answer_quality,
        "evaluation": run.evaluation,
        "ref": _run_scope_ref_from_state(debug_state) if debug_state else None,
    }


@router.post("/runs/{run_id}/reviews")
def save_manual_review(run_id: str, payload: ReviewPayload):
    answer_quality = STORE.apply_review(run_id, payload.model_dump())
    if answer_quality is None:
        return {"run_id": run_id, "status": "missing"}
    return {
        "run_id": run_id,
        "status": "updated",
        "answer_quality": answer_quality,
    }


@router.get("/runs/{run_id}/debug-stream")
def get_debug_stream(run_id: str, pipeline_id: str, pipeline_run_id: str):
    state = STORE.get_debug_run_state(run_id)
    pipeline_steps = _debug_pipeline_steps(pipeline_id)
    if not state:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "pipeline_steps": pipeline_steps,
            "events": [],
        }
    if state.pipeline_id != pipeline_id or state.pipeline_run_id != pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "pipeline_steps": pipeline_steps,
            "events": [],
        }
    events = STORE.list_debug_events(run_id)
    availability = _control_availability(state=state, events=events, pipeline_steps=pipeline_steps)
    return {
        "run_id": run_id,
        "pipeline_id": pipeline_id,
        "pipeline_run_id": pipeline_run_id,
        "status": "ok",
        "execution_mode": state.execution_mode,
        "execution_state": state.execution_state,
        "control_availability": availability,
        "pipeline_steps": pipeline_steps,
        "events": [_serialize_debug_event(event) for event in events],
    }


@router.get("/runs/{run_id}/debug-step/{step_id}/detail")
def get_debug_step_detail(run_id: str, step_id: str):
    """Drill-through API: return canonical step detail payload."""
    events = STORE.list_debug_events(run_id)
    matched_event = None
    for ev in events:
        if ev.step_id == step_id:
            matched_event = ev
            break
    if matched_event is None:
        raise HTTPException(status_code=404, detail={"status": "not_found", "step_id": step_id})

    state = STORE.get_debug_run_state(run_id)
    if state is None:
        raise HTTPException(status_code=404, detail={"status": "not_found", "run_id": run_id})

    runtime_context = STORE.get_debug_runtime_context(run_id)
    snapshots = (runtime_context or {}).get("step_output_snapshots", {})
    if isinstance(snapshots, dict) and isinstance(snapshots.get(step_id), dict):
        step_outputs = snapshots.get(step_id, {})
    else:
        step_outputs = (runtime_context or {}).get("step_outputs", {})
    run_record = STORE.get_run(run_id)

    payload = _serialize_step_detail(
        event=matched_event,
        state=state,
        step_outputs=step_outputs,
        runtime_context=runtime_context or {},
        run_record=run_record,
        events=events,
        step_output_snapshots=snapshots if isinstance(snapshots, dict) else {},
    )
    _validate_canonical_step_detail(payload)
    return payload


@router.get("/runs/{run_id}/inspection")
def get_inspection_subgraph(run_id: str, ref: str, max_depth: int = 1):
    subgraph = _resolve_inspection_subgraph(run_id=run_id, ref=ref, max_depth=max_depth)
    return subgraph.model_dump(mode="json")


def _find_runtime_asset(runtime_context: dict[str, Any], artifact_id: str) -> dict[str, Any] | None:
    payloads = runtime_context.get("asset_payloads")
    if not isinstance(payloads, list):
        return None
    normalized = artifact_id.strip()
    normalized_basename = normalized.rsplit("/", 1)[-1]
    for item in payloads:
        if not isinstance(item, dict):
            continue
        item_artifact_id = str(item.get("artifact_id") or "")
        item_file_name = str(item.get("filename") or item.get("file_name") or "")
        if (
            item_artifact_id == normalized
            or item_file_name == normalized
            or item_file_name.rsplit("/", 1)[-1] == normalized_basename
        ):
            return item
    return None


def _build_artifact_preview(asset: dict[str, Any]) -> dict[str, Any]:
    file_name = str(asset.get("filename") or asset.get("artifact_id") or "unknown")
    mime_type = str(asset.get("mime_type") or "application/octet-stream")
    if mime_type == "application/octet-stream" and file_name.endswith(".docx"):
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if mime_type == "application/octet-stream" and file_name.endswith(".xlsx"):
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if mime_type == "application/octet-stream" and file_name.endswith(".pptx"):
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    payload = asset.get("payload")
    if not isinstance(payload, (bytes, bytearray)):
        raise HTTPException(status_code=422, detail={"status": "invalid_asset_payload", "file_name": file_name})

    payload_bytes = bytes(payload)
    size_bytes = len(payload_bytes)

    if mime_type in {"text/markdown", "text/plain"} or file_name.endswith((".md", ".txt")):
        return {
            "kind": "text",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes},
            "preview": {"text": payload_bytes.decode("utf-8", errors="ignore")[:16000]},
        }

    if mime_type == "application/json" or file_name.endswith(".json"):
        text = payload_bytes.decode("utf-8", errors="ignore")
        parsed = json.loads(text)
        return {
            "kind": "json",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes},
            "preview": {"parsed": parsed},
        }

    if mime_type == "application/pdf" or file_name.endswith(".pdf"):
        return {
            "kind": "pdf",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes},
            "preview": {
                "encoding": "base64",
                "bytes_b64": base64.b64encode(payload_bytes).decode("ascii"),
            },
        }

    if mime_type.startswith("image/"):
        return {
            "kind": "image",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes},
            "preview": {
                "encoding": "data_url",
                "data_url": f"data:{mime_type};base64,{base64.b64encode(payload_bytes).decode('ascii')}",
            },
        }

    if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or file_name.endswith(".xlsx"):
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(payload_bytes), data_only=True, read_only=True)
        sheets: list[dict[str, Any]] = []
        for sheet_name in workbook.sheetnames:
            ws = workbook[sheet_name]
            rows: list[list[str]] = []
            for row in ws.iter_rows(min_row=1, max_row=20, max_col=12, values_only=True):
                rows.append(["" if cell is None else str(cell) for cell in row])
            sheets.append({"sheet_name": sheet_name, "rows": rows})
        return {
            "kind": "table",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes, "sheet_count": len(sheets)},
            "preview": {"sheets": sheets},
        }

    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_name.endswith(".pptx"):
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(payload_bytes))
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(presentation.slides, start=1):
            text_lines: list[str] = []
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text.strip()
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = (paragraph.text or "").strip()
                    if text:
                        text_lines.append(text)
            slides.append({"index": index, "title": title or f"Slide {index}", "lines": text_lines[:24]})
        return {
            "kind": "slides",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes, "slide_count": len(slides)},
            "preview": {"slides": slides},
        }

    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_name.endswith(".docx"):
        paragraphs: list[str] = []
        try:
            from docx import Document  # type: ignore

            document = Document(io.BytesIO(payload_bytes))
            paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text][:120]
        except Exception:
            with zipfile.ZipFile(io.BytesIO(payload_bytes), mode="r") as zf:
                xml_bytes = zf.read("word/document.xml")
            root = ElementTree.fromstring(xml_bytes)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            extracted: list[str] = []
            for paragraph in root.findall(".//w:p", ns):
                runs = [node.text or "" for node in paragraph.findall(".//w:t", ns)]
                text = "".join(runs).strip()
                if text:
                    extracted.append(text)
            paragraphs = extracted[:120]
        return {
            "kind": "doc",
            "mime_type": mime_type,
            "file_name": file_name,
            "metadata": {"size_bytes": size_bytes, "paragraph_count": len(paragraphs)},
            "preview": {"paragraphs": paragraphs},
        }

    return {
        "kind": "binary",
        "mime_type": mime_type,
        "file_name": file_name,
        "metadata": {"size_bytes": size_bytes},
        "preview": {},
    }


@router.get("/runs/{run_id}/artifacts/preview")
def get_artifact_preview(run_id: str, artifact_id: str):
    runtime_context = STORE.get_debug_runtime_context(run_id)
    if runtime_context is None:
        raise HTTPException(status_code=404, detail={"status": "not_found", "run_id": run_id})
    asset = _find_runtime_asset(runtime_context, artifact_id)
    if asset is None:
        raise HTTPException(status_code=404, detail={"status": "not_found", "artifact_id": artifact_id})
    return _build_artifact_preview(asset)


def _node_kind_from_payload(payload: dict[str, Any]) -> str:
    meta = payload.get("meta") if isinstance(payload.get("meta"), dict) else {}
    record_type = str(meta.get("record_type") or "")
    mime_type = str(payload.get("mime_type") or "")
    if "verification" in record_type or "verification" in mime_type:
        return "verification"
    if "reconstruction" in record_type or "reconstruction" in mime_type:
        return "program"
    if "normalized" in record_type:
        return "normalized"
    if "ir" in record_type or "ir" in mime_type:
        return "ir"
    return "surface"


def _serialize_trace_topic_sequence(value: Any) -> list[dict[str, str]]:
    if not isinstance(value, list):
        return []
    items: list[dict[str, str]] = []
    for item in value:
        if not isinstance(item, dict):
            continue
        topic = str(item.get("topic") or "").strip()
        event_type = str(item.get("event_type") or "").strip()
        status = str(item.get("status") or "").strip()
        if not topic or not event_type or not status:
            continue
        items.append(
            {
                "topic": topic,
                "event_type": event_type,
                "status": status,
            }
        )
    return items


def _serialize_trace_timeline(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    items: list[dict[str, Any]] = []
    for item in value:
        if not isinstance(item, dict):
            continue
        topic = str(item.get("topic") or "").strip()
        event_type = str(item.get("event_type") or "").strip()
        status = str(item.get("status") or "").strip()
        occurred_at = str(item.get("occurred_at") or "").strip()
        payload = item.get("payload") if isinstance(item.get("payload"), dict) else {}
        if not topic or not event_type or not status:
            continue
        items.append(
            {
                "topic": topic,
                "event_type": event_type,
                "status": status,
                "occurred_at": occurred_at,
                "payload": payload,
            }
        )
    return items


def _serialize_raw_events(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    items: list[dict[str, Any]] = []
    for item in value:
        if not isinstance(item, dict):
            continue
        topic = str(item.get("topic") or "").strip()
        payload = item.get("payload") if isinstance(item.get("payload"), dict) else {}
        if not topic:
            continue
        items.append({"topic": topic, "payload": payload})
    return items


def _serialize_step_logs(event: DebugStepEvent) -> dict[str, list[str]] | None:
    log_events = _serialize_log_events(event)
    if log_events is not None:
        return _legacy_logs_from_log_events(log_events)
    if not isinstance(event.metrics, dict):
        return None
    raw_logs = event.metrics.get("logs")
    logs = raw_logs if isinstance(raw_logs, dict) else {}
    stdout_lines = [
        str(item)
        for item in (logs.get("stdout_lines") or [])
        if isinstance(item, str) and item.strip()
    ]
    stderr_lines = [
        str(item)
        for item in (logs.get("stderr_lines") or [])
        if isinstance(item, str) and item.strip()
    ]
    if not stdout_lines and not stderr_lines:
        return None
    return {
        "stdout_lines": stdout_lines,
        "stderr_lines": stderr_lines,
    }


def _serialize_named_step_logs(event: DebugStepEvent, key: str) -> dict[str, list[str]] | None:
    log_events = _serialize_log_events(event)
    if log_events is not None and key in {"executor_logs", "system_logs"}:
        source = "executor" if key == "executor_logs" else "system"
        return _named_logs_from_log_events(log_events, source=source)
    if not isinstance(event.metrics, dict):
        return None
    raw_logs = event.metrics.get(key)
    logs = raw_logs if isinstance(raw_logs, dict) else {}
    stdout_lines = [
        str(item)
        for item in (logs.get("stdout_lines") or [])
        if isinstance(item, str) and item.strip()
    ]
    stderr_lines = [
        str(item)
        for item in (logs.get("stderr_lines") or [])
        if isinstance(item, str) and item.strip()
    ]
    if not stdout_lines and not stderr_lines:
        legacy_logs = _serialize_step_logs(event)
        if legacy_logs is None:
            return None
        if key == "executor_logs":
            return legacy_logs
        if key == "system_logs":
            return {"stdout_lines": [], "stderr_lines": []}
        return None
    return {
        "stdout_lines": stdout_lines,
        "stderr_lines": stderr_lines,
    }


def _extract_runtime_trace_logs(event: DebugStepEvent) -> dict[str, list[str]] | None:
    if not isinstance(event.metrics, dict):
        return None
    raw_trace = event.metrics.get("trace")
    trace = raw_trace if isinstance(raw_trace, dict) else {}
    stdout_lines: list[str] = []
    stderr_lines: list[str] = []
    seen_stdout: set[str] = set()
    seen_stderr: set[str] = set()

    def append_lines(payload: Any) -> None:
        if not isinstance(payload, dict):
            return
        for item in payload.get("stdout_lines") or []:
            if isinstance(item, str) and item.strip() and item not in seen_stdout:
                seen_stdout.add(item)
                stdout_lines.append(item)
        for item in payload.get("stderr_lines") or []:
            if isinstance(item, str) and item.strip() and item not in seen_stderr:
                seen_stderr.add(item)
                stderr_lines.append(item)

    for item in trace.get("timeline") or []:
        if isinstance(item, dict):
            append_lines(item.get("payload"))
    for item in trace.get("raw_events") or []:
        if isinstance(item, dict):
            append_lines(item.get("payload"))

    if not stdout_lines and not stderr_lines:
        return None
    return {
        "stdout_lines": stdout_lines,
        "stderr_lines": stderr_lines,
    }


def _serialize_log_events(event: DebugStepEvent) -> list[dict[str, Any]] | None:
    if not isinstance(event.metrics, dict):
        return None
    raw_log_events = event.metrics.get("log_events")
    if not isinstance(raw_log_events, list):
        return None

    items: list[dict[str, Any]] = []
    for item in raw_log_events:
        if not isinstance(item, dict):
            continue
        seq = item.get("seq")
        at = str(item.get("at") or "").strip()
        source = str(item.get("source") or "").strip()
        stream = str(item.get("stream") or "").strip()
        message = str(item.get("message") or "").strip()
        if not isinstance(seq, int):
            continue
        if not at or source not in {"executor", "system"} or stream not in {"stdout", "stderr"} or not message:
            continue
        items.append(
            {
                "seq": seq,
                "at": at,
                "source": source,
                "stream": stream,
                "message": message,
            }
        )
    if not items:
        return None
    return items


def _serialize_step_trace(event: DebugStepEvent) -> dict[str, Any] | None:
    if not isinstance(event.metrics, dict):
        return None

    raw_trace = event.metrics.get("trace")
    trace = raw_trace if isinstance(raw_trace, dict) else {}
    raw_details = event.metrics.get("details")
    details = raw_details if isinstance(raw_details, dict) else {}

    payload = {
        "workflow_id": str(trace.get("workflow_id") or "").strip(),
        "request_id": str(trace.get("request_id") or "").strip(),
        "executor_id": str(trace.get("executor_id") or "").strip(),
        "executor_kind": str(trace.get("executor_kind") or "").strip(),
        "transition_id": str(trace.get("transition_id") or details.get("transition_id") or "").strip(),
        "marking_before_ref": str(trace.get("marking_before_ref") or details.get("marking_before_ref") or "").strip(),
        "marking_after_ref": str(trace.get("marking_after_ref") or details.get("marking_after_ref") or "").strip(),
        "enabled_transition_ids": [
            str(item)
            for item in (trace.get("enabled_transition_ids") or details.get("enabled_transition_ids") or [])
            if isinstance(item, str) and item.strip()
        ],
        "topic_sequence": _serialize_trace_topic_sequence(trace.get("topic_sequence")),
        "timeline": _serialize_trace_timeline(trace.get("timeline")),
        "raw_events": _serialize_raw_events(trace.get("raw_events")),
        "trace_id": str(trace.get("trace_id") or "").strip(),
        "trace_fragment_id": str(trace.get("trace_fragment_id") or trace.get("committed_trace_fragment_id") or "").strip(),
    }

    has_value = any(
        [
            payload["workflow_id"],
            payload["request_id"],
            payload["executor_id"],
            payload["executor_kind"],
            payload["transition_id"],
            payload["marking_before_ref"],
            payload["marking_after_ref"],
            payload["enabled_transition_ids"],
            payload["topic_sequence"],
            payload["timeline"],
            payload["raw_events"],
            payload["trace_id"],
            payload["trace_fragment_id"],
        ]
    )
    if not has_value:
        return None
    return payload


def _serialize_step_detail(
    *,
    event: DebugStepEvent,
    state: DebugRunState,
    step_outputs: dict,
    runtime_context: dict,
    run_record: BenchmarkRunRecord | None,
    events: list[DebugStepEvent] | None = None,
    step_output_snapshots: dict[str, Any] | None = None,
) -> dict:
    """Build canonical step detail response for the selected event."""
    raw_step_name = event.step_name
    step_name = raw_step_name
    legacy_step_name = _implementation_step_name(state.pipeline_id, raw_step_name)

    source_bytes = runtime_context.get("source_bytes", b"")
    artifact_id = str(runtime_context.get("artifact_id") or f"artifact:{event.run_id}")

    inputs: dict[str, Any] = {
        "artifact_ids": [artifact_id],
        "fragment_ids": [],
        "program_ids": [],
    }
    outputs: dict[str, Any] = {
        "artifact_ids": [],
        "fragment_ids": [],
        "program_ids": [],
        "pair_ids": [],
    }
    operation_ref = "modelado/operators/noop"
    operation_params: dict[str, Any] = {}
    checks: list[dict[str, Any]] = []
    why_summary = "Step executed with canonical debug pipeline semantics."
    why_policy_name = step_name
    why_policy_params: dict[str, Any] = {}

    environment_fragments = STORE.list_environment_fragments(
        run_id=event.run_id,
        env_type=event.env_type,
        env_id=event.env_id,
    )
    scope_ref = _legacy_scope_ref(event.env_type, event.env_id)

    raw_manifest = runtime_context.get("asset_manifest")
    manifest = list(raw_manifest) if isinstance(raw_manifest, list) else []
    resolved_mime_by_artifact_id: dict[str, str] = {}
    if isinstance(step_outputs.get("asset_decomposition_statuses"), list):
        for item in step_outputs.get("asset_decomposition_statuses", []):
            if not isinstance(item, dict):
                continue
            artifact_ref = str(item.get("artifact_id") or "")
            mime_ref = str(item.get("mime_type") or "")
            if artifact_ref and mime_ref:
                resolved_mime_by_artifact_id[artifact_ref] = mime_ref
    artifact_file_names: dict[str, str] = {}
    if manifest:
        inputs["artifact_ids"] = [str(item.get("artifact_id")) for item in manifest if item.get("artifact_id")]
        artifact_file_names = {
            str(item.get("artifact_id")): str(item.get("filename"))
            for item in manifest
            if item.get("artifact_id") and item.get("filename")
        }
        lineage_root_ids = [f"artifact:{str(item.get('artifact_id'))}" for item in manifest if item.get("artifact_id")]
        lineage_nodes = [
            {
                "node_id": f"artifact:{str(item.get('artifact_id'))}",
                "kind": "artifact",
                "fragment_id": str(item.get("artifact_id")),
                "cas_id": None,
                "mime_type": resolved_mime_by_artifact_id.get(str(item.get("artifact_id") or ""), item.get("mime_type")),
                "label": str(item.get("filename") or item.get("artifact_id")),
                "meta": {
                    "record_type": "artifact",
                    "filename": item.get("filename"),
                    "size_bytes": item.get("size_bytes"),
                    "value_preview": None,
                },
            }
            for item in manifest
            if item.get("artifact_id")
        ]

        for node in lineage_nodes:
            if not isinstance(node, dict) or node.get("kind") != "artifact":
                continue
            fragment_id = str(node.get("fragment_id") or "")
            if not fragment_id:
                continue
            asset = _find_runtime_asset(runtime_context, fragment_id)
            if asset is None:
                continue
            try:
                preview_payload = _build_artifact_preview(asset)
            except Exception:
                continue
            meta = node.get("meta") if isinstance(node.get("meta"), dict) else {}
            kind = str(preview_payload.get("kind") or "")
            preview = preview_payload.get("preview") if isinstance(preview_payload.get("preview"), dict) else {}
            value_preview: Any = None
            if kind == "text":
                text_value = preview.get("text")
                if isinstance(text_value, str):
                    value_preview = text_value
            elif kind == "doc":
                paragraphs = preview.get("paragraphs")
                if isinstance(paragraphs, list):
                    snippet = [p for p in paragraphs if isinstance(p, str) and p.strip()]
                    if snippet:
                        value_preview = "\n\n".join(snippet[:8])
            elif kind == "json":
                parsed = preview.get("parsed")
                if isinstance(parsed, dict):
                    value_preview = parsed
            if value_preview is not None:
                meta["value_preview"] = value_preview
                node["meta"] = meta
    else:
        lineage_root_ids = [f"artifact:{artifact_id}"]
        lineage_nodes = [
            {
                "node_id": f"artifact:{artifact_id}",
                "kind": "artifact",
                "fragment_id": artifact_id,
                "cas_id": None,
                "mime_type": runtime_context.get("mime_type"),
                "label": artifact_id,
                "meta": {"record_type": "artifact"},
            }
        ]
    lineage_edges: list[dict[str, Any]] = []

    seen_fragment_ids: set[str] = set()
    for payload in environment_fragments:
        fragment_id = str(payload.get("id") or payload.get("cas_id") or "")
        if not fragment_id:
            continue
        if fragment_id in seen_fragment_ids:
            continue
        seen_fragment_ids.add(fragment_id)
        meta = payload.get("meta") if isinstance(payload.get("meta"), dict) else {}
        raw_value = payload.get("value")
        payload_mime = payload.get("mime_type")
        value_preview: Any = None
        if isinstance(raw_value, dict):
            if isinstance(raw_value.get("text"), str):
                text_value = raw_value.get("text")
                value_preview = text_value if len(text_value) <= 12000 else f"{text_value[:12000]}\n...[truncated]"
            elif payload_mime == "application/vnd.ikam.claim-ir+json" or all(key in raw_value for key in ("subject", "predicate", "object")):
                value_preview = raw_value
            else:
                value_preview = {"keys": sorted(raw_value.keys())[:16]}
        elif isinstance(raw_value, str):
            value_preview = raw_value if len(raw_value) <= 12000 else f"{raw_value[:12000]}\n...[truncated]"
        lineage_nodes.append(
            {
                "node_id": f"fragment:{fragment_id}",
                "kind": _node_kind_from_payload(payload),
                "fragment_id": fragment_id,
                "cas_id": payload.get("cas_id"),
                "mime_type": payload.get("mime_type"),
                "label": fragment_id,
                "meta": {
                    **meta,
                    "value_preview": value_preview,
                },
            }
        )
        lineage_edges.append(
            {
                "from": f"artifact:{meta.get('artifact_id') if isinstance(meta.get('artifact_id'), str) else artifact_id}",
                "to": f"fragment:{fragment_id}",
                "relation": "contains",
                "step_name": step_name,
            }
        )

    def _fragment_preview(fragment) -> dict[str, object]:
        fragment_id = getattr(fragment, "id", None) or getattr(fragment, "cas_id", None)
        cas_id = getattr(fragment, "cas_id", None)
        mime_type = getattr(fragment, "mime_type", None)
        value = getattr(fragment, "value", None)
        value_keys = sorted(value.keys()) if isinstance(value, dict) else []
        if isinstance(value, dict):
            if isinstance(value.get("text"), str):
                value_preview = value.get("text")
            elif mime_type == "application/vnd.ikam.claim-ir+json" or all(key in value for key in ("subject", "predicate", "object")):
                value_preview = value
            else:
                value_preview = {"keys": value_keys[:16]}
        elif isinstance(value, str):
            value_preview = value
        else:
            value_preview = None
        return {
            "fragment_id": fragment_id,
            "cas_id": cas_id,
            "mime_type": mime_type,
            "value_keys": value_keys,
            "value_preview": value_preview,
        }

    if raw_step_name == "init.initialize":
        raw_manifest = runtime_context.get("asset_manifest")
        manifest = list(raw_manifest) if isinstance(raw_manifest, list) else []
        if not manifest:
            manifest.append(
                {
                    "artifact_id": artifact_id,
                    "filename": artifact_id,
                    "mime_type": str(runtime_context.get("mime_type", "")),
                    "size_bytes": len(source_bytes) if isinstance(source_bytes, (bytes, bytearray)) else 0,
                    "role": "source",
                    "included": True,
                    "exclusion_reason": None,
                }
            )

        outputs["artifact_ids"] = [item["artifact_id"] for item in manifest]
        outputs["artifact_manifest"] = manifest
        outputs["artifact_count"] = len(manifest)
        outputs["total_size_bytes"] = sum(int(item.get("size_bytes") or 0) for item in manifest)
        type_counts: dict[str, int] = {}
        for item in manifest:
            mime = str(item.get("mime_type") or "application/octet-stream")
            type_counts[mime] = type_counts.get(mime, 0) + 1
        outputs["artifact_type_counts"] = type_counts
        lineage_root_ids = []
        lineage_nodes = []
        for item in manifest:
            aid = str(item.get("artifact_id") or "")
            if not aid:
                continue
            node_id = f"artifact:{aid}"
            lineage_root_ids.append(node_id)
            lineage_nodes.append(
                {
                    "node_id": node_id,
                    "kind": "artifact",
                    "fragment_id": aid,
                    "cas_id": None,
                    "mime_type": item.get("mime_type"),
                    "label": str(item.get("filename") or aid),
                    "meta": {
                        "record_type": "artifact",
                        "filename": item.get("filename"),
                        "size_bytes": item.get("size_bytes"),
                    },
                }
            )
        checks.append({"name": "case_registry_loaded", "status": "pass", "details": {"artifact_count": len(manifest)}})
        checks.append({"name": "artifact_manifest_complete", "status": "pass", "details": {"all_included": True}})
        why_summary = "Loaded the case fixture and prepared full artifact manifest for incremental ingestion."
        why_policy_params = {"artifact_count": len(manifest), "total_size_bytes": outputs["total_size_bytes"]}

    if isinstance(event.metrics, dict):
        metric_operation_ref = event.metrics.get("operation_ref")
        metric_operation_params = event.metrics.get("operation_params")
        if isinstance(metric_operation_ref, str) and metric_operation_ref.strip():
            operation_ref = metric_operation_ref
        if isinstance(metric_operation_params, dict):
            operation_params = metric_operation_params
    if operation_ref == "modelado/operators/noop":
        operation_ref = "modelado/operators/unknown"

    produced_fragment_ids: list[str] = []
    raw_fragment_ids = outputs.get("fragment_ids")
    if isinstance(raw_fragment_ids, list):
        produced_fragment_ids.extend([str(item) for item in raw_fragment_ids if isinstance(item, str) and item])
    lifted = step_outputs.get("lifted")
    if isinstance(lifted, list):
        for item in lifted:
            if not isinstance(item, dict):
                continue
            item_id = item.get("id")
            if isinstance(item_id, str) and item_id:
                produced_fragment_ids.append(item_id)
    produced_fragment_ids = list(dict.fromkeys(produced_fragment_ids))

    # Surface any captured debug sink events and operation telemetry
    debug_sink_events = step_outputs.get("debug_sink_events", [])
    if debug_sink_events:
        outputs["debug_events"] = debug_sink_events

    if isinstance(event.metrics, dict):
        details = event.metrics.get("details", {})
        if isinstance(details, dict) and "operation_telemetry" in details:
            outputs["operation_telemetry"] = details["operation_telemetry"]
        elif "operation_telemetry" in step_outputs:
            outputs["operation_telemetry"] = step_outputs["operation_telemetry"]

    if isinstance(outputs.get("operation_telemetry"), dict):
        operation_telemetry = dict(outputs["operation_telemetry"])
        executor_identity = _event_executor_identity(event)
        operation_telemetry.setdefault("executor_id", executor_identity.get("executor_id") or "")
        operation_telemetry.setdefault("executor_kind", executor_identity.get("executor_kind") or "")
        operation_telemetry.setdefault("operation_name", _presentation_step_name(event.pipeline_id, legacy_step_name))
        outputs["operation_telemetry"] = operation_telemetry

    if isinstance(step_outputs.get("documents"), list):
        outputs["documents"] = [item for item in step_outputs["documents"] if isinstance(item, dict)]
    if isinstance(step_outputs.get("document_fragment_refs"), list):
        outputs["document_fragment_refs"] = [item for item in step_outputs["document_fragment_refs"] if isinstance(item, str)]
    if isinstance(step_outputs.get("document_loads"), list):
        outputs["document_loads"] = [item for item in step_outputs["document_loads"] if isinstance(item, dict)]
    if isinstance(step_outputs.get("fragment_ids"), list) and not outputs.get("fragment_ids"):
        outputs["fragment_ids"] = [item for item in step_outputs["fragment_ids"] if isinstance(item, str) and item]
    if isinstance(step_outputs.get("program_ids"), list) and not outputs.get("program_ids"):
        outputs["program_ids"] = [item for item in step_outputs["program_ids"] if isinstance(item, str) and item]
    if isinstance(step_outputs.get("artifact_ids"), list) and not outputs.get("artifact_ids"):
        outputs["artifact_ids"] = [item for item in step_outputs["artifact_ids"] if isinstance(item, str) and item]

    produced_fragment_ids = list(
        dict.fromkeys(
            [
                *[str(item) for item in outputs.get("fragment_ids", []) if isinstance(item, str) and item],
                *produced_fragment_ids,
            ]
        )
    )

    if legacy_step_name == "map.conceptual.lift.surface_fragments":
        decomposition = step_outputs.get("decomposition")
        fragment_artifact_map = step_outputs.get("fragment_artifact_map") if isinstance(step_outputs.get("fragment_artifact_map"), dict) else {}
        artifact_metrics = step_outputs.get("decomposition_artifact_metrics") if isinstance(step_outputs.get("decomposition_artifact_metrics"), list) else []
        reconstruction_verdicts = step_outputs.get("reconstruction_verdicts") if isinstance(step_outputs.get("reconstruction_verdicts"), list) else []
        structural_map = step_outputs.get("structural_map") if isinstance(step_outputs.get("structural_map"), dict) else {}
        map_dna = step_outputs.get("map_dna") if isinstance(step_outputs.get("map_dna"), dict) else {}
        map_outline_nodes = step_outputs.get("map_outline_nodes") if isinstance(step_outputs.get("map_outline_nodes"), list) else []
        map_root_node_id = step_outputs.get("map_root_node_id") if isinstance(step_outputs.get("map_root_node_id"), str) else ""
        map_node_summaries = step_outputs.get("map_node_summaries") if isinstance(step_outputs.get("map_node_summaries"), dict) else {}
        map_node_constituents = step_outputs.get("map_node_constituents") if isinstance(step_outputs.get("map_node_constituents"), dict) else {}
        map_relationships = step_outputs.get("map_relationships") if isinstance(step_outputs.get("map_relationships"), list) else []
        map_subgraph = step_outputs.get("map_subgraph") if isinstance(step_outputs.get("map_subgraph"), dict) else {}
        map_segment_anchors = step_outputs.get("map_segment_anchors") if isinstance(step_outputs.get("map_segment_anchors"), dict) else {}
        map_segment_candidates = step_outputs.get("map_segment_candidates") if isinstance(step_outputs.get("map_segment_candidates"), list) else []
        map_profile_candidates = step_outputs.get("map_profile_candidates") if isinstance(step_outputs.get("map_profile_candidates"), dict) else {}
        map_generation_provenance = step_outputs.get("map_generation_provenance") if isinstance(step_outputs.get("map_generation_provenance"), dict) else {}
        map_agent_review = step_outputs.get("map_agent_review") if isinstance(step_outputs.get("map_agent_review"), dict) else {}
        map_elicitation = step_outputs.get("map_elicitation") if isinstance(step_outputs.get("map_elicitation"), dict) else {}
        map_agent_spec = step_outputs.get("map_agent_spec") if isinstance(step_outputs.get("map_agent_spec"), dict) else {}
        map_status = str(step_outputs.get("map_status") or "ok")
        map_error_reason = step_outputs.get("map_error_reason")
        fragments: list[dict[str, Any]] = []
        root_fragment_ids: list[str] = []
        reconstruction_fragments: list[dict[str, Any]] = []
        for frag in getattr(decomposition, "root_fragments", []):
            fragment_id = getattr(frag, "id", None) or getattr(frag, "cas_id", None)
            if isinstance(fragment_id, str):
                root_fragment_ids.append(fragment_id)
                reconstruction_fragments.append(
                    {
                        "fragment_id": fragment_id,
                        "mime_type": getattr(frag, "mime_type", None),
                        "is_surface": bool(fragment_id in (set(outputs["fragment_ids"]) if outputs["fragment_ids"] else set())),
                    }
                )
        for frag in getattr(decomposition, "structural", []):
            preview = _fragment_preview(frag)
            fragments.append(preview)
            if isinstance(preview.get("fragment_id"), str):
                outputs["fragment_ids"].append(preview["fragment_id"])
        structural_ids = list(outputs["fragment_ids"])
        structural_id_set = set(structural_ids)
        if reconstruction_fragments:
            for item in reconstruction_fragments:
                item["is_surface"] = bool(item.get("fragment_id") in structural_id_set)
        if not root_fragment_ids:
            root_fragment_ids = structural_ids
        links: list[dict[str, str]] = []
        for fragment_id in outputs["fragment_ids"]:
            artifact_ref = fragment_artifact_map.get(str(fragment_id)) if fragment_artifact_map else None
            if artifact_ref:
                links.append({"artifact_id": str(artifact_ref), "surface_fragment_id": str(fragment_id)})
        outputs["decomposition"] = {
            "root_fragment_ids": root_fragment_ids,
            "structural_fragment_ids": structural_ids,
            "reconstruction_fragments": reconstruction_fragments,
            "surface_fragments": fragments,
            "artifact_surface_links": links,
            "artifact_count": len(inputs.get("artifact_ids", [])) or (len({item["artifact_id"] for item in links}) if links else 1),
            "asset_statuses": step_outputs.get("asset_decomposition_statuses", []),
            "boundary_diagnostics": [
                build_boundary_diagnostic(item)
                for item in artifact_metrics
                if isinstance(item, dict)
            ],
            "reconstruction_verdicts": [item for item in reconstruction_verdicts if isinstance(item, dict)],
        }
        outputs["map"] = {
            "status": map_status,
            "map_subgraph": map_subgraph,
            "structural_map": structural_map,
            "map_dna": map_dna,
            "outline_nodes": [node for node in map_outline_nodes if isinstance(node, dict)],
            "root_node_id": map_root_node_id,
            "node_summaries": {
                str(node_id): str(summary)
                for node_id, summary in map_node_summaries.items()
                if isinstance(node_id, str) and isinstance(summary, str)
            },
            "node_constituents": {
                str(node_id): [str(fragment_id) for fragment_id in fragment_ids if isinstance(fragment_id, (str, int))]
                for node_id, fragment_ids in map_node_constituents.items()
                if isinstance(node_id, str) and isinstance(fragment_ids, list)
            },
            "relationships": [item for item in map_relationships if isinstance(item, dict)],
            "segment_anchors": {
                str(segment_id): [item for item in anchors if isinstance(item, dict)]
                for segment_id, anchors in map_segment_anchors.items()
                if isinstance(segment_id, str) and isinstance(anchors, list)
            },
            "segment_candidates": [item for item in map_segment_candidates if isinstance(item, dict)],
            "profile_candidates": {
                str(segment_id): [str(profile) for profile in profiles if isinstance(profile, (str, int))]
                for segment_id, profiles in map_profile_candidates.items()
                if isinstance(segment_id, str) and isinstance(profiles, list)
            },
            "generation_provenance": map_generation_provenance,
            "agent_review": map_agent_review,
            "elicitation": map_elicitation,
            "agent_spec": map_agent_spec,
            "preview_mode_default": "semantic_map",
            "error_reason": str(map_error_reason) if map_error_reason is not None else None,
        }
        map_nodes_by_id: dict[str, dict[str, Any]] = {
            str(node.get("id")): node
            for node in map_outline_nodes
            if isinstance(node, dict) and isinstance(node.get("id"), str)
        }
        existing_node_ids = {
            str(node.get("node_id"))
            for node in lineage_nodes
            if isinstance(node, dict) and isinstance(node.get("node_id"), str)
        }
        existing_edges = {
            (
                str(edge.get("from")),
                str(edge.get("to")),
                str(edge.get("relation")),
            )
            for edge in lineage_edges
            if isinstance(edge, dict)
        }
        for node_id, node_payload in map_nodes_by_id.items():
            if node_id not in existing_node_ids:
                parent_id = node_payload.get("parent_id") if isinstance(node_payload.get("parent_id"), str) else None
                lineage_nodes.append(
                    {
                        "node_id": node_id,
                        "kind": "map_root" if node_id == map_root_node_id else "map_node",
                        "fragment_id": None,
                        "cas_id": None,
                        "mime_type": "application/vnd.ikam.map-node+json",
                        "label": str(node_payload.get("title") or node_id),
                        "meta": {
                            "record_type": "map_node",
                            "map_root": bool(node_id == map_root_node_id),
                            "kind": str(node_payload.get("kind") or "unknown"),
                            "level": int(node_payload.get("level") or 0),
                            "parent_id": parent_id,
                            "artifact_ids": [str(aid) for aid in (node_payload.get("artifact_ids") or []) if isinstance(aid, (str, int))],
                        },
                    }
                )
                existing_node_ids.add(node_id)

            parent_id = node_payload.get("parent_id") if isinstance(node_payload.get("parent_id"), str) else None
            if parent_id:
                edge_key = (parent_id, node_id, "map_contains")
                if edge_key not in existing_edges:
                    lineage_edges.append({"from": parent_id, "to": node_id, "relation": "map_contains", "step_name": step_name})
                    existing_edges.add(edge_key)

            for artifact_ref in node_payload.get("artifact_ids") or []:
                artifact_id_ref = str(artifact_ref)
                edge_key = (node_id, f"artifact:{artifact_id_ref}", "map_to_artifact")
                if edge_key not in existing_edges:
                    lineage_edges.append(
                        {
                            "from": node_id,
                            "to": f"artifact:{artifact_id_ref}",
                            "relation": "map_to_artifact",
                            "step_name": step_name,
                        }
                    )
                    existing_edges.add(edge_key)

            node_kind = str(node_payload.get("kind") or "")
            if node_kind == "surface_fragment":
                fragment_id = None
                if node_id.startswith("map:surface:"):
                    fragment_id = node_id[len("map:surface:") :]
                if fragment_id:
                    edge_key = (node_id, f"fragment:{fragment_id}", "map_to_surface")
                    if edge_key not in existing_edges:
                        lineage_edges.append(
                            {
                                "from": node_id,
                                "to": f"fragment:{fragment_id}",
                                "relation": "map_to_surface",
                                "step_name": step_name,
                            }
                        )
                        existing_edges.add(edge_key)

        if map_root_node_id and map_root_node_id in map_nodes_by_id and map_root_node_id not in lineage_root_ids:
            lineage_root_ids.append(map_root_node_id)
        if inputs.get("artifact_ids"):
            linked_artifacts = {item["artifact_id"] for item in links}
            outputs["decomposition"]["artifacts_without_surface_fragments"] = [
                aid for aid in inputs["artifact_ids"] if aid not in linked_artifacts
            ]
        checks.append(
            {
                "name": "decomposition_nonempty",
                "status": "pass" if len(outputs["fragment_ids"]) > 0 else "fail",
                "details": {"fragment_count": len(outputs["fragment_ids"])},
            }
        )
        checks.append({"name": "artifact_contains_edges_emitted", "status": "pass", "details": {"edge_count": len(outputs["fragment_ids"])}})
        checks.append(
            {
                "name": "map_render_equivalence",
                "status": "fail" if any((isinstance(item, dict) and item.get("status") == "fail") for item in reconstruction_verdicts) else "pass",
                "details": {
                    "verdict_count": len(reconstruction_verdicts),
                    "failures": [
                        item for item in reconstruction_verdicts if isinstance(item, dict) and item.get("status") == "fail"
                    ],
                },
            }
        )
        checks.append(
            {
                "name": "map_structural_contract",
                "status": "pass"
                if isinstance(outputs["map"].get("structural_map"), dict)
                and isinstance(outputs["map"].get("map_dna"), dict)
                and isinstance(outputs["map"].get("root_node_id"), str)
                else "fail",
                "details": {
                    "map_status": outputs["map"].get("status"),
                    "outline_node_count": len(outputs["map"].get("outline_nodes") or []),
                },
            }
        )
        checks.append(
            {
                "name": "map_generation_provenance_present",
                "status": "pass"
                if isinstance(outputs["map"].get("generation_provenance"), dict)
                and bool((outputs["map"].get("generation_provenance") or {}).get("provider"))
                and bool((outputs["map"].get("generation_provenance") or {}).get("model"))
                else "fail",
                "details": {
                    "provenance": outputs["map"].get("generation_provenance"),
                },
            }
        )
        map_outline_ids = {
            str(node.get("id"))
            for node in outputs["map"].get("outline_nodes") or []
            if isinstance(node, dict) and isinstance(node.get("id"), str)
        }
        non_leaf_outline_ids = {
            str(node.get("id"))
            for node in outputs["map"].get("outline_nodes") or []
            if isinstance(node, dict)
            and isinstance(node.get("id"), str)
            and str(node.get("kind") or "") != "surface_fragment"
        }
        map_summaries = outputs["map"].get("node_summaries")
        summary_covered = {
            node_id
            for node_id in non_leaf_outline_ids
            if isinstance(map_summaries, dict)
            and isinstance(map_summaries.get(node_id), str)
            and str(map_summaries.get(node_id)).strip()
        }
        checks.append(
            {
                "name": "map_node_summary_coverage",
                "status": "pass" if non_leaf_outline_ids.issubset(summary_covered) else "fail",
                "details": {
                    "required": sorted(non_leaf_outline_ids),
                    "covered": sorted(summary_covered),
                },
            }
        )
        map_constituents = outputs["map"].get("node_constituents")
        constituent_keys = {
            node_id
            for node_id, fragment_ids in map_constituents.items()
            if isinstance(node_id, str)
            and isinstance(fragment_ids, list)
            and all(isinstance(fragment_id, str) for fragment_id in fragment_ids)
        } if isinstance(map_constituents, dict) else set()
        checks.append(
            {
                "name": "map_constituent_links_complete",
                "status": "pass" if map_outline_ids.issubset(constituent_keys) else "fail",
                "details": {
                    "required": sorted(map_outline_ids),
                    "covered": sorted(constituent_keys),
                },
            }
        )
        max_outline_level = max(
            [
                int(node.get("level") or 0)
                for node in outputs["map"].get("outline_nodes") or []
                if isinstance(node, dict)
            ]
            or [0]
        )
        checks.append(
            {
                "name": "map_semantic_outline_nontrivial",
                "status": "pass" if len(outputs["map"].get("outline_nodes") or []) >= 3 and max_outline_level >= 1 else "fail",
                "details": {
                    "outline_node_count": len(outputs["map"].get("outline_nodes") or []),
                    "max_outline_level": max_outline_level,
                },
            }
        )
        why_summary = "Decomposed case artifacts into surface fragments and emitted artifact-to-fragment links."
        why_policy_params = {
            "surface_fragment_count": len(outputs["fragment_ids"]),
            "artifact_count": outputs["decomposition"]["artifact_count"],
        }

    elif legacy_step_name == "map.conceptual.embed.discovery_index":
        surface_embeddings = step_outputs.get("surface_embeddings", {})
        surface_clusters = step_outputs.get("surface_clusters", [])
        details = event.metrics.get("details") if isinstance(event.metrics, dict) else {}
        cluster_threshold = float(details.get("cluster_threshold")) if isinstance(details, dict) and isinstance(details.get("cluster_threshold"), (int, float)) else 0.7
        embedded_ids = sorted(surface_embeddings.keys()) if isinstance(surface_embeddings, dict) else []
        outputs["fragment_ids"] = embedded_ids
        outputs["embedding_dimensions"] = 0
        if isinstance(surface_embeddings, dict) and surface_embeddings:
            first_val = next(iter(surface_embeddings.values()))
            outputs["embedding_dimensions"] = len(first_val) if isinstance(first_val, (list, tuple)) else 0
        outputs["clusters"] = []
        cluster_members_by_id: dict[str, list[str]] = {}
        cluster_similarity_by_id: dict[str, float | None] = {}
        if isinstance(surface_clusters, list):
            for index, cluster in enumerate(surface_clusters):
                if isinstance(cluster, dict):
                    members_raw = cluster.get("members", [])
                    members = [str(member) for member in members_raw] if isinstance(members_raw, (list, tuple, set)) else []
                    centroid_id = str(cluster.get("centroid_id")) if cluster.get("centroid_id") is not None else None
                    avg_similarity = float(cluster.get("avg_similarity")) if isinstance(cluster.get("avg_similarity"), (int, float)) else None
                else:
                    members = [str(member) for member in cluster] if isinstance(cluster, (list, tuple, set)) else []
                    centroid_id = None
                    avg_similarity = None
                cluster_id = f"cluster-{index}"
                cluster_members_by_id[cluster_id] = members
                cluster_similarity_by_id[cluster_id] = avg_similarity
                outputs["clusters"].append(
                    {
                        "cluster_id": cluster_id,
                        "member_fragment_ids": members,
                        "centroid_id": centroid_id,
                        "avg_similarity": avg_similarity,
                    }
                )

        valid_vectors: dict[str, list[float]] = {}
        for fragment_id in embedded_ids:
            vector = surface_embeddings.get(fragment_id) if isinstance(surface_embeddings, dict) else None
            if isinstance(vector, (list, tuple)) and all(isinstance(v, (int, float)) for v in vector):
                valid_vectors[str(fragment_id)] = [float(v) for v in vector]

        ordered_ids = sorted(valid_vectors.keys())
        matrix: list[list[float]] = []
        if ordered_ids:
            import math

            norms = {fid: math.sqrt(sum(v * v for v in valid_vectors[fid])) or 1.0 for fid in ordered_ids}
            for source_id in ordered_ids:
                row: list[float] = []
                source_vec = valid_vectors[source_id]
                for target_id in ordered_ids:
                    target_vec = valid_vectors[target_id]
                    dot = sum(a * b for a, b in zip(source_vec, target_vec))
                    score = dot / (norms[source_id] * norms[target_id])
                    row.append(float(max(-1.0, min(1.0, score))))
                matrix.append(row)

        fragment_to_cluster: dict[str, str] = {}
        for cluster_id, members in cluster_members_by_id.items():
            for member in members:
                fragment_to_cluster.setdefault(member, cluster_id)

        points: list[dict[str, Any]] = []
        if ordered_ids:
            try:
                import numpy as np

                vectors = np.array([valid_vectors[fid] for fid in ordered_ids], dtype=float)
                centered = vectors - vectors.mean(axis=0, keepdims=True)
                _, _, vh = np.linalg.svd(centered, full_matrices=False)
                components = vh[:2]
                coords = centered @ components.T
                if coords.shape[1] == 1:
                    coords = np.column_stack([coords[:, 0], np.zeros(coords.shape[0])])
                elif coords.shape[1] == 0:
                    coords = np.zeros((coords.shape[0], 2))
                for idx, fragment_id in enumerate(ordered_ids):
                    points.append(
                        {
                            "fragment_id": fragment_id,
                            "x": float(coords[idx, 0]),
                            "y": float(coords[idx, 1]),
                            "cluster_id": fragment_to_cluster.get(fragment_id),
                        }
                    )
                projection_method = "pca_2d"
            except Exception:
                count = len(ordered_ids)
                for idx, fragment_id in enumerate(ordered_ids):
                    x = float(idx) / float(max(1, count - 1))
                    points.append(
                        {
                            "fragment_id": fragment_id,
                            "x": x,
                            "y": 0.0,
                            "cluster_id": fragment_to_cluster.get(fragment_id),
                        }
                    )
                projection_method = "pca_2d"
        else:
            projection_method = "pca_2d"

        point_lookup = {str(item["fragment_id"]): item for item in points}
        centroids: list[dict[str, Any]] = []
        for cluster_id, members in cluster_members_by_id.items():
            cluster_points = [point_lookup[mid] for mid in members if mid in point_lookup]
            if not cluster_points:
                continue
            x = sum(float(item["x"]) for item in cluster_points) / len(cluster_points)
            y = sum(float(item["y"]) for item in cluster_points) / len(cluster_points)
            centroids.append(
                {
                    "cluster_id": cluster_id,
                    "x": x,
                    "y": y,
                    "size": len(cluster_points),
                    "avg_similarity": cluster_similarity_by_id.get(cluster_id),
                }
            )

        decomposition = step_outputs.get("decomposition")
        expected_fragment_ids: list[str] = []
        if decomposition is not None:
            for fragment in getattr(decomposition, "structural", []):
                fragment_id = getattr(fragment, "id", None) or getattr(fragment, "cas_id", None)
                if isinstance(fragment_id, str):
                    expected_fragment_ids.append(fragment_id)
        if not expected_fragment_ids:
            expected_fragment_ids = list(ordered_ids)
        expected_set = set(expected_fragment_ids)
        embedded_set = set(ordered_ids)
        missing_fragment_ids = sorted(expected_set - embedded_set)

        singleton_clusters = sum(1 for members in cluster_members_by_id.values() if len(members) <= 1)
        matrix_values = [value for row in matrix for value in row]

        mode = details["embedding_mode"] if isinstance(details, dict) and isinstance(details.get("embedding_mode"), str) else None

        outputs["embedding_projection"] = {
            "method": projection_method,
            "points": points,
            "centroids": centroids,
        }
        outputs["pairwise_similarity"] = {
            "fragment_ids": ordered_ids,
            "matrix": matrix,
            "min": min(matrix_values) if matrix_values else None,
            "max": max(matrix_values) if matrix_values else None,
            "threshold": cluster_threshold,
        }
        outputs["embedding_debug"] = {
            "expected_count": len(expected_set),
            "embedded_count": len(embedded_set),
            "coverage_ratio": float(len(embedded_set) / len(expected_set)) if expected_set else 1.0,
            "missing_fragment_ids": missing_fragment_ids,
            "singleton_clusters": singleton_clusters,
            "cluster_count": len(cluster_members_by_id),
            "threshold": cluster_threshold,
            "embedding_mode": mode,
        }

        if len([node for node in lineage_nodes if isinstance(node, dict) and str(node.get("node_id", "")).startswith("fragment:")]) == 0:
            decomposition = step_outputs.get("decomposition")
            fragment_artifact_map = step_outputs.get("fragment_artifact_map") if isinstance(step_outputs.get("fragment_artifact_map"), dict) else {}
            for fragment in getattr(decomposition, "structural", []):
                preview = _fragment_preview(fragment)
                fragment_id = preview.get("fragment_id")
                if not isinstance(fragment_id, str):
                    continue
                source_artifact_id = str(fragment_artifact_map.get(fragment_id) or artifact_id)
                line_meta = {
                    "record_type": "surface_fragment",
                    "artifact_id": source_artifact_id,
                    "value_preview": preview.get("value_preview"),
                }
                if source_artifact_id in artifact_file_names:
                    line_meta["file_name"] = artifact_file_names[source_artifact_id]
                lineage_nodes.append(
                    {
                        "node_id": f"fragment:{fragment_id}",
                        "kind": "surface",
                        "fragment_id": fragment_id,
                        "cas_id": preview.get("cas_id"),
                        "mime_type": preview.get("mime_type"),
                        "label": fragment_id,
                        "meta": line_meta,
                    }
                )
                lineage_edges.append(
                    {
                        "from": f"artifact:{source_artifact_id}",
                        "to": f"fragment:{fragment_id}",
                        "relation": "contains",
                        "step_name": step_name,
                    }
                )
        checks.append(
            {
                "name": "surface_embedding_coverage",
                "status": "pass" if len(embedded_ids) > 0 else "warn",
                "details": {"embedded_count": len(embedded_ids)},
            }
        )
        why_summary = "Embedded mapped surface fragments to prepare candidate matching." 
        why_policy_params = {"embedding_count": len(embedded_ids), "dimensions": outputs["embedding_dimensions"]}

    elif legacy_step_name == "map.conceptual.lift.entities_and_relationships":
        entity_relationships = step_outputs.get("entity_relationships") if isinstance(step_outputs.get("entity_relationships"), list) else []
        entity_relationship_set = step_outputs.get("entity_relationship_set") if isinstance(step_outputs.get("entity_relationship_set"), dict) else {}
        summary = step_outputs.get("summary") if isinstance(step_outputs.get("summary"), dict) else {}
        outputs["fragment_ids"] = [str(item) for item in (step_outputs.get("fragment_ids") or []) if isinstance(item, str) and item]
        outputs["entity_relationships"] = [item for item in entity_relationships if isinstance(item, dict)]
        outputs["entity_relationship_set"] = entity_relationship_set
        outputs["summary"] = summary
        checks.append(
            {
                "name": "entity_relationships_emitted",
                "status": "pass" if len(outputs["fragment_ids"]) > 0 else "warn",
                "details": {"entity_relationship_count": len(outputs["fragment_ids"])},
            }
        )
        why_summary = "Derived entity and relationship rows from chunk extractions for downstream claims parsing."
        why_policy_params = {"entity_relationship_count": len(outputs["fragment_ids"])}

    elif legacy_step_name == "map.conceptual.lift.claims":
        claims = step_outputs.get("claims") if isinstance(step_outputs.get("claims"), list) else []
        claim_set = step_outputs.get("claim_set") if isinstance(step_outputs.get("claim_set"), dict) else {}
        summary = step_outputs.get("summary") if isinstance(step_outputs.get("summary"), dict) else {}
        outputs["fragment_ids"] = [str(item) for item in (step_outputs.get("fragment_ids") or []) if isinstance(item, str) and item]
        outputs["claims"] = [item for item in claims if isinstance(item, dict)]
        outputs["claim_set"] = claim_set
        outputs["summary"] = summary
        checks.append(
            {
                "name": "claims_emitted",
                "status": "pass" if len(outputs["fragment_ids"]) > 0 else "warn",
                "details": {"claim_count": len(outputs["fragment_ids"])} ,
            }
        )
        why_summary = "Derived claim rows from entity relationship rows for downstream claim-set inspection and persistence."
        why_policy_params = {"claim_count": len(outputs["fragment_ids"])}

    elif legacy_step_name == "map.conceptual.normalize.discovery":
        ir_fragments = step_outputs.get("ir_fragments", [])
        lifted_from_map = step_outputs.get("lifted_from_map") if isinstance(step_outputs.get("lifted_from_map"), dict) else {}
        fragment_artifact_map = step_outputs.get("fragment_artifact_map") if isinstance(step_outputs.get("fragment_artifact_map"), dict) else {}
        decomposition = step_outputs.get("decomposition")
        ir_ids: list[str] = []
        for frag in ir_fragments:
            preview = _fragment_preview(frag)
            if isinstance(preview.get("fragment_id"), str):
                ir_ids.append(preview["fragment_id"])
        outputs["fragment_ids"] = ir_ids

        surface_to_ir: dict[str, list[str]] = {}
        for ir_id, surface_id in lifted_from_map.items():
            if ir_id is None or surface_id is None:
                continue
            surface_to_ir.setdefault(str(surface_id), []).append(str(ir_id))

        surface_ids: list[str] = []
        if decomposition is not None:
            for fragment in getattr(decomposition, "structural", []):
                fragment_id = getattr(fragment, "id", None) or getattr(fragment, "cas_id", None)
                if isinstance(fragment_id, str):
                    surface_ids.append(fragment_id)
        if not surface_ids:
            surface_ids = sorted(surface_to_ir.keys())

        transformations: list[dict[str, Any]] = []
        for surface_id in surface_ids:
            ir_targets = sorted(surface_to_ir.get(surface_id, []))
            source_artifact_id = fragment_artifact_map.get(surface_id)
            transformations.append(
                {
                    "surface_fragment_id": surface_id,
                    "source_artifact_id": str(source_artifact_id) if source_artifact_id else None,
                    "ir_fragment_ids": ir_targets,
                    "lift_status": "lifted" if ir_targets else "surface_only",
                    "lift_reason": None if ir_targets else "no_ir_generated",
                }
            )

        outputs["lift_map"] = [
            {
                "surface_fragment_id": item["surface_fragment_id"],
                "ir_fragment_ids": item["ir_fragment_ids"],
            }
            for item in transformations
        ]
        outputs["lift_transformations"] = transformations

        checks.append({"name": "lift_provenance_complete", "status": "pass", "details": {"ir_fragment_count": len(ir_ids)}})
        why_summary = "Lifted surface fragments into IR fragments for semantic operations."
        why_policy_params = {"ir_fragment_count": len(ir_ids)}

    elif legacy_step_name == "map.reconstructable.embed":
        embeddings = step_outputs.get("embeddings", {})
        outputs["fragment_ids"] = sorted(embeddings.keys()) if isinstance(embeddings, dict) else []
        outputs["embedding_dimensions"] = 0
        if isinstance(embeddings, dict) and embeddings:
            first_val = next(iter(embeddings.values()))
            outputs["embedding_dimensions"] = len(first_val) if isinstance(first_val, (list, tuple)) else 0

        details = event.metrics.get("details") if isinstance(event.metrics, dict) else {}
        cluster_threshold = float(details.get("cluster_threshold")) if isinstance(details, dict) and isinstance(details.get("cluster_threshold"), (int, float)) else 0.7

        valid_vectors: dict[str, list[float]] = {}
        if isinstance(embeddings, dict):
            for fragment_id, vector in embeddings.items():
                if isinstance(vector, (list, tuple)) and all(isinstance(v, (int, float)) for v in vector):
                    valid_vectors[str(fragment_id)] = [float(v) for v in vector]

        ordered_ids = sorted(valid_vectors.keys())
        matrix: list[list[float]] = []
        if ordered_ids:
            import math

            norms = {fid: math.sqrt(sum(v * v for v in valid_vectors[fid])) or 1.0 for fid in ordered_ids}
            for source_id in ordered_ids:
                row: list[float] = []
                source_vec = valid_vectors[source_id]
                for target_id in ordered_ids:
                    target_vec = valid_vectors[target_id]
                    dot = sum(a * b for a, b in zip(source_vec, target_vec))
                    score = dot / (norms[source_id] * norms[target_id])
                    row.append(float(max(-1.0, min(1.0, score))))
                matrix.append(row)

        matrix_values = [value for row in matrix for value in row]
        outputs["pairwise_similarity"] = {
            "fragment_ids": ordered_ids,
            "matrix": matrix,
            "min": min(matrix_values) if matrix_values else None,
            "max": max(matrix_values) if matrix_values else None,
            "threshold": cluster_threshold,
        }

        ir_fragments = step_outputs.get("ir_fragments") if isinstance(step_outputs.get("ir_fragments"), list) else []
        expected_ir_ids: list[str] = []
        for frag in ir_fragments:
            preview = _fragment_preview(frag)
            if isinstance(preview.get("fragment_id"), str):
                expected_ir_ids.append(preview["fragment_id"])

        expected_set = set(expected_ir_ids) if expected_ir_ids else set(ordered_ids)
        embedded_set = set(ordered_ids)
        missing_fragment_ids = sorted(expected_set - embedded_set)
        mode = details["embedding_mode"] if isinstance(details, dict) and isinstance(details.get("embedding_mode"), str) else None
        outputs["embedding_debug"] = {
            "expected_count": len(expected_set),
            "embedded_count": len(embedded_set),
            "coverage_ratio": float(len(embedded_set) / len(expected_set)) if expected_set else 1.0,
            "missing_fragment_ids": missing_fragment_ids,
            "singleton_clusters": len(ordered_ids),
            "cluster_count": len(ordered_ids),
            "threshold": cluster_threshold,
            "embedding_mode": mode,
        }

        existing_node_ids = {
            str(node.get("node_id"))
            for node in lineage_nodes
            if isinstance(node, dict) and isinstance(node.get("node_id"), str)
        }
        lifted_from_map = step_outputs.get("lifted_from_map") if isinstance(step_outputs.get("lifted_from_map"), dict) else {}
        for frag in ir_fragments:
            preview = _fragment_preview(frag)
            fragment_id = preview.get("fragment_id")
            if not isinstance(fragment_id, str):
                continue
            node_id = f"fragment:{fragment_id}"
            if node_id not in existing_node_ids:
                source_surface_id = lifted_from_map.get(fragment_id)
                line_meta = {
                    "record_type": "ir_fragment",
                    "source_surface_fragment_id": str(source_surface_id) if source_surface_id else None,
                    "artifact_id": artifact_id,
                    "value_preview": preview.get("value_preview"),
                }
                lineage_nodes.append(
                    {
                        "node_id": node_id,
                        "kind": "ir",
                        "fragment_id": fragment_id,
                        "cas_id": preview.get("cas_id"),
                        "mime_type": preview.get("mime_type"),
                        "label": fragment_id,
                        "meta": line_meta,
                    }
                )
                existing_node_ids.add(node_id)

        for ir_id, surface_id in lifted_from_map.items():
            if not ir_id or not surface_id:
                continue
            lineage_edges.append(
                {
                    "from": f"fragment:{surface_id}",
                    "to": f"fragment:{ir_id}",
                    "relation": "lifted_from",
                    "step_name": step_name,
                }
            )

        checks.append(
            {
                "name": "ir_embedding_coverage",
                "status": "pass" if len(outputs["fragment_ids"]) > 0 else "warn",
                "details": {"embedded_count": len(outputs["fragment_ids"])},
            }
        )
        why_summary = "Embedded lifted IR fragments for candidate search."
        why_policy_params = {"embedding_count": len(outputs["fragment_ids"]), "dimensions": outputs["embedding_dimensions"]}

    elif legacy_step_name == "map.reconstructable.search.dependency_resolution":
        candidates = step_outputs.get("candidates", [])
        outputs["candidates"] = []
        for index, c in enumerate(candidates):
            payload = c if isinstance(c, dict) else c.__dict__
            pair_id = f"pair-{index}"
            outputs["pair_ids"].append(pair_id)
            outputs["candidates"].append(
                {
                    "pair_id": pair_id,
                    "source_fragment_id": payload.get("source_id"),
                    "target_fragment_id": payload.get("target_id"),
                    "similarity": payload.get("similarity"),
                    "tier": payload.get("tier"),
                }
            )
        checks.append({"name": "candidate_threshold_applied", "status": "pass", "details": {"candidate_count": len(outputs["pair_ids"])}})
        why_summary = "Scored candidate fragment pairs using embedding similarity thresholds."
        why_policy_params = {"candidate_count": len(outputs["pair_ids"])}

    elif legacy_step_name == "map.reconstructable.normalize":
        normalized = step_outputs.get("normalized_fragments", [])
        outputs["fragment_ids"] = [
            str(getattr(fragment, "id", None) or getattr(fragment, "cas_id", None))
            for fragment in normalized
            if getattr(fragment, "id", None) is not None or getattr(fragment, "cas_id", None) is not None
        ]
        programs = step_outputs.get("reconstruction_programs", [])
        outputs["program_ids"] = [
            str(getattr(program, "id", None) or getattr(program, "cas_id", None))
            for program in programs
            if getattr(program, "id", None) is not None or getattr(program, "cas_id", None) is not None
        ]
        nmap = step_outputs.get("normalized_from_map", {})
        outputs["normalized_from_map"] = [
            {"normalized_fragment_id": str(key), "source_ir_fragment_ids": [str(item) for item in value]}
            for key, value in (dict(nmap).items() if isinstance(nmap, dict) else [])
        ]
        checks.append({"name": "normalization_lineage_complete", "status": "pass", "details": {"normalized_count": len(outputs["fragment_ids"])}})
        why_summary = "Normalized candidate fragments and produced reconstruction programs."
        why_policy_params = {"normalized_count": len(outputs["fragment_ids"]), "program_count": len(outputs["program_ids"])}

    elif legacy_step_name == "map.reconstructable.compose.reconstruction_programs":
        proposal = step_outputs.get("proposal", {})
        if isinstance(proposal, dict):
            outputs["proposal"] = {
                "commit_mode": proposal.get("commit_mode"),
                "proposed_fragment_ids": proposal.get("normalized_fragment_ids", []),
                "proposed_program_ids": proposal.get("program_ids", []),
                "fragment_count": proposal.get("fragment_count"),
            }
        else:
            outputs["proposal"] = {"commit_mode": getattr(proposal, "commit_mode", None), "proposed_fragment_ids": [], "proposed_program_ids": []}
        checks.append({"name": "proposal_complete_for_mode", "status": "pass", "details": {"commit_mode": outputs["proposal"].get("commit_mode")}})
        why_summary = "Composed commit proposal from normalized outputs and reconstruction programs."
        why_policy_params = {"commit_mode": outputs["proposal"].get("commit_mode")}

    elif legacy_step_name == "map.conceptual.verify.discovery_gate":
        verification = step_outputs.get("verification", {})
        if isinstance(verification, dict):
            outputs["verification"] = {
                "passed": verification.get("passed"),
                "metric": verification.get("policy") or verification.get("metric"),
                "tolerance": verification.get("tolerance"),
                "measured_drift": verification.get("measured_drift"),
                "diff_summary": verification.get("reason") or verification.get("diff_summary"),
                "reconstruction_chain_ids": verification.get("reconstruction_chain_ids", []),
            }
        else:
            outputs["verification"] = {"passed": getattr(verification, "passed", None)}
        if event.error and event.error.get("drift_at"):
            outputs["retry_target_step_name"] = event.error.get("drift_at")
        check_status = "pass" if bool(outputs["verification"].get("passed")) else "fail"
        checks.append({"name": "proof_gate", "status": check_status, "details": {"verification": outputs["verification"]}})
        why_summary = "Verified reconstructed output against source artifact constraints."
        why_policy_params = {"passed": outputs["verification"].get("passed")}

    elif legacy_step_name == "map.conceptual.commit.semantic_only":
        commit = step_outputs.get("commit", {})
        if isinstance(commit, dict):
            outputs["commit"] = {
                "mode": commit.get("mode"),
                "target_ref": commit.get("target_ref") or "refs/heads/main",
                "promoted_fragment_ids": commit.get("promoted_fragment_ids") or commit.get("committed_fragment_ids", []),
                "promoted_program_ids": commit.get("promoted_program_ids") or commit.get("committed_program_ids", []),
            }
        else:
            outputs["commit"] = {
                "mode": getattr(commit, "mode", None),
                "target_ref": "refs/heads/main",
                "promoted_fragment_ids": [],
                "promoted_program_ids": [],
            }
        checks.append({"name": "promotion_scope_valid", "status": "pass", "details": {"mode": outputs["commit"].get("mode")}})
        why_summary = "Promoted verified proposal into target ref scope with promoted fragments and programs."
        why_policy_params = {"mode": outputs["commit"].get("mode")}

    elif legacy_step_name == "map.reconstructable.build_subgraph.reconstruction":
        graph = step_outputs.get("graph_projection", {})
        if isinstance(graph, dict):
            nodes = graph.get("nodes", [])
            edges = graph.get("edges", [])
            outputs["graph_delta"] = {
                "nodes_added": len(nodes),
                "edges_added": len(edges),
                "node_ids": [str(node.get("id") or node.get("fragment_id") or f"node-{index}") for index, node in enumerate(nodes) if isinstance(node, dict)],
                "edge_ids": [str(edge.get("id") or f"edge-{index}") for index, edge in enumerate(edges) if isinstance(edge, dict)],
            }
            outputs["graph_summary"] = {
                "total_nodes": len(nodes),
                "total_edges": len(edges),
            }
        else:
            outputs["graph_delta"] = {"nodes_added": 0, "edges_added": 0, "node_ids": [], "edge_ids": []}
            outputs["graph_summary"] = {"total_nodes": 0, "total_edges": 0}
        checks.append({"name": "projection_predicates_present", "status": "pass", "details": outputs["graph_summary"]})
        why_summary = "Projected promoted fragments and relations into graph snapshot deltas for the target ref."
        why_policy_params = outputs["graph_summary"]

    # enrich lineage using explicit mappings for selected steps
    if legacy_step_name == "map.conceptual.lift.surface_fragments" and artifact_file_names:
        for node in lineage_nodes:
            if not isinstance(node, dict) or not str(node.get("node_id", "")).startswith("fragment:"):
                continue
            meta = node.get("meta") if isinstance(node.get("meta"), dict) else {}
            artifact_ref = meta.get("artifact_id") if isinstance(meta.get("artifact_id"), str) else None
            if artifact_ref and artifact_ref in artifact_file_names:
                meta["file_name"] = artifact_file_names[artifact_ref]
                node["meta"] = meta

    if raw_step_name == "map.conceptual.normalize.discovery":
        for item in outputs.get("lift_map", []):
            surface_id = item.get("surface_fragment_id")
            for ir_id in item.get("ir_fragment_ids", []):
                lineage_edges.append(
                    {
                        "from": f"fragment:{surface_id}",
                        "to": f"fragment:{ir_id}",
                        "relation": "lifted_from",
                        "step_name": step_name,
                    }
                )
    if raw_step_name == "map.reconstructable.normalize":
        for item in outputs.get("normalized_from_map", []):
            normalized_id = item.get("normalized_fragment_id")
            for source_id in item.get("source_ir_fragment_ids", []):
                lineage_edges.append(
                    {
                        "from": f"fragment:{source_id}",
                        "to": f"fragment:{normalized_id}",
                        "relation": "normalized_from",
                        "step_name": step_name,
                    }
                )

    step_trace = _serialize_step_trace(event)

    def _normalized_boundary_ids(raw: dict[str, Any] | None) -> dict[str, list[str]]:
        raw = raw if isinstance(raw, dict) else {}
        return {
            "artifact_ids": [str(item) for item in raw.get("artifact_ids", []) if isinstance(item, str) and item],
            "fragment_ids": [str(item) for item in raw.get("fragment_ids", []) if isinstance(item, str) and item],
            "program_ids": [str(item) for item in raw.get("program_ids", []) if isinstance(item, str) and item],
            "verification_ids": [str(item) for item in raw.get("verification_ids", []) if isinstance(item, str) and item],
        }

    def _boundary_summary(ids: dict[str, list[str]]) -> dict[str, int]:
        return {
            "artifact_count": len(ids["artifact_ids"]),
            "fragment_count": len(ids["fragment_ids"]),
            "program_count": len(ids["program_ids"]),
            "verification_count": len(ids["verification_ids"]),
        }

    input_ids = _normalized_boundary_ids(inputs)
    output_ids = _normalized_boundary_ids(outputs)
    if produced_fragment_ids:
        output_ids["fragment_ids"] = list(dict.fromkeys([*output_ids["fragment_ids"], *produced_fragment_ids]))

    scoped_fragment_ids = sorted(
        {
            str(payload.get("id") or payload.get("cas_id") or "")
            for payload in environment_fragments
            if isinstance(payload, dict) and str(payload.get("id") or payload.get("cas_id") or "")
        }
    )
    scoped_fragment_ids = [fragment_id for fragment_id in scoped_fragment_ids if fragment_id]

    trace_fragment_id = ""
    marking_before_ref = ""
    marking_after_ref = ""
    executor_id = ""
    executor_kind = ""
    if isinstance(step_trace, dict):
        trace_fragment_id = str(step_trace.get("trace_fragment_id") or "").strip()
        marking_before_ref = str(step_trace.get("marking_before_ref") or "").strip()
        marking_after_ref = str(step_trace.get("marking_after_ref") or "").strip()
        executor_id = str(step_trace.get("executor_id") or "").strip()
        executor_kind = str(step_trace.get("executor_kind") or "").strip()

    snapshots = step_output_snapshots if isinstance(step_output_snapshots, dict) else {}
    ordered_events = list(events or [])
    next_event = None
    for index, candidate in enumerate(ordered_events):
        if candidate.step_id != event.step_id:
            continue
        for later in ordered_events[index + 1 :]:
            if later.run_id == event.run_id:
                next_event = later
                break
        break

    next_snapshot = snapshots.get(next_event.step_id) if next_event and isinstance(snapshots.get(next_event.step_id), dict) else {}
    next_inputs = next_snapshot.get("inputs") if isinstance(next_snapshot, dict) and isinstance(next_snapshot.get("inputs"), dict) else None
    if isinstance(next_inputs, dict):
        handoff_ids = _normalized_boundary_ids(next_inputs)
        handoff_source = "next_step_inputs"
    else:
        handoff_ids = {
            "artifact_ids": list(output_ids["artifact_ids"]),
            "fragment_ids": list(output_ids["fragment_ids"]),
            "program_ids": list(output_ids["program_ids"]),
            "verification_ids": list(output_ids["verification_ids"]),
        }
        handoff_source = "declared_outputs_fallback"

    persisted_only_keys = []
    if output_ids["fragment_ids"] and set(output_ids["fragment_ids"]) - set(handoff_ids["fragment_ids"]):
        persisted_only_keys.append("fragment_ids")
    if output_ids["program_ids"] and set(output_ids["program_ids"]) - set(handoff_ids["program_ids"]):
        persisted_only_keys.append("program_ids")
    if output_ids["verification_ids"] and set(output_ids["verification_ids"]) - set(handoff_ids["verification_ids"]):
        persisted_only_keys.append("verification_ids")

    payload: dict[str, Any] = {
        "schema_version": "v1",
        "pipeline_id": state.pipeline_id,
        "pipeline_run_id": state.pipeline_run_id,
        "run_id": event.run_id,
        "step_id": event.step_id,
        "step_name": step_name,
        "attempt_index": event.attempt_index,
        "outcome": {
            "status": event.status,
            "duration_ms": event.duration_ms,
            "ref": scope_ref,
        },
        "why": {
            "summary": why_summary,
            "policy": {
                "name": why_policy_name,
                "params": why_policy_params,
            },
        },
        "inputs": inputs,
        "outputs": outputs,
        "checks": checks,
        "operation_ref": operation_ref,
        "operation_params": operation_params,
        "produced_fragment_ids": produced_fragment_ids,
        "lineage": {
            "roots": lineage_root_ids,
            "nodes": lineage_nodes,
            "edges": lineage_edges,
        },
        "step_boundaries": {
            "input_boundary": {
                **input_ids,
                "payload_summary": _boundary_summary(input_ids),
                "raw_payload": inputs,
            },
            "transition": {
                "step_name": step_name,
                "attempt_index": event.attempt_index,
                "operator_ref": operation_ref,
                "operator_params": operation_params,
                "executor_id": executor_id,
                "executor_kind": executor_kind,
                "duration_ms": event.duration_ms,
                "policy": {
                    "name": why_policy_name,
                    "params": why_policy_params,
                },
                "checks": checks,
                "why_summary": why_summary,
            },
            "output_boundary": {
                **output_ids,
                "produced_fragment_ids": produced_fragment_ids,
                "promoted_fragment_ids": [str(item) for item in outputs.get("promoted_fragment_ids", []) if isinstance(item, str) and item],
                "target_ref": outputs.get("target_ref") if isinstance(outputs.get("target_ref"), str) else None,
                "executor_outputs": outputs,
                "ikam_writes": {
                    "fragment_ids": list(output_ids["fragment_ids"]),
                    "program_ids": list(output_ids["program_ids"]),
                    "verification_ids": list(output_ids["verification_ids"]),
                },
                "payload_summary": _boundary_summary(output_ids),
                "raw_payload": outputs,
            },
            "ikam_environment_before": {
                "active_ref": scope_ref,
                "marking_ref": marking_before_ref,
                "trace_fragment_id": trace_fragment_id,
                "fragment_count": len(scoped_fragment_ids),
                "verification_count": len(output_ids["verification_ids"]),
                "reconstruction_program_count": len(output_ids["program_ids"]),
                "visible_artifact_ids": list(input_ids["artifact_ids"]),
                "visible_fragment_ids": list(dict.fromkeys([*input_ids["fragment_ids"], *scoped_fragment_ids])),
                "visible_program_ids": list(input_ids["program_ids"]),
                "snapshot_mode": "estimated_from_step_outputs",
            },
            "ikam_environment_after": {
                "active_ref": scope_ref,
                "marking_ref": marking_after_ref,
                "trace_fragment_id": trace_fragment_id,
                "fragment_count": len(scoped_fragment_ids),
                "verification_count": len(output_ids["verification_ids"]),
                "reconstruction_program_count": len(output_ids["program_ids"]),
                "visible_artifact_ids": list(input_ids["artifact_ids"]),
                "visible_fragment_ids": list(dict.fromkeys([*scoped_fragment_ids, *output_ids["fragment_ids"]])),
                "visible_program_ids": list(output_ids["program_ids"]),
                "snapshot_mode": "post_step_snapshot",
            },
            "handoff_to_next": {
                "next_step_name": next_event.step_name if next_event else None,
                "source": handoff_source,
                "forwarded_artifact_ids": handoff_ids["artifact_ids"],
                "forwarded_fragment_ids": handoff_ids["fragment_ids"],
                "forwarded_program_ids": handoff_ids["program_ids"],
                "forwarded_verification_ids": handoff_ids["verification_ids"],
                "dropped_executor_local_keys": [],
                "persisted_only_keys": persisted_only_keys,
                "notes": "Derived from recorded inputs of next step" if handoff_source == "next_step_inputs" else "Fallback to declared outputs",
            },
        },
    }
    persisted_transition_validations = runtime_context.get("step_transition_validations")
    persisted_input_validations = runtime_context.get("step_input_validations")
    persisted_output_validations = runtime_context.get("step_output_validations")
    input_validation = (
        persisted_input_validations.get(event.step_id)
        if isinstance(persisted_input_validations, dict) and isinstance(persisted_input_validations.get(event.step_id), dict)
        else None
    )
    output_validation = (
        persisted_output_validations.get(event.step_id)
        if isinstance(persisted_output_validations, dict) and isinstance(persisted_output_validations.get(event.step_id), dict)
        else None
    )
    transition_validation = (
        persisted_transition_validations.get(event.step_id)
        if isinstance(persisted_transition_validations, dict) and isinstance(persisted_transition_validations.get(event.step_id), dict)
        else None
    )
    if input_validation is None or _validation_has_unresolved_refs(input_validation, direction="input"):
        input_validation = build_runtime_transition_validation_for_direction(
            validators=_workflow_validators(state.pipeline_id, raw_step_name),
            artifact_id=str(runtime_context.get("artifact_id") or step_outputs.get("artifact_id") or ""),
            mime_type=str(runtime_context.get("mime_type") or ""),
            fixture_path=(
                str(runtime_context.get("fixture_path") or "")
                or (
                    str(case_fixture_dir(run_record.case_id).resolve())
                    if run_record is not None and isinstance(run_record.case_id, str) and run_record.case_id
                    else None
                )
            ),
            run_id=event.run_id,
            step_id=event.step_id,
            step_outputs=step_outputs,
            environment_fragments=environment_fragments,
            direction="input",
        )
    if output_validation is None or _validation_has_unresolved_refs(output_validation, direction="output"):
        output_validation = build_runtime_transition_validation_for_direction(
            validators=_workflow_validators(state.pipeline_id, raw_step_name),
            artifact_id=str(runtime_context.get("artifact_id") or step_outputs.get("artifact_id") or ""),
            mime_type=str(runtime_context.get("mime_type") or ""),
            fixture_path=(
                str(runtime_context.get("fixture_path") or "")
                or (
                    str(case_fixture_dir(run_record.case_id).resolve())
                    if run_record is not None and isinstance(run_record.case_id, str) and run_record.case_id
                    else None
                )
            ),
            run_id=event.run_id,
            step_id=event.step_id,
            step_outputs=_normalize_step_outputs_for_output_validation(
                step_outputs=step_outputs,
                pre_execution_outputs=step_outputs,
            ),
            environment_fragments=environment_fragments,
            direction="output",
        )
    if transition_validation is None or _validation_has_unresolved_refs(transition_validation, direction="output"):
        transition_validation = build_runtime_transition_validation(
            validators=_workflow_validators(state.pipeline_id, raw_step_name),
            artifact_id=str(runtime_context.get("artifact_id") or step_outputs.get("artifact_id") or ""),
            mime_type=str(runtime_context.get("mime_type") or ""),
            fixture_path=(
                str(runtime_context.get("fixture_path") or "")
                or (
                    str(case_fixture_dir(run_record.case_id).resolve())
                    if run_record is not None and isinstance(run_record.case_id, str) and run_record.case_id
                    else None
                )
            ),
            run_id=event.run_id,
            step_id=event.step_id,
            step_outputs=step_outputs,
            environment_fragments=environment_fragments,
        )
    fallback_input_validation, fallback_output_validation = _split_transition_validation_by_direction(transition_validation)
    if input_validation is None or _validation_is_effectively_empty(input_validation, direction="input"):
        input_validation = fallback_input_validation
    if output_validation is None or _validation_is_effectively_empty(output_validation, direction="output"):
        output_validation = fallback_output_validation
    if transition_validation is not None:
        payload["transition_validation"] = transition_validation
    if input_validation is not None:
        payload["input_validation"] = input_validation
    if output_validation is not None:
        payload["output_validation"] = output_validation
    next_transition_id = event.metrics.get("next_transition_id") if isinstance(event.metrics, dict) else None
    if isinstance(next_transition_id, str) and next_transition_id:
        payload["next_transition_id"] = next_transition_id

    if run_record is not None:
        payload["graph_id"] = run_record.graph.graph_id

    if step_trace is not None:
        payload["trace"] = step_trace

    log_events = _serialize_log_events(event)
    if log_events is not None:
        payload["log_events"] = log_events

    executor_logs = (
        _named_logs_from_log_events(log_events, source="executor")
        if log_events is not None
        else (_extract_runtime_trace_logs(event) or _serialize_named_step_logs(event, "executor_logs"))
    )
    if executor_logs is not None:
        payload["executor_logs"] = executor_logs

    system_logs = (
        _named_logs_from_log_events(log_events, source="system")
        if log_events is not None
        else _serialize_named_step_logs(event, "system_logs")
    )
    if system_logs is not None:
        payload["system_logs"] = system_logs

    step_logs = _legacy_logs_from_log_events(log_events) if log_events is not None else _serialize_step_logs(event)
    if step_logs is not None:
        payload["logs"] = step_logs

    return payload


def _validate_canonical_step_detail(payload: dict[str, Any]) -> None:
    required_top = {
        "schema_version",
        "pipeline_id",
        "pipeline_run_id",
        "run_id",
        "step_id",
        "step_name",
        "attempt_index",
        "outcome",
        "why",
        "inputs",
        "outputs",
        "checks",
        "lineage",
    }
    missing = [key for key in required_top if key not in payload]
    if missing:
        raise HTTPException(status_code=422, detail={"status": "schema_invalid", "missing": missing})

    lineage = payload.get("lineage")
    if not isinstance(lineage, dict):
        raise HTTPException(status_code=422, detail={"status": "schema_invalid", "missing": ["lineage"]})
    for key in ("roots", "nodes", "edges"):
        if key not in lineage:
            raise HTTPException(status_code=422, detail={"status": "schema_invalid", "missing": [f"lineage.{key}"]})


@router.get("/runs/{run_id}/env/fragments")
def get_env_fragments(
    run_id: str,
    pipeline_id: str,
    pipeline_run_id: str,
    step_id: str,
    attempt_index: int,
    ref: str | None = None,
    env_type: str | None = None,
    env_id: str | None = None,
):
    scope = _resolve_scope_query(ref=ref, env_type=env_type, env_id=env_id)
    state = STORE.get_debug_run_state(run_id)
    if not state:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "fragments": [],
        }
    if state.pipeline_id != pipeline_id or state.pipeline_run_id != pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "fragments": [],
        }
    fragments = STORE.list_scoped_fragments(
        run_id=run_id,
        ref=scope.ref,
        step_id=step_id,
        attempt_index=attempt_index,
    )
    return {
        "run_id": run_id,
        "pipeline_id": pipeline_id,
        "pipeline_run_id": pipeline_run_id,
        "status": "ok",
        "scope": {
            "ref": scope.ref,
            "step_id": step_id,
            "attempt_index": attempt_index,
        },
        "fragments": fragments,
    }


@router.get("/runs/{run_id}/verification")
def get_verification_records(
    run_id: str,
    pipeline_id: str,
    pipeline_run_id: str,
    step_id: str,
    attempt_index: int,
    ref: str | None = None,
    env_type: str | None = None,
    env_id: str | None = None,
):
    scope = _resolve_scope_query(ref=ref, env_type=env_type, env_id=env_id)
    state = STORE.get_debug_run_state(run_id)
    if not state or state.pipeline_id != pipeline_id or state.pipeline_run_id != pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "verification_records": [],
        }
    records = STORE.list_scoped_verification_records(
        run_id=run_id,
        ref=scope.ref,
        step_id=step_id,
        attempt_index=attempt_index,
    )
    return {
        "run_id": run_id,
        "pipeline_id": pipeline_id,
        "pipeline_run_id": pipeline_run_id,
        "status": "ok",
        "verification_records": records,
    }


@router.get("/runs/{run_id}/reconstruction-program")
def get_reconstruction_programs(
    run_id: str,
    pipeline_id: str,
    pipeline_run_id: str,
    step_id: str,
    attempt_index: int,
    ref: str | None = None,
    env_type: str | None = None,
    env_id: str | None = None,
):
    scope = _resolve_scope_query(ref=ref, env_type=env_type, env_id=env_id)
    state = STORE.get_debug_run_state(run_id)
    if not state or state.pipeline_id != pipeline_id or state.pipeline_run_id != pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "reconstruction_programs": [],
        }
    records = STORE.list_scoped_reconstruction_programs(
        run_id=run_id,
        ref=scope.ref,
        step_id=step_id,
        attempt_index=attempt_index,
    )
    return {
        "run_id": run_id,
        "pipeline_id": pipeline_id,
        "pipeline_run_id": pipeline_run_id,
        "status": "ok",
        "reconstruction_programs": records,
    }


@router.get("/runs/{run_id}/env-summary")
def get_environment_summary(
    run_id: str,
    pipeline_id: str,
    pipeline_run_id: str,
    ref: str | None = None,
    env_type: str | None = None,
    env_id: str | None = None,
):
    scope = _resolve_scope_query(ref=ref, env_type=env_type, env_id=env_id)
    state = STORE.get_debug_run_state(run_id)
    if not state or state.pipeline_id != pipeline_id or state.pipeline_run_id != pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": pipeline_id,
            "pipeline_run_id": pipeline_run_id,
            "status": "missing",
            "summary": {
                "fragment_count": 0,
                "verification_count": 0,
                "reconstruction_program_count": 0,
            },
        }
    summary = STORE.summarize_environment_scope(
        run_id=run_id,
        ref=scope.ref,
    )
    summary = {key: value for key, value in summary.items() if key not in {"env_type", "env_id"}}
    return {
        "run_id": run_id,
        "pipeline_id": pipeline_id,
        "pipeline_run_id": pipeline_run_id,
        "status": "ok",
        "summary": summary,
    }


@router.post("/runs/{run_id}/control")
async def control_run(run_id: str, payload: RunControlPayload):
    if payload.action == "inject_verify_fail":
        _require_test_gate()

    state = STORE.get_debug_run_state(run_id)
    if not state or state.pipeline_id != payload.pipeline_id or state.pipeline_run_id != payload.pipeline_run_id:
        return {
            "run_id": run_id,
            "pipeline_id": payload.pipeline_id,
            "pipeline_run_id": payload.pipeline_run_id,
            "status": "missing",
        }
        
    pipeline_steps = _debug_pipeline_steps(payload.pipeline_id)

    is_new = STORE.register_control_command(run_id, payload.command_id)
    if not is_new:
        return {
            "run_id": run_id,
            "pipeline_id": payload.pipeline_id,
            "pipeline_run_id": payload.pipeline_run_id,
            "status": "duplicate",
            "state": _serialize_debug_state(state),
        }

    next_state = state
    if payload.action == "set_mode" and payload.mode:
        next_state.execution_mode = payload.mode
    elif payload.action == "pause":
        next_state.execution_state = "paused"
    elif payload.action == "resume":
        next_state.execution_state = "running"
    elif payload.action == "next_step":
        if _async_next_step_enabled():
            existing_task = _RUN_STEP_TASKS.get(run_id)
            if existing_task is not None and not existing_task.done():
                availability = _control_availability(state=next_state, events=STORE.list_debug_events(run_id), pipeline_steps=pipeline_steps)
                return {
                    "run_id": run_id,
                    "pipeline_id": payload.pipeline_id,
                    "pipeline_run_id": payload.pipeline_run_id,
                    "status": "busy",
                    "control_availability": availability,
                    "state": _serialize_debug_state(next_state),
                }

            prepared_step, terminal_event = _prepare_next_pipeline_step(run_id=run_id, state=next_state)
            if terminal_event is not None:
                event = terminal_event
                if event.status != "succeeded" and next_state.execution_state != "budget_exhausted":
                    next_state.execution_state = "paused"
                elif event.status == "succeeded" and next_state.current_step_name == pipeline_steps[-1]:
                    next_state.execution_state = "completed"
                STORE.set_debug_run_state(run_id, next_state)
                availability = _control_availability(state=next_state, events=STORE.list_debug_events(run_id), pipeline_steps=pipeline_steps)
                return {
                    "run_id": run_id,
                    "pipeline_id": payload.pipeline_id,
                    "pipeline_run_id": payload.pipeline_run_id,
                    "status": "accepted",
                    "control_availability": availability,
                    "state": _serialize_debug_state(next_state),
                }

            next_state.execution_state = "running"
            STORE.set_debug_run_state(run_id, next_state)
            _RUN_STEP_TASKS[run_id] = asyncio.create_task(
                _execute_next_pipeline_step_async(run_id=run_id, state=next_state, prepared_step=prepared_step)
            )

            availability = _control_availability(state=next_state, events=STORE.list_debug_events(run_id), pipeline_steps=pipeline_steps)
            return {
                "run_id": run_id,
                "pipeline_id": payload.pipeline_id,
                "pipeline_run_id": payload.pipeline_run_id,
                "status": "accepted",
                "control_availability": availability,
                "state": _serialize_debug_state(next_state),
            }

        event = await _execute_next_pipeline_step(run_id=run_id, state=next_state)
        if event.status != "succeeded" and next_state.execution_state != "budget_exhausted":
            next_state.execution_state = "paused"
        elif event.status == "succeeded" and next_state.current_step_name == pipeline_steps[-1]:
            next_state.execution_state = "completed"
    elif payload.action == "inject_verify_fail":
        runtime_context = STORE.get_debug_runtime_context(run_id)
        if runtime_context is None:
            raise RuntimeError("Missing debug runtime context for injection")
        step_outputs = runtime_context.get("step_outputs") or {}
        step_outputs["_injected_verify_fail"] = {
            "drift_at": payload.drift_at or "map.conceptual.lift.surface_fragments",
        }
        runtime_context["step_outputs"] = step_outputs
        STORE.set_debug_runtime_context(run_id, runtime_context)

    resume_events: list[DebugStepEvent] = []
    if payload.action == "resume":
        while next_state.execution_state == "running" and next_state.current_step_name != pipeline_steps[-1]:
            event = await _execute_next_pipeline_step(run_id=run_id, state=next_state)
            resume_events.append(event)
            if event.status != "succeeded":
                next_state.execution_state = "paused"
                break
        if next_state.execution_state == "running" and next_state.current_step_name == pipeline_steps[-1]:
            next_state.execution_state = "completed"

    STORE.set_debug_run_state(run_id, next_state)
    availability = _control_availability(state=next_state, events=STORE.list_debug_events(run_id), pipeline_steps=pipeline_steps)
    result: dict = {
        "run_id": run_id,
        "pipeline_id": payload.pipeline_id,
        "pipeline_run_id": payload.pipeline_run_id,
        "status": "ok",
        "control_availability": availability,
        "state": _serialize_debug_state(next_state),
    }
    if payload.action == "resume":
        result["events"] = [_serialize_debug_event(e) for e in resume_events]
    return result


@router.get("/runs/{run_id}/debug-state")
def get_debug_state(run_id: str):
    """Return current DebugRunState for a run (no pipeline identity required)."""
    state = STORE.get_debug_run_state(run_id)
    if not state:
        return {"run_id": run_id, "status": "missing"}
    events = STORE.list_debug_events(run_id)
    pipeline_steps = _debug_pipeline_steps(state.pipeline_id)
    availability = _control_availability(state=state, events=events, pipeline_steps=pipeline_steps)
    return {
        "run_id": run_id,
        "status": "ok",
        "pipeline_id": state.pipeline_id,
        "pipeline_run_id": state.pipeline_run_id,
        "ref": _legacy_scope_ref(state.env_type, state.env_id),
        "execution_state": state.execution_state,
        "execution_mode": state.execution_mode,
        "current_step_name": state.current_step_name,
        "pipeline_steps": pipeline_steps,
        "current_attempt_index": state.current_attempt_index,
        "control_availability": availability,
    }


@router.get("/runs/{run_id}/embedding-info")
def get_embedding_info(run_id: str):
    """Query pgvector DB for embedding metadata about a completed run."""
    if STORE.get_run(run_id) is None:
        raise HTTPException(status_code=404, detail={"run_id": run_id, "status": "not_found"})

    database_url = os.getenv("DATABASE_URL")
    if not database_url:
        return {"run_id": run_id, "status": "db_unavailable"}

    runtime_context = STORE.get_debug_runtime_context(run_id)
    artifact_id = (runtime_context or {}).get("artifact_id", "")
    op_id = artifact_id.split(":")[-1] if ":" in artifact_id else artifact_id

    if not op_id:
        return {"run_id": run_id, "status": "no_runtime_context"}

    import psycopg

    try:
        with psycopg.connect(database_url) as conn:
            with conn.cursor() as cur:
                cur.execute(
                    """
                    SELECT embedding_model, vector_dims(embedding) AS vector_dims, COUNT(*) AS fragment_count
                    FROM ikam_fragment_store
                    WHERE operation_id = %s
                    GROUP BY embedding_model, vector_dims(embedding)
                    """,
                    (op_id,),
                )
                rows = cur.fetchall()
    except Exception as db_exc:
        return {"run_id": run_id, "status": "db_error", "detail": str(db_exc)}

    if len(rows) == 0:
        return {"run_id": run_id, "embedding_model": None, "vector_dims": None, "fragment_count": 0}

    if len(rows) > 1:
        raise RuntimeError(
            f"embedding-info: expected 1 embedding group, got {len(rows)} for op_id={op_id!r}"
        )

    embedding_model, vector_dims, fragment_count = rows[0]
    return {
        "run_id": run_id,
        "status": "ok",
        "embedding_model": embedding_model,
        "vector_dims": vector_dims,
        "fragment_count": fragment_count,
    }
