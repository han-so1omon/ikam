from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Sunbeam Market Collective — Mission, Vision, Values

## Mission
Bring good food to neighborhoods with consistent quality, fair sourcing, and a calm operational backbone.

## Vision
Be the Bay Area’s most trusted specialty grocer for prepared foods and community-first retail.

## Values
- Freshness is a promise
- Clear definitions (numbers and recipes)
- Respect the floor
- Fix the system, not the person
""",
    )

    w(
        "brand-guide.md",
        """# Sunbeam Market Collective — Brand Guide (v0.1)

## Palette
- Sun Yellow: #FBBF24
- Deep Green: #1F3D2B
- Off-white: #F7F5F0
- Slate: #64748B

## Voice
- Bright, helpful, specific.
- Avoid hype. Prefer “what it is” + “why it’s good”.

## Known drift
- Meal-kit pilot name varies across teams.
- Promo naming is inconsistent.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Camille Torres (COO)
Last updated: 2026-02-02

## Priorities
1) Reduce shrink/waste by aligning measurement and incentives.
2) Standardize prepared foods SOP (v3 rollout).
3) Get one revenue definition used by Store Ops and Finance.
4) Decide whether meal-kit pilot scales or shuts down.
""",
    )

    w(
        "campaign-brief-spring.md",
        """# Campaign Brief — Spring Refresh

Owner: Imani Brooks

## Objective
Drive prepared foods attachment rate and bring lapsed customers back.

## Key messages
- Seasonal bowls
- “Lunch solved”

## Risks
- Promo calendar drift between stores
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Fresh, daily, real ingredients
2) Neighborhood convenience
3) Prepared foods that taste homemade
4) Transparent sourcing
""",
    )

    w(
        "product-style-guide.md",
        """# Product Style Guide

## Prepared foods naming
- Name + key ingredient + dietary tag (if applicable)

## Labeling
- Ingredient list required
- Allergen callouts

## Photography
- Natural light; food-forward; avoid clutter
""",
    )

    w(
        "sop-prepared-foods-v2.md",
        """# SOP — Prepared Foods (v2)

Effective: 2024-12-01

- Batch sheets printed daily
- Temp logs on paper
- Waste recorded in store-level sheet
""",
    )

    w(
        "sop-prepared-foods-v3.md",
        """# SOP — Prepared Foods (v3)

Drafted: 2026-01-18

Changes:
- Waste recorded in central system (weekly)
- Temp logs moved to shared form

Known issue:
- Stores still using v2 printouts
""",
    )

    w(
        "store-manager-call-2026-01-08.md",
        """# Store Manager Call — 2026-01-08

- Promo calendar mismatch noted (Store #2 running old promo)
- Meal kit pilot name confusion continues
""",
    )

    w(
        "store-manager-call-2026-02-05.md",
        """# Store Manager Call — 2026-02-05

- CFO requests consistent revenue definition
- Commissary claims waste down 20%; store managers disagree
""",
    )

    w(
        "voice-note-transcript-2026-02-12.md",
        """# Voice Note Transcript — 2026-02-12 (Jonah)

- Store ops keeps reporting gross sales; finance reports net of returns/voids.
- Waste metrics don’t match because the commissary counts “rework” as saved.
- HR headcount depends on whether we include temps.
""",
    )

    w(
        "vendor-scorecard-template.md",
        """# Vendor Scorecard Template

Vendor:
Category:

Scores (1–5):
- Quality
- Fill rate
- Lead time
- Communication

Notes:
""",
    )

    w(
        "promo-postmortem-template.md",
        """# Promo Postmortem Template

Promo:
Dates:

## What happened

## Results

## Issues

## Fixes
""",
    )

    w(
        "performance-review-template.md",
        """# Performance Review Template

Role:
Period:

## Outcomes

## Strengths

## Growth areas

## Next goals
""",
    )


def make_json():
    wjson(
        "org-structure.json",
        {
            "generated": str(date.today()),
            "roles": {
                "COO": "Camille Torres",
                "CFO": "Jonah Feld",
                "Marketing": "Imani Brooks",
                "SupplyChain": "Wes Park",
                "HR": "Hana Kim",
                "Stores": ["Priya", "Marco", "Elena", "Devon", "Sam", "Tori"],
            },
        },
    )

    wjson(
        "vendor-list.json",
        {
            "generated": str(date.today()),
            "vendors": [
                {"id": "V001", "name": "BayHarvest Produce", "category": "produce", "risk": "seasonal"},
                {"id": "V002", "name": "GoldenCrust Bakery", "category": "bakery", "risk": "capacity"},
                {"id": "V003", "name": "NorthCoast Dairy", "category": "dairy", "risk": "low"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "Revenue (ops)", "definition": "gross sales", "notes": "includes voids until reconciled"},
                {"name": "Revenue (finance)", "definition": "net sales", "notes": "excludes returns/voids"},
                {"name": "Waste (commissary)", "definition": "waste minus rework savings", "notes": "optimistic"},
                {"name": "Shrink/Waste (finance)", "definition": "shrink + waste per inventory adjustment", "notes": "conservative"},
            ],
        },
    )

    wjson(
        "catering-order-intake-form.json",
        {
            "schema_version": "0.2",
            "title": "Catering Order Intake",
            "fields": [
                {"id": "company_name", "type": "string", "required": True},
                {"id": "contact_email", "type": "string", "required": True},
                {"id": "event_date", "type": "date", "required": True},
                {"id": "headcount", "type": "integer", "required": True},
                {"id": "dietary_notes", "type": "string"},
                {"id": "delivery", "type": "enum", "values": ["pickup", "delivery"], "required": True},
            ],
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "employee-handbook-excerpt.pdf",
        "Employee Handbook Excerpt (summary)",
        [
            "Food safety expectations",
            "Attendance policy",
            "Harassment-free workplace",
            "Temps policy varies by store (note: inconsistent)",
        ],
    )


def make_xlsx():
    import openpyxl
    from openpyxl.styles import Font, PatternFill

    # promo calendar
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="H1 2026"
    ws.append(["Week","Promo name","Stores","Notes"])
    rows=[("2026-W02","Winter Bowl Week","All",""),("2026-W06","Meal Kit Pilot Push","Store 1,2","name varies"),("2026-W10","Spring Refresh","All","")]
    for r in rows: ws.append(list(r))
    for c in ws[1]:
        c.fill=PatternFill("solid", fgColor="FBBF24"); c.font=Font(bold=True)
    wb.save(BASE/"promo-calendar-2026-h1.xlsx")

    # revenue history
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="Quarterly"
    ws2.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q3",8200000,""),(2024,"Q4",9600000,""),(2025,"Q1",9000000,""),(2025,"Q2",9800000,""),(2025,"Q3",10400000,""),(2025,"Q4",12100000,"holiday incident")]:
        ws2.append(list(r))
    for i in range(2, ws2.max_row+1): ws2.cell(i,3).number_format="$#,##0"
    wb2.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    # projection
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="2026"
    ws3.append(["Month","Revenue","Notes"])
    for r in [("2026-01",3200000,""),("2026-02",3000000,""),("2026-03",3400000,"")]: ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,2).number_format="$#,##0"
    wb3.save(BASE/"projected-revenue-2026.xlsx")

    # weekly store performance (ops gross vs finance net)
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-02"
    ws4.append(["Week","Ops Gross","Finance Net","Notes"])
    for r in [("2026-W05",780000,742000,"voids/returns"),("2026-W06",810000,768000,"promo mix")]: ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="$#,##0"; ws4.cell(i,3).number_format="$#,##0"
    wb4.save(BASE/"store-performance-weekly-2026-02.xlsx")

    # shrink/waste
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="2025"
    ws5.append(["Month","Commissary Waste %","Finance Shrink/Waste %","Notes"])
    for r in [("2025-10",0.031,0.035,""),("2025-11",0.028,0.036,"holiday prep"),("2025-12",0.025,0.035,"commissary claims -20%")]: ws5.append(list(r))
    for i in range(2, ws5.max_row+1):
        ws5.cell(i,2).number_format="0.0%"; ws5.cell(i,3).number_format="0.0%"
    wb5.save(BASE/"shrink-waste-summary-2025.xlsx")

    # replenishment plan
    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="Replenishment"
    ws6.append(["SKU","Item","Vendor","Reorder point","Target stock","Notes"])
    for r in [("PF-BOWL-01","Seasonal bowl","Commissary",120,300,""),("DAI-MILK-02","Whole milk","NorthCoast Dairy",80,200,"")]: ws6.append(list(r))
    wb6.save(BASE/"inventory-replenishment-plan.xlsx")

    # labor scheduling template
    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="Template"
    ws7.append(["Store","Role","Mon","Tue","Wed","Thu","Fri","Sat","Sun","Total hrs"])
    ws7.append(["Store 1","Prepared foods",8,8,8,8,8,10,10,None])
    ws7.cell(2,10).value="=SUM(C2:I2)"
    wb7.save(BASE/"labor-scheduling-template.xlsx")

    # hiring plan
    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="2026"
    ws8.append(["Dept","Role","Count","Priority","Notes"])
    for r in [("Commissary","Prep cook",6,"High",""),("Stores","Shift lead",4,"Med","temps counted differently")]: ws8.append(list(r))
    wb8.save(BASE/"hiring-plan-2026.xlsx")

    # incident log
    wb9=openpyxl.Workbook(); ws9=wb9.active; ws9.title="Q1 2026"
    ws9.append(["Date","Location","Incident","Severity","Notes"])
    for r in [("2026-01-12","Store 3","Prepared foods temp log missing","Med","SOP drift"),("2026-02-02","Commissary","Stockout on bowls","High","promo demand")]: ws9.append(list(r))
    wb9.save(BASE/"incident-log-2026-q1.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    Y=RGBColor(0xFB,0xBF,0x24)
    G=RGBColor(0x1F,0x3D,0x2B)
    OFF=RGBColor(0xF7,0xF5,0xF0)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=G; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=Y; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Sunbeam Market Collective"; p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
