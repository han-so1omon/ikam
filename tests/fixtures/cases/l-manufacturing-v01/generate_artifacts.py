from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Cascadia Fasteners & Forming — Mission, Vision, Values

## Mission
Make reliable components with short lead times by running disciplined manufacturing processes and honest customer communication.

## Vision
Be the first-call regional manufacturer for standard and custom fasteners/components across construction and light industrial supply chains.

## Values
- **Safety first.**
- **Ship what we promise.**
- **Document the work.**
- **Fix root causes.**
- **Respect the floor.**
""",
    )

    w(
        "brand-guide.md",
        """# Cascadia Fasteners & Forming — Brand Guide (v0.1)

## Voice + tone
- Direct, technical when needed, no fluff.
- Prefer specs, lead times, and clear commitments.

## Palette (suggested)
- Steel: #2F3A44
- White: #F7F7F7
- Safety orange: #F97316
- Slate: #64748B

## Slide style
- Steel header bar
- Safety orange accent line
- Simple diagrams over stock photos

## Naming
- Always include: SKU + plain-English name
- Sites: Tacoma Plant, Kent 3PL
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Dana Whitaker (GM)
Last updated: 2026-02-04

## Priorities
1) Improve on-time delivery (OTD) by tightening definitions and reducing partial shipments.
2) Reduce scrap + rework via press setup discipline and better measurement.
3) Inventory accuracy: reconcile ERP vs spreadsheet adjustments.

## Targets (directional)
- OTD (customer view): 95%+
- Internal OTD (strict): 93%+
- Scrap (incl. rework): -15% YoY

## Risks
- SOP revision drift (v3 vs v4 on the floor)
- Vendor lead time shocks (steel coil)
""",
    )

    w(
        "product-catalog-overview.md",
        """# Product Catalog Overview

## Standard SKUs
- BOLT-M8-30 — M8 bolt, 30mm (zinc plated)
- BRKT-L-02 — L-bracket (2")
- CLIP-HVAC-07 — HVAC mounting clip

## Custom runs
- Small formed brackets/clips per drawing
- Short lead times for repeat tooling

## Notes
- Pricing depends on steel index + plating vendor.
""",
    )

    w(
        "sop-press-line-setup-v3.md",
        """# SOP — Press Line Setup (v3)

Effective: 2025-03-10
Owner: Luis Ortega

## Steps
1) Lockout/tagout
2) Install die set
3) Set feeder tension
4) Run first-article parts (n=5)
5) QA checks dimensions; record on paper sheet

## Notes
- Scrap tracking: record scrap; rework tracked separately.
""",
    )

    w(
        "sop-press-line-setup-v4.md",
        """# SOP — Press Line Setup (v4)

Drafted: 2026-01-12
Owner: Anika Patel (QA)

Changes vs v3:
- Add measurement gauge calibration step
- First-article sample size n=8
- Record results in shared spreadsheet + upload photo

## Known issue
Floor still using v3 checklist printouts.
""",
    )

    w(
        "weekly-ops-review-2025-12-03.md",
        """# Weekly Ops Review — 2025-12-03

Attendees: Dana, Luis, Anika, Samir, Renee

## Metrics
- OTD: 92% (internal)
- Scrap: 3.8% (ops)

## Notes
- Partial shipments counted differently in customer scorecard.
- Inventory adjustments logged in a side spreadsheet again.
""",
    )

    w(
        "weekly-ops-review-2026-01-14.md",
        """# Weekly Ops Review — 2026-01-14

## Updates
- SOP v4 drafted; rollout plan unclear.
- Plating vendor quality improving.

## Actions
- Define OTD once (and store the definition next to the KPI).
- Inventory cycle counts weekly for top 20 SKUs.
""",
    )

    w(
        "voice-note-transcript-2026-01-22.md",
        """# Voice Note Transcript — 2026-01-22 (Renee)

- Finance inventory valuation is standard cost. Ops keeps adjusting using last purchase price on a few SKUs.
- Customer scorecard says 96% OTD because partial shipments count as on-time if first box arrives.
- Internal dashboard says 92% because we count full order completion date.
- We need the definitions in writing.
""",
    )

    w(
        "supplier-evaluation-template.md",
        """# Supplier Evaluation Template

Supplier:
Category:

Scores (1–5):
- Lead time reliability:
- Quality:
- Pricing stability:
- Responsiveness:

Notes:

Decision:
- Keep / Trial / Replace
""",
    )

    w(
        "nonconformance-report-template.md",
        """# Nonconformance Report (NCR) Template

NCR ID:
Date:
SKU / Part:
Issue description:
Containment action:
Root cause (5 whys):
Corrective action:
Owner:
Due date:
Verification:
""",
    )


def make_json():
    wjson(
        "vendor-list.json",
        {
            "generated": str(date.today()),
            "vendors": [
                {"vendor_id": "V-STEEL-01", "name": "Rainier Steel Supply", "category": "steel coil", "risk_notes": "Lead times volatile"},
                {"vendor_id": "V-PLATE-02", "name": "BlueHarbor Plating", "category": "plating", "risk_notes": "New vendor 2024-11; early quality issues"},
                {"vendor_id": "V-PACK-01", "name": "NorthPack", "category": "packaging", "risk_notes": "Stable"},
            ],
        },
    )

    wjson(
        "top-customers.json",
        {
            "generated": str(date.today()),
            "customers": [
                {"account": "Summit Build Supply", "segment": "Distributor", "region": "PNW"},
                {"account": "Evergreen HVAC Systems", "segment": "OEM", "region": "US"},
                {"account": "Harborline Contractors", "segment": "Construction", "region": "WA"},
                {"account": "Cascadia Industrial Parts", "segment": "Distributor", "region": "PNW"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-01-20",
            "kpis": [
                {"name": "OTD (customer)", "definition": "On-time if first shipment arrives by requested date", "notes": "May hide partial shipment delays"},
                {"name": "OTD (internal)", "definition": "On-time if full order completed by requested date", "notes": "Stricter"},
                {"name": "Scrap rate (ops)", "definition": "Scrap incl rework / total parts", "notes": "Includes rework as scrap"},
                {"name": "Scrap rate (QA)", "definition": "True scrap only / total parts", "notes": "Excludes rework"},
            ],
        },
    )

    wjson(
        "expedite-request-form.json",
        {
            "schema_version": "0.1",
            "title": "Expedite Request",
            "fields": [
                {"id": "customer", "type": "string", "required": True},
                {"id": "po_number", "type": "string", "required": True},
                {"id": "sku", "type": "string", "required": True},
                {"id": "quantity", "type": "integer", "required": True, "min": 1},
                {"id": "need_by", "type": "date", "required": True},
                {"id": "reason", "type": "string", "required": False},
            ],
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "customer-scorecard-2025-q4.pdf",
        "Customer Scorecard — Q4 2025 (Summary)",
        [
            "OTD (customer definition): 96%",
            "Quality: 2 customer NCRs (down from 5)",
            "Notes: Partial shipments counted as on-time if first shipment meets requested date.",
        ],
    )


def make_xlsx():
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill

    # Quarterly revenue
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Quarterly"
    ws.append(["Year", "Quarter", "Revenue", "Gross Margin %", "Notes"])
    rows = [
        (2024, "Q3", 2100000, 0.31, "Stable"),
        (2024, "Q4", 2450000, 0.30, "Holiday demand"),
        (2025, "Q1", 2250000, 0.32, ""),
        (2025, "Q2", 2380000, 0.33, "Press Line 4 ramp"),
        (2025, "Q3", 2500000, 0.32, "OTD definition dispute begins"),
        (2025, "Q4", 2700000, 0.31, "Inventory adjustments in side sheet"),
    ]
    for r in rows:
        ws.append(list(r))

    for c in ws[1]:
        c.fill = PatternFill("solid", fgColor="2F3A44")
        c.font = Font(color="FFFFFF", bold=True)
        c.alignment = Alignment(horizontal="center")

    for r in range(2, ws.max_row + 1):
        ws.cell(r, 3).number_format = "$#,##0"
        ws.cell(r, 4).number_format = "0.0%"

    wb.save(BASE / "quarterly-revenue-history-2024-2025.xlsx")

    # Projection 2026 (simple)
    wb2 = openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "2026"
    ws2.append(["Month", "Revenue", "Assumed GM%", "Gross Profit", "Notes"])
    months = [
        ("2026-01", 820000, 0.315, None, "SOP rollout attempt"),
        ("2026-02", 780000, 0.315, None, "Capex request for QA equipment"),
        ("2026-03", 860000, 0.315, None, "Steel index volatility"),
    ]
    for m, rev, gm, gp, notes in months:
        ws2.append([m, rev, gm, None, notes])

    for c in ws2[1]:
        c.fill = PatternFill("solid", fgColor="F97316")
        c.font = Font(color="FFFFFF", bold=True)

    for r in range(2, ws2.max_row + 1):
        rev = ws2.cell(r, 2).coordinate
        gm = ws2.cell(r, 3).coordinate
        ws2.cell(r, 4).value = f"={rev}*{gm}"
        ws2.cell(r, 2).number_format = "$#,##0"
        ws2.cell(r, 3).number_format = "0.0%"
        ws2.cell(r, 4).number_format = "$#,##0"

    wb2.save(BASE / "projected-revenue-2026.xlsx")

    # Inventory valuation
    wb3 = openpyxl.Workbook()
    ws3 = wb3.active
    ws3.title = "Inventory"
    ws3.append(["SKU", "Description", "Qty", "Std Cost", "Last Purchase", "Valuation (Finance)", "Valuation (Ops)", "Notes"])
    items = [
        ("BOLT-M8-30", "M8 bolt 30mm", 120000, 0.07, 0.08, None, None, "Ops uses last purchase for this SKU"),
        ("BRKT-L-02", "L bracket 2in", 48000, 0.22, 0.21, None, None, ""),
        ("CLIP-HVAC-07", "HVAC mounting clip", 65000, 0.11, 0.12, None, None, ""),
    ]
    for it in items:
        ws3.append(list(it))

    for r in range(2, ws3.max_row + 1):
        qty = ws3.cell(r, 3).coordinate
        std = ws3.cell(r, 4).coordinate
        last = ws3.cell(r, 5).coordinate
        ws3.cell(r, 6).value = f"={qty}*{std}"
        ws3.cell(r, 7).value = f"={qty}*{last}"
        for c in [4, 5]:
            ws3.cell(r, c).number_format = "$0.00"
        for c in [6, 7]:
            ws3.cell(r, c).number_format = "$#,##0"

    wb3.save(BASE / "inventory-valuation-2025-12.xlsx")

    # Capex plan
    wb4 = openpyxl.Workbook()
    ws4 = wb4.active
    ws4.title = "Capex 2026"
    ws4.append(["Item", "Dept", "Cost", "Justification", "Status"])
    capex = [
        ("Optical comparator", "QA", 45000, "Reduce measurement ambiguity; supports SOP v4", "Requested"),
        ("Die maintenance kit", "Ops", 12000, "Reduce setup time + scrap", "Approved"),
    ]
    for c in capex:
        ws4.append(list(c))
    for r in range(2, ws4.max_row + 1):
        ws4.cell(r, 3).number_format = "$#,##0"

    wb4.save(BASE / "capex-plan-2026.xlsx")

    # Production plan H1
    wb5 = openpyxl.Workbook()
    ws5 = wb5.active
    ws5.title = "H1 2026"
    ws5.append(["Week", "Press Line", "Primary SKU", "Planned Qty", "OEE Target", "Notes"])
    plan = [
        ("2026-W02", "Line 4", "BOLT-M8-30", 50000, 0.72, ""),
        ("2026-W03", "Line 2", "BRKT-L-02", 22000, 0.70, "Setup change"),
        ("2026-W04", "Line 4", "CLIP-HVAC-07", 30000, 0.71, ""),
    ]
    for p in plan:
        ws5.append(list(p))
    for r in range(2, ws5.max_row + 1):
        ws5.cell(r, 4).number_format = "#,##0"
        ws5.cell(r, 5).number_format = "0.0%"

    wb5.save(BASE / "production-plan-2026-h1.xlsx")

    # QA incidents log
    wb6 = openpyxl.Workbook()
    ws6 = wb6.active
    ws6.title = "Incidents 2025"
    ws6.append(["Date", "SKU", "Issue", "Severity", "Disposition", "Notes"])
    incidents = [
        ("2025-11-18", "BOLT-M8-30", "Plating thickness out of spec", "High", "Contain + replate", "Vendor BlueHarbor"),
        ("2025-12-09", "BRKT-L-02", "Burrs above tolerance", "Med", "Rework", "Setup drift"),
    ]
    for inc in incidents:
        ws6.append(list(inc))
    wb6.save(BASE / "qa-incidents-log-2025.xlsx")

    # Price list
    wb7 = openpyxl.Workbook()
    ws7 = wb7.active
    ws7.title = "Q1 2026"
    ws7.append(["SKU", "Description", "Unit Price", "Min Order Qty", "Notes"])
    prices = [
        ("BOLT-M8-30", "M8 bolt 30mm", 0.18, 5000, "Steel index clause"),
        ("BRKT-L-02", "L bracket 2in", 0.65, 2000, ""),
        ("CLIP-HVAC-07", "HVAC mounting clip", 0.34, 3000, ""),
    ]
    for pr in prices:
        ws7.append(list(pr))
    for r in range(2, ws7.max_row + 1):
        ws7.cell(r, 3).number_format = "$0.00"
        ws7.cell(r, 4).number_format = "#,##0"
    wb7.save(BASE / "price-list-2026-q1.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    STEEL = RGBColor(0x2F, 0x3A, 0x44)
    WHITE = RGBColor(0xF7, 0xF7, 0xF7)
    ORANGE = RGBColor(0xF9, 0x73, 0x16)

    prs = Presentation()

    def bg(slide):
        f = slide.background.fill
        f.solid(); f.fore_color.rgb = WHITE

    def header(slide):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = STEEL
        bar.line.fill.background()
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb = ORANGE
        acc.line.fill.background()

    def title_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg(s); header(s)
        tx = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        p.text = "Cascadia Fasteners & Forming"
        p.font.size = Pt(44); p.font.bold = True; p.font.name = "Georgia"; p.font.color.rgb = STEEL
        sub = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11), Inches(0.8))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = "Regional manufacturing • standard + custom components • short lead times"
        p2.font.size = Pt(20); p2.font.name = "Arial"; p2.font.color.rgb = STEEL

    def bullets(title, items):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg(s); header(s)
        t = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
        p = t.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(34); p.font.bold = True; p.font.name = "Georgia"; p.font.color.rgb = STEEL
        box = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(4.2))
        tf = box.text_frame
        tf.clear()
        for i, it in enumerate(items):
            pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            pp.text = it
            pp.font.name = "Arial"; pp.font.size = Pt(20); pp.font.color.rgb = STEEL

    title_slide()
    bullets("What we make", [
        "Standard fasteners + formed parts",
        "Custom runs with repeat tooling",
        "Kitting/packaging for distributors",
    ])
    bullets("Operations priorities", [
        "Align OTD definitions (customer vs internal)",
        "Reduce scrap + rework via press setup discipline",
        "Improve inventory accuracy (ERP vs spreadsheets)",
    ])
    bullets("Capacity", [
        "Tacoma plant; Press Line 4 added 2025",
        "Kent 3PL for overflow + shipping",
    ])

    prs.save(BASE / "marketing-pitch-deck.pptx")


def main():
    make_markdown()
    make_json()
    make_pdfs()
    make_xlsx()
    make_pptx()


if __name__ == "__main__":
    main()
