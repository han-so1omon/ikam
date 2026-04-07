from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Meridian SignalWorks — Mission, Vision, Values

## Mission
Turn messy operational signals into clean, trustworthy workflows.

## Vision
Be the default operational signal layer for mid-market teams that need accountability and speed.

## Values
- **Clarity wins.**
- **Definitions matter.**
- **Build for operators.**
- **Secure by default.**
- **Own the incident.**
""",
    )

    w(
        "brand-guide.md",
        """# Meridian SignalWorks — Brand Guide (v0.1)

## Voice
- Calm, precise, trustworthy.
- Use concrete nouns; avoid hype.

## Palette
- Deep navy: #0B1220
- Off-white: #F6F7FB
- Teal: #21B6B8
- Slate: #64748B

## UI / diagram style
- Rounded rectangles, subtle shadows
- Teal for highlights, slate for secondary

## Naming mess (known)
- Feature names vs codenames vs SKUs drift; include aliases in internal docs.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Mina Okafor
Last updated: 2026-01-30 (draft)

## 2026 bets
1) Stabilize ingestion pipelines (queue backlog risk)
2) Ship Feature Orchid (scope expanded)
3) Make KPI definitions explicit in-product

## Risks
- Roadmap drift (slides vs Jira)
- Security policy drift (policy vs runbook)
""",
    )

    w(
        "product-overview.md",
        """# Product Overview

Meridian SignalWorks ingests signals (events, docs, logs), normalizes entities, and triggers workflows.

Modules:
- Meridian Core — ingestion + normalization
- Meridian Automations — workflows
- Meridian Insights — dashboards

Feature Orchid (codename): new reconciliation layer + “definition registry”.
""",
    )

    w(
        "prd-orchid-v1.md",
        """# PRD — Orchid (v1)

Status: draft
Date: 2025-12-05

## Problem
KPI definitions drift across teams and customers; dashboards become untrustworthy.

## Proposal
A definition registry + reconciliation flows.

## Non-goals
- Full MDM
""",
    )

    w(
        "prd-orchid-v3-FINAL.md",
        """# PRD — Orchid (v3 FINAL)

Status: FINAL (circulated)
Date: 2026-01-22

## Changes since v1
- Scope expanded: adds queue safeguards + incident timeline view.

## Ship target
- External roadmap claims Q2
- Internal Jira estimates Q3 (not reflected here)
""",
    )

    w(
        "release-notes-2025-q4.md",
        """# Release Notes — 2025 Q4

- Improved ingestion retries (partial)
- New dashboard: Atlas Metrics (beta)
- Security: policy update drafted (not fully rolled out)
""",
    )

    w(
        "architecture-overview.md",
        """# Architecture Overview (high-level)

Key components:
- Ingestion API
- Queue (events)
- Normalizer service
- Workflow engine
- Metrics pipeline

Known risk: queue backlog under burst traffic.
""",
    )

    w(
        "runbook-queue-backlog.md",
        """# Runbook — Queue Backlog

## Symptoms
- ingestion latency increases
- dashboard lag

## Actions
- scale consumers
- apply rate limits
- drain backlog

## Notes
- Runbook assumes 180-day key rotation cadence.
""",
    )

    w(
        "incident-2026-01-17-postmortem.md",
        """# Incident Postmortem — 2026-01-17 (Queue backlog)

Summary: elevated ingestion latency and delayed workflow triggers.

## Narrative
Primary factor: upstream vendor latency amplified our queue.

## Contributing factors
- Consumer autoscaling lag
- Missing alert on queue depth

## Action items
- Add queue-depth SLO
- Add backpressure controls
""",
    )

    w(
        "security-policy-access-control.md",
        """# Security Policy — Access Control

- Access reviews: quarterly
- Key rotation: every 90 days
- Break-glass access: logged + reviewed
""",
    )

    w(
        "key-rotation-runbook.md",
        """# Key Rotation Runbook

Cadence: rotate keys every 180 days (operationally realistic).

Notes:
- Policy document says 90 days; we have not operationalized that yet.
""",
    )

    w(
        "exec-notes-2026-01-05.md",
        """# Exec Notes — 2026-01-05

- Orchid scope creep: now includes incident timeline + definition registry.
- Finance preparing board update; ARR number still debated.
""",
    )

    w(
        "product-review-notes-2026-01-19.md",
        """# Product Review — 2026-01-19

- Roadmap slide says Orchid Q2.
- Jules says Jira slipped to Q3.
- Mina wants external comms to stay Q2 “if possible”.
""",
    )

    w(
        "voice-note-transcript-2026-02-02.md",
        """# Voice Note Transcript — 2026-02-02 (Riley)

- Board deck says 24.8M ARR.
- My model says 23.9M because I exclude churned annual prepaid that’s still in cash.
- Everyone keeps mixing ARR with run-rate bookings.
""",
    )

    w(
        "incident-postmortem-template.md",
        """# Incident Postmortem Template

## Summary

## Impact

## Timeline

## Root cause

## Contributing factors

## What went well

## What didn't

## Action items
""",
    )

    w(
        "prd-template.md",
        """# PRD Template

## Problem

## Users

## Goals / Non-goals

## Requirements

## Risks

## Metrics
""",
    )


def make_json():
    wjson(
        "org-chart.json",
        {
            "generated": str(date.today()),
            "org": {
                "CEO": "Harper Lin",
                "CTO": "Theo Ramirez",
                "VP Product": "Mina Okafor",
                "Dir Platform Eng": "Jules Chen",
                "Head Security": "Sana Iqbal",
                "Head CS": "Owen Patel",
                "Finance Lead": "Riley Morgan",
            },
        },
    )

    wjson(
        "kanban-snapshot-2026-01.json",
        {
            "snapshot_date": "2026-01-21",
            "board": "Platform",
            "columns": {
                "Backlog": ["ORCHID-101", "ORCHID-122", "QUEUE-88"],
                "In Progress": ["ORCHID-140", "QUEUE-91"],
                "Review": ["SEC-33"],
                "Done": ["ATLAS-12"],
            },
        },
    )

    wjson(
        "incident-2026-01-17-metrics.json",
        {
            "incident_id": "INC-2026-01-17",
            "signals": {
                "vendor_latency_ms_p95": 240,
                "queue_depth_max": 980000,
                "consumer_lag_seconds_max": 7400,
                "internal_processing_errors": 0.02,
            },
            "notes": "Metrics suggest backlog depth and consumer lag were primary drivers; vendor latency was secondary.",
        },
    )

    wjson(
        "top-customers.json",
        {
            "generated": str(date.today()),
            "customers": [
                {"account": "Evergreen Ledger", "aka": ["Evergreen Financial"], "segment": "fintech"},
                {"account": "Harborline Health Services", "segment": "health services"},
                {"account": "Cascadia Freight Ops", "segment": "logistics"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "ARR (board)", "definition": "Annualized recurring revenue incl annual prepaid until renewal date", "notes": "Optimistic"},
                {"name": "ARR (finance)", "definition": "Annualized recurring revenue excluding churned annual prepaid", "notes": "Conservative"},
                {"name": "NRR", "definition": "Net revenue retention", "notes": "Definition varies by product add-ons"},
            ],
        },
    )

    wjson(
        "customer-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Customer Intake",
            "fields": [
                {"id": "legal_name", "type": "string", "required": True},
                {"id": "preferred_name", "type": "string", "required": False},
                {"id": "industry", "type": "string", "required": True},
                {"id": "kpis_in_scope", "type": "array", "items": {"type": "string"}},
                {"id": "data_sources", "type": "array", "items": {"type": "string"}},
            ],
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "board-update-2026-02.pdf",
        "Board Update — Feb 2026 (Summary)",
        [
            "ARR: $24.8M (board definition)",
            "Reliability: queue backlog incident in Jan",
            "Roadmap: Orchid targeted Q2 (public) — internal slip risk",
        ],
    )


def make_xlsx():
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill

    # Revenue history
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Quarterly"
    ws.append(["Year", "Quarter", "Revenue", "Notes"])
    rows = [
        (2024, "Q4", 4100000, ""),
        (2025, "Q1", 4600000, "Pricing change prep"),
        (2025, "Q2", 5200000, "Usage add-on introduced"),
        (2025, "Q3", 5600000, "Churn event"),
        (2025, "Q4", 6000000, "Orchid scope expanded"),
    ]
    for r in rows:
        ws.append(list(r))
    for c in ws[1]:
        c.fill = PatternFill("solid", fgColor="0B1220")
        c.font = Font(color="FFFFFF", bold=True)
        c.alignment = Alignment(horizontal="center")
    for r in range(2, ws.max_row + 1):
        ws.cell(r, 3).number_format = "$#,##0"
    wb.save(BASE / "quarterly-revenue-history-2024-2025.xlsx")

    # ARR model
    wb2 = openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "ARR Model"
    ws2.append(["Customer", "MRR", "Annual Prepaid?", "Churned?", "ARR (board)", "ARR (finance)"])
    data = [
        ("Evergreen Ledger", 620000, "Y", "N"),
        ("Harborline Health Services", 410000, "N", "N"),
        ("Cascadia Freight Ops", 360000, "Y", "Y"),
    ]
    for cust, mrr, prepaid, churned in data:
        ws2.append([cust, mrr, prepaid, churned, None, None])

    for c in ws2[1]:
        c.fill = PatternFill("solid", fgColor="21B6B8")
        c.font = Font(color="FFFFFF", bold=True)

    for r in range(2, ws2.max_row + 1):
        mrr = ws2.cell(r, 2).coordinate
        prepaid = ws2.cell(r, 3).coordinate
        churned = ws2.cell(r, 4).coordinate
        # Board includes annual prepaid even if churned until renewal (simplified as include always if prepaid=='Y')
        ws2.cell(r, 5).value = f"=IF({prepaid}=\"Y\",{mrr}*12,{mrr}*12)"
        # Finance excludes churned prepaid
        ws2.cell(r, 6).value = f"=IF(AND({prepaid}=\"Y\",{churned}=\"Y\"),0,{mrr}*12)"
        for col in [2, 5, 6]:
            ws2.cell(r, col).number_format = "$#,##0"

    # Totals
    end = ws2.max_row + 1
    ws2.cell(end, 4).value = "TOTAL"
    ws2.cell(end, 5).value = f"=SUM(E2:E{ws2.max_row})"
    ws2.cell(end, 6).value = f"=SUM(F2:F{ws2.max_row})"
    ws2.cell(end, 5).number_format = "$#,##0"
    ws2.cell(end, 6).number_format = "$#,##0"

    wb2.save(BASE / "arr-model-2026.xlsx")

    # Sprint log
    wb3 = openpyxl.Workbook()
    ws3 = wb3.active
    ws3.title = "Sprints Q4 2025"
    ws3.append(["Sprint", "Start", "End", "Committed", "Completed", "Carryover", "Notes"])
    sprints = [
        ("S-2025-20", "2025-10-06", "2025-10-17", 120, 98, None, ""),
        ("S-2025-21", "2025-10-20", "2025-10-31", 130, 110, None, "Orchid prep"),
        ("S-2025-22", "2025-11-03", "2025-11-14", 140, 92, None, "Scope churn"),
        ("S-2025-23", "2025-11-17", "2025-11-28", 125, 88, None, "Holiday"),
        ("S-2025-24", "2025-12-01", "2025-12-12", 150, 101, None, "Orchid expanded"),
    ]
    for sp in sprints:
        ws3.append(list(sp))
    for r in range(2, ws3.max_row + 1):
        comm = ws3.cell(r, 4).coordinate
        comp = ws3.cell(r, 5).coordinate
        ws3.cell(r, 6).value = f"={comm}-{comp}"
    wb3.save(BASE / "sprint-log-2025-q4.xlsx")

    # Pipeline + renewal risk
    wb4 = openpyxl.Workbook()
    ws4 = wb4.active
    ws4.title = "Pipeline Q1 2026"
    ws4.append(["Account", "Stage", "Close", "Amount", "Probability", "Weighted"])
    deals = [
        ("JuniperPay", "Negotiation", "2026-03-15", 2400000, 0.55),
        ("NorthHarbor Ops", "Discovery", "2026-02-28", 1200000, 0.35),
    ]
    for a, st, cl, amt, pr in deals:
        ws4.append([a, st, cl, amt, pr, None])
    for r in range(2, ws4.max_row + 1):
        amt = ws4.cell(r, 4).coordinate
        pr = ws4.cell(r, 5).coordinate
        ws4.cell(r, 6).value = f"={amt}*{pr}"
        ws4.cell(r, 4).number_format = "$#,##0"
        ws4.cell(r, 6).number_format = "$#,##0"
        ws4.cell(r, 5).number_format = "0%"
    wb4.save(BASE / "pipeline-2026-q1.xlsx")

    wb5 = openpyxl.Workbook()
    ws5 = wb5.active
    ws5.title = "Renewal Risk"
    ws5.append(["Customer", "Renewal", "ARR", "Health", "Risk drivers"])
    rows = [
        ("Cascadia Freight Ops", "2026-04-01", 0, "Red", "Churned but still in cash view"),
        ("Evergreen Ledger", "2026-09-01", 7440000, "Yellow", "Reliability concerns"),
    ]
    for r in rows:
        ws5.append(list(r))
    for r in range(2, ws5.max_row + 1):
        ws5.cell(r, 3).number_format = "$#,##0"
    wb5.save(BASE / "customer-renewal-risk.xlsx")

    # Vendor risk register
    wb6 = openpyxl.Workbook()
    ws6 = wb6.active
    ws6.title = "Vendors"
    ws6.append(["Vendor", "Service", "Risk", "Notes"])
    vs = [
        ("QuasarQueue", "Queue infrastructure", "Med", "Latency spikes suspected"),
        ("NimbusAuth", "Auth provider", "Low", "Stable"),
    ]
    for v in vs:
        ws6.append(list(v))
    wb6.save(BASE / "vendor-risk-register.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    NAVY = RGBColor(0x0B, 0x12, 0x20)
    OFF = RGBColor(0xF6, 0xF7, 0xFB)
    TEAL = RGBColor(0x21, 0xB6, 0xB8)

    def deck(path: Path, title: str, subtitle: str, bullets_by_slide: list[tuple[str, list[str]]]):
        prs = Presentation()

        def header(slide):
            f = slide.background.fill
            f.solid(); f.fore_color.rgb = OFF
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
            bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
            bar.line.fill.background()
            acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
            acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
            acc.line.fill.background()

        # Title
        s = prs.slides.add_slide(prs.slide_layouts[6])
        header(s)
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Georgia"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = NAVY
        sb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11), Inches(0.7))
        p2 = sb.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Arial"; p2.font.size = Pt(20); p2.font.color.rgb = NAVY

        for t, items in bullets_by_slide:
            s2 = prs.slides.add_slide(prs.slide_layouts[6])
            header(s2)
            tt = s2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
            pp = tt.text_frame.paragraphs[0]
            pp.text = t
            pp.font.name = "Georgia"; pp.font.size = Pt(34); pp.font.bold = True; pp.font.color.rgb = NAVY
            box = s2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(4.2))
            tf = box.text_frame
            tf.clear()
            for i, it in enumerate(items):
                pr = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pr.text = it
                pr.font.name = "Arial"; pr.font.size = Pt(20); pr.font.color.rgb = NAVY

        prs.save(path)

    deck(
        BASE / "roadmap-2026-h1.pptx",
        "Roadmap — H1 2026",
        "Meridian SignalWorks",
        [
            ("Q1", ["Stabilize ingestion", "Queue backlog SLO", "Atlas Metrics rollout"]),
            ("Q2 (public)", ["Orchid ships (public) — risk: internal slip", "Definition registry beta"]),
        ],
    )

    deck(
        BASE / "sales-deck-v2.pptx",
        "Meridian SignalWorks",
        "Operational signal parsing + workflow automation",
        [
            ("What it does", ["Ingest signals", "Normalize entities", "Trigger workflows", "Explain definitions"]),
            ("Proof points", ["Atlas Metrics dashboard", "QueueGuard reliability work", "Orchid reconciliation"]),
            ("Security", ["Access control policy exists", "Runbook drift acknowledged"]),
        ],
    )

    deck(
        BASE / "cs-qbr-template.pptx",
        "QBR Template",
        "Customer Success",
        [
            ("Health + outcomes", ["SLOs", "Adoption", "Risks"]),
            ("Roadmap", ["Customer asks", "Product commitments"]),
        ],
    )


def main():
    make_markdown()
    make_json()
    make_pdfs()
    make_xlsx()
    make_pptx()


if __name__ == "__main__":
    main()
