from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Juniper & Juno Goods — Mission, Vision, Values

## Mission
Make gifting feel thoughtful and local by curating home goods that people actually use.

## Vision
A small chain boutique known for warm stores, seasonal drops, and reliable fulfillment.

## Values
- Curate with taste
- Keep it human
- Margins matter
- Don’t let promos break ops
""",
    )

    w(
        "brand-guide.md",
        """# Juniper & Juno Goods — Brand Guide (v0.1)

## Palette
- Sage: #5B7A67
- Cream: #F4EFE6
- Terracotta: #C46A4A
- Ink: #1A1A1A

## Voice
- Warm, specific, not precious.

## Known drift
- Promo names vary across teams.
- SKU naming partially standardized.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Ben Carter
Last updated: 2026-02-01

## Priorities
1) Inventory accuracy across stores and warehouse
2) Promo discipline (postmortems required)
3) Reconcile gross vs net revenue reporting
""",
    )

    w(
        "campaign-brief-spring.md",
        """# Campaign Brief — Spring Collection

Owner: Keisha Wong

## Objective
Increase repeat purchase and drive foot traffic to Venice store.

## Promo naming (mess)
- Spring Reset
- Refresh Week
- New Season Drop
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Seasonal drops
2) Giftable staples
3) Warm in-store experience
4) Simple shipping
""",
    )

    w(
        "product-style-guide.md",
        """# Product Style Guide

## Product naming
- Name + material + size

## Photography
- Natural light
- Minimal props
- Avoid brand-y backgrounds
""",
    )

    w(
        "store-manager-call-2026-01-09.md",
        """# Store Manager Call — 2026-01-09

- Venice store opened strong but inventory accuracy is rough.
- Warehouse sheet says candle backorder; stores say overstock.
""",
    )

    w(
        "store-manager-call-2026-02-06.md",
        """# Store Manager Call — 2026-02-06

- Ops weekly sales uses gross; bookkeeping uses net.
- Promo naming drift is causing customer confusion.
""",
    )

    w(
        "voice-note-transcript-2026-02-13.md",
        """# Voice Note Transcript — 2026-02-13 (West Ledger)

- Bookkeeping reports net sales; ops sends gross.
- Returns spike during major promos.
- We need one definition in the KPI sheet.
""",
    )

    w(
        "vendor-scorecard-template.md",
        """# Vendor Scorecard Template

Vendor:
Category:

Scores (1–5):
- Quality
- Lead time
- Packaging consistency
- Responsiveness

Notes:
""",
    )

    w(
        "promo-postmortem-template.md",
        """# Promo Postmortem Template

Promo:
Dates:

## What happened

## Results

## Issues

## Fixes
""",
    )


def make_json():
    wjson(
        "vendor-list.json",
        {
            "generated": str(date.today()),
            "vendors": [
                {"id": "V001", "name": "SageCo Candles", "category": "candles", "risk": "lead time"},
                {"id": "V002", "name": "PaperRill Stationery", "category": "stationery", "risk": "low"},
                {"id": "V003", "name": "TerraCeram", "category": "ceramics", "risk": "breakage"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "Revenue (ops)", "definition": "gross sales", "notes": "includes discounts before reconciliation"},
                {"name": "Revenue (bookkeeping)", "definition": "net sales", "notes": "excludes returns/voids"},
                {"name": "Sell-through", "definition": "units sold / units received"},
                {"name": "AOV", "definition": "net sales / orders"},
            ],
        },
    )

    wjson(
        "wholesale-order-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Wholesale Order Intake",
            "fields": [
                {"id": "store_name", "type": "string", "required": True},
                {"id": "contact_email", "type": "string", "required": True},
                {"id": "sku_list", "type": "array", "items": {"type": "string"}},
                {"id": "ship_by", "type": "date"},
            ],
        },
    )


def make_xlsx():
    import openpyxl
    from openpyxl.styles import Font, PatternFill

    # promo calendar
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="H1 2026"
    ws.append(["Week","Promo name","Channels","Notes"])
    for r in [("2026-W03","Spring Reset","All","aka Refresh Week"),("2026-W07","New Season Drop","Stores + IG","alias used in ads")]:
        ws.append(list(r))
    for c in ws[1]:
        c.fill=PatternFill("solid", fgColor="5B7A67"); c.font=Font(color="FFFFFF", bold=True)
    wb.save(BASE/"promo-calendar-2026-h1.xlsx")

    # revenue history
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="Quarterly"
    ws2.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q4",1400000,"2 stores"),(2025,"Q1",1550000,""),(2025,"Q2",1680000,""),(2025,"Q3",1750000,"online growth"),(2025,"Q4",2100000,"promo returns spike")]:
        ws2.append(list(r))
    for i in range(2, ws2.max_row+1): ws2.cell(i,3).number_format="$#,##0"
    wb2.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    # projection
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="2026"
    ws3.append(["Month","Revenue","Notes"])
    for r in [("2026-01",520000,""),("2026-02",490000,""),("2026-03",560000,"")]:
        ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,2).number_format="$#,##0"
    wb3.save(BASE/"projected-revenue-2026.xlsx")

    # weekly performance gross vs net
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-02"
    ws4.append(["Week","Ops Gross","Bookkeeping Net","Notes"])
    for r in [("2026-W05",168000,158500,"discounts/returns"),("2026-W06",175500,164200,"promo naming drift")]:
        ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="$#,##0"; ws4.cell(i,3).number_format="$#,##0"
    wb4.save(BASE/"store-performance-weekly-2026-02.xlsx")

    # replenishment plan
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Replenishment"
    ws5.append(["SKU","Item","Vendor","Reorder point","Target stock","Notes"])
    for r in [("CND-SGE-08","Sage candle 8oz","SageCo",40,120,"stores claim overstock"),("STN-NTE-01","Notebook","PaperRill",30,90,"")]:
        ws5.append(list(r))
    wb5.save(BASE/"inventory-replenishment-plan.xlsx")

    # inventory adjustments
    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="2026-01"
    ws6.append(["Date","Location","SKU","Adjustment","Reason","Notes"])
    for r in [("2026-01-14","Venice","CND-SGE-08",-12,"damage",""),("2026-01-20","Warehouse","CND-SGE-08",+24,"late PO","warehouse shows backorder")]:
        ws6.append(list(r))
    wb6.save(BASE/"inventory-adjustments-2026-01.xlsx")

    # labor scheduling template
    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="Template"
    ws7.append(["Store","Role","Mon","Tue","Wed","Thu","Fri","Sat","Sun","Total hrs"])
    ws7.append(["Venice","Retail associate",8,8,8,8,8,10,10,None])
    ws7.cell(2,10).value="=SUM(C2:I2)"
    wb7.save(BASE/"labor-scheduling-template.xlsx")

    # incident log
    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="Q1 2026"
    ws8.append(["Date","Location","Incident","Severity","Notes"])
    for r in [("2026-01-22","Online","Shipping delay spike","Med","carrier issue"),("2026-02-03","Silver Lake","POS outage 30m","Low","")]:
        ws8.append(list(r))
    wb8.save(BASE/"incident-log-2026-q1.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    SAGE=RGBColor(0x5B,0x7A,0x67)
    CRM=RGBColor(0xF4,0xEF,0xE6)
    TERR=RGBColor(0xC4,0x6A,0x4A)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=CRM
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=SAGE; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=TERR; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Juniper & Juno Goods"; p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
