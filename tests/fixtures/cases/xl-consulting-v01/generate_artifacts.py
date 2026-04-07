from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_md():
    w(
        "mission-vision-values.md",
        """# Alderpoint Strategy & Operations — Mission, Vision, Values

Mission: help large organizations change with measurable outcomes.
Vision: the consulting partner known for pragmatic transformation.
Values: clarity, integrity, disciplined definitions.
""",
    )

    w(
        "brand-guide.md",
        """# ASO — Brand Guide (v0.1)

Palette:
- Deep blue: #0B1F3A
- Graphite: #1F2937
- White: #F8FAFC
- Gold accent: #D4AF37

Voice: confident, concise, specific.
""",
    )

    w(
        "program-charter.md",
        """# Program Charter — GraniteWorks Utilities — Operations Transformation 2026

Program director: Devin Holt
PMO lead: Talia Reed

Objectives:
- Reduce operating cost
- Improve reliability KPIs

Savings reporting: executive deck uses annualized gross; finance model uses net.
""",
    )

    w(
        "client-email-thread-2026-03-18.md",
        """# Client Email Thread — 2026-03-18

Subject: Phase 1 milestone date

Client: "Confirm Phase 1 complete date?"
ASO PMO: "Re-baselined to 2026-03-29 due to dependency delays."

Note: RAID log still shows 2026-03-15.
""",
    )

    w(
        "slide-template-guidelines.md",
        """# Slide Template Guidelines

- Use deep blue header bar, light background.
- Avoid rewriting numbers by hand; link to the model.
- Mark DRAFT vs FINAL explicitly.
""",
    )

    w(
        "meeting-notes-template.md",
        """# Meeting Notes Template

Attendees:
Agenda:
Decisions:
Actions:
Risks:
""",
    )

    w(
        "raid-template.md",
        """# RAID Template

Risk | Assignee | Due | Status
Issue | Owner | Next step
""",
    )

    w(
        "issue-escalation-template.md",
        """# Issue Escalation Template

Issue:
Impact:
Decision needed:
Options:
Owner:
""",
    )


def make_json():
    wjson(
        "org-structure.json",
        {
            "generated": str(date.today()),
            "leaders": {
                "ManagingPartner": "Colleen Park",
                "ProgramDirector": "Devin Holt",
                "OpsLead": "Mira Shah",
                "FinanceLead": "Andre Kim",
                "PMO": "Talia Reed",
            },
            "practices": ["Strategy", "Operations", "Finance", "Analytics"],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-03-10",
            "kpis": [
                {"name": "Savings (exec)", "definition": "annualized gross run-rate"},
                {"name": "Savings (finance)", "definition": "net savings after one-time costs"},
                {"name": "Phase 1 complete", "definition": "milestone date per re-baselined plan"},
            ],
        },
    )

    wjson(
        "deliverable-index.json",
        {
            "program": "GraniteWorks Utilities — Ops Transformation 2026",
            "latest": {
                "exec_deck": "exec-deck-2026-03.pptx",
                "ops_workstream": "workstream-deck-ops-v7.pptx",
                "finance_workstream": "workstream-deck-finance-v5.pptx",
            },
            "notes": "Some teams still circulate older 'final' decks.",
        },
    )


def make_xlsx():
    import openpyxl

    # Change log
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="SOW Changes"
    ws.append(["Date","Change","Approved By","Notes"])
    ws.append(["2026-02-20","Added analytics enablement","Client sponsor","" ])
    wb.save(BASE/"statement-of-work-change-log.xlsx")

    # RAID log: milestone date conflict
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="RAID"
    ws2.append(["Type","ID","Description","Owner","Due","Status"])
    ws2.append(["Milestone","M1","Phase 1 complete","PMO","2026-03-15","At risk" ])
    ws2.append(["Risk","R12","Baseline dispute","Andre","2026-03-08","Open" ])
    wb2.save(BASE/"pmo-raID-log.xlsx")

    # Master program plan
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Plan"
    ws3.append(["Workstream","Milestone","Date","Notes"])
    ws3.append(["PMO","Phase 1 complete","2026-03-29","re-baselined" ])
    wb3.save(BASE/"master-program-plan-2026.xlsx")

    # Savings model contradiction
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="Savings"
    ws4.append(["Initiative","Gross annualized","One-time cost","Net savings"])
    ws4.append(["Field dispatch optimization", 8_200_000, 1_100_000, None])
    ws4.append(["Procurement renegotiation", 6_900_000, 850_000, None])
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,4).value=f"=B{i}-C{i}"
        for c in [2,3,4]: ws4.cell(i,c).number_format="$#,##0"
    # total
    ws4.append(["TOTAL", None, None, None])
    t=ws4.max_row
    ws4.cell(t,2).value=f"=SUM(B2:B{t-1})"
    ws4.cell(t,3).value=f"=SUM(C2:C{t-1})"
    ws4.cell(t,4).value=f"=SUM(D2:D{t-1})"
    for c in [2,3,4]: ws4.cell(t,c).number_format="$#,##0"
    wb4.save(BASE/"savings-model-2026.xlsx")

    # Benefits tracker shows exec number (gross)
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="2026-03"
    ws5.append(["Metric","Value","Definition"])
    ws5.append(["Annualized savings", 18_400_000, "gross (exec deck)" ])
    ws5.append(["Net savings", 14_900_000, "net (finance model)" ])
    ws5.cell(2,2).number_format="$#,##0"; ws5.cell(3,2).number_format="$#,##0"
    wb5.save(BASE/"benefits-tracker-2026-03.xlsx")

    # Timesheet totals (2,120 hours)
    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="Timesheet"
    ws6.append(["Week","Role","Hours"])
    rows=[("2026-W10","Consultants",540),("2026-W11","Consultants",560),("2026-W12","Consultants",520),("2026-W13","Consultants",500)]
    for r in rows: ws6.append(list(r))
    ws6.append(["TOTAL","",None])
    ws6.cell(ws6.max_row,3).value=f"=SUM(C2:C{ws6.max_row-1})"
    wb6.save(BASE/"timesheet-2026-03.xlsx")

    # Revenue history
    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="Quarterly"
    ws7.append(["Year","Quarter","Revenue"])
    for r in [(2024,"Q4",12_800_000),(2025,"Q1",13_200_000),(2025,"Q2",14_100_000),(2025,"Q3",13_900_000),(2025,"Q4",15_600_000)]:
        ws7.append(list(r)); ws7.cell(ws7.max_row,3).number_format="$#,##0"
    wb7.save(BASE/"revenue-history-2024-2025.xlsx")

    # Pipeline
    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="Pipeline"
    ws8.append(["Client","Stage","Close","Value","Probability","Weighted"])
    ws8.append(["Rivermark Insurance","Negotiation","2026-06-30",3_200_000,0.55,None])
    ws8.append(["Sable Foods","Discovery","2026-08-15",1_800_000,0.35,None])
    for i in range(2, ws8.max_row+1):
        ws8.cell(i,6).value=f"=D{i}*E{i}"
        ws8.cell(i,4).number_format="$#,##0"; ws8.cell(i,6).number_format="$#,##0"; ws8.cell(i,5).number_format="0%"
    wb8.save(BASE/"pipeline-enterprise.xlsx")


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "practice-catalog.pdf",
        "ASO — Practice Catalog",
        ["Strategy", "Operations", "Finance", "Analytics", "PMO"],
    )

    pdf(
        BASE / "msa.pdf",
        "Master Services Agreement (Sample) — ASO",
        ["Confidentiality", "Billing terms", "IP", "Limitations"],
    )

    pdf(
        BASE / "sow-master.pdf",
        "Statement of Work — GraniteWorks Utilities",
        [
            "Scope: operations transformation + PMO.",
            "Savings reporting definitions vary by stakeholder.",
            "Out-of-scope items require change control.",
        ],
    )

    # invoice mismatch: billed 1,980 vs timesheet 2,120
    pdf(
        BASE / "invoice-summary-2026-03.pdf",
        "Invoice Summary — 2026-03",
        [
            "Timesheet hours: 2,120",
            "Billed hours: 1,980 (cap + write-offs)",
            "Notes: client requested cap in March; non-billable internal rework excluded.",
        ],
    )


def make_pptx():
    from pptx import Presentation
    prs=Presentation(); prs.save(BASE/"exec-deck-2026-03.pptx")
    prs=Presentation(); prs.save(BASE/"workstream-deck-ops-v7.pptx")
    prs=Presentation(); prs.save(BASE/"workstream-deck-finance-v5.pptx")


def main():
    make_md(); make_json(); make_xlsx(); make_pdfs(); make_pptx()


if __name__ == "__main__":
    main()
