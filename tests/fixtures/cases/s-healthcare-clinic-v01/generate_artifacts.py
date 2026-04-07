from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Pinecrest Family Clinic — Mission, Vision, Values

## Mission
Provide calm, accessible primary care with clear communication.

## Vision
A neighborhood clinic known for consistency and patient trust.

## Values
- Plain-language care
- Privacy by default
- Fix the system
""",
    )

    w(
        "brand-guide.md",
        """# Pinecrest Family Clinic — Brand Guide (v0.1)

Palette:
- Blue: #1D4ED8
- Teal: #14B8A6
- White: #F8FAFC

Voice:
- Calm, reassuring, direct.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Jordan Kim
Last updated: 2026-01-08

Priorities:
- Keep billing clean (denials)
- Maintain training completion
- Reduce reminder-related no-shows
""",
    )

    w(
        "clinic-ops-handbook.md",
        """# Clinic Ops Handbook

- Check-in workflow
- Reminder cadence (SMS + phone)
- No-show follow-up
""",
    )

    w(
        "billing-process-map.md",
        """# Billing Process Map

1) Eligibility
2) Coding
3) Claim submission
4) Denials
5) Patient statements
""",
    )

    w(
        "privacy-policy.md",
        """# Privacy Policy

- Annual HIPAA training required
- Minimum necessary access
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Friendly primary care
2) Same-week appointments
3) Clear billing expectations
""",
    )

    w(
        "staff-huddle-notes-2026-01-09.md",
        """# Staff Huddle Notes — 2026-01-09

- Reminder workflow switched to new SMS tool (simpler).
- Track missed reminders manually for now.
""",
    )

    w(
        "staff-huddle-notes-2026-02-06.md",
        """# Staff Huddle Notes — 2026-02-06

- Denials categorized differently by billing vendor.
- Collections lag for one payer due to processing delay.
""",
    )

    w(
        "patient-complaint-template.md",
        """# Patient Complaint Template

Date:
Complaint:
Resolution:
Owner:
""",
    )

    w(
        "incident-report-template.md",
        """# Incident Report Template

Date:
Incident:
Severity:
Actions:
Owner:
""",
    )

    w(
        "performance-review-template.md",
        """# Performance Review Template

Role:
Period:

Outcomes:
Strengths:
Growth:
""",
    )


def make_json():
    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-01-08",
            "kpis": [
                {"name": "Collections", "definition": "payments received / amounts billed"},
                {"name": "Denial rate", "definition": "denied claims / submitted claims"},
            ],
        },
    )

    wjson(
        "patient-intake-form.json",
        {
            "schema_version": "0.2",
            "title": "Patient Intake",
            "fields": [
                {"id": "full_name", "type": "string", "required": True},
                {"id": "dob", "type": "date", "required": True},
                {"id": "insurance", "type": "string"},
                {"id": "reason", "type": "string"},
            ],
            "notes": "Synthetic only; avoid PHI in benchmark fills.",
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "employee-handbook-excerpt.pdf",
        "Employee Handbook Excerpt (summary)",
        ["Privacy expectations", "Attendance", "Training requirements"],
    )

    pdf(
        BASE / "service-line-one-pager.pdf",
        "Services — Overview",
        ["Primary care", "Preventive visits", "Basic labs/referrals"],
    )


def make_xlsx():
    import openpyxl

    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Template"
    ws.append(["Role","Mon","Tue","Wed","Thu","Fri","Total hrs"])
    ws.append(["Front desk",8,8,8,8,8,None])
    ws.append(["MA",8,8,8,8,8,None])
    for r in [2,3]:
        ws.cell(r,7).value=f"=SUM(B{r}:F{r})"
    wb.save(BASE/"staffing-schedule-template.xlsx")

    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="Q1 2026"
    ws2.append(["Date","Incident","Severity","Notes"])
    for r in [("2026-01-28","Reminder SMS outage 2h","Low",""),("2026-02-10","Eligibility denials spike","Med","payer portal issues")]:
        ws2.append(list(r))
    wb2.save(BASE/"incident-log-2026-q1.xlsx")

    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="2026-01"
    ws3.append(["Metric","Value","Notes"])
    ws3.append(["Collections",0.91,"one payer lag"])
    ws3.cell(2,2).number_format="0.0%"
    wb3.save(BASE/"collections-summary-2026-01.xlsx")

    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-01"
    ws4.append(["Payer","Denial rate","Top reason","Notes"])
    for r in [("Payer A",0.058,"eligibility",""),("Payer B",0.044,"coding","category changed")]:
        ws4.append(list(r))
    for i in range(2, ws4.max_row+1): ws4.cell(i,2).number_format="0.0%"
    wb4.save(BASE/"denials-tracker-2026-01.xlsx")

    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="2026"
    ws5.append(["Dept","Required","Completed","Completion %"])
    ws5.append(["Clinic",12,12,None])
    ws5.cell(2,4).value="=C2/B2"; ws5.cell(2,4).number_format="0.0%"
    wb5.save(BASE/"hipaa-training-tracker-2026.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    BLUE=RGBColor(0x1D,0x4E,0xD8)
    TEAL=RGBColor(0x14,0xB8,0xA6)
    OFF=RGBColor(0xF8,0xFA,0xFC)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=BLUE; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=TEAL; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Pinecrest Family Clinic"; p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
