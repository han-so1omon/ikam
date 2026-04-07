from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Northlake Advisory Group — Mission, Vision, Values

## Mission
Help mid-market operators make decisions faster by turning messy operational reality into clear, executable plans.

## Vision
Be the consulting partner that clients trust for outcomes — not decks.

## Values
- **Operator first:** practical over theoretical.
- **Earn the right:** show your work; document assumptions.
- **Truth over harmony:** surface the real constraints.
- **Templates are tools:** not a substitute for thinking.
- **Own the handoff:** implementation beats recommendations.
""",
    )

    w(
        "brand-guide.md",
        """# Northlake Advisory Group — Brand Guide (v0.1)

## Voice + tone
- Clear, specific, slightly blunt (in a helpful way)
- Avoid hype words. Prefer verbs and numbers.

## Visual palette (suggested)
- Navy: #0B1F3A
- Off-white: #F7F5F0
- Teal: #2AA6A1
- Gray: #6B7280

## Slide style
- Title bar in Navy
- One accent line in Teal
- Lots of whitespace
- Use simple icons, not stock photos

## Naming + file hygiene
- Client aliases exist (legal name vs codename). Always include both in internal artifacts.
- Deck versions proliferate. Don’t overwrite; append a date or version.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Elena Park
Last updated: 2026-01-29 (draft; not fully adopted)

## Where we win
- Mid-market operators who are drowning in ad-hoc reporting.
- Engagements where we can do both diagnosis and implementation.

## What we sell (in practice)
- Diagnostics (2–4 weeks)
- 90-day transformation sprints (marketed as standardized; delivered variably)
- Retainers (advisory + light analytics)

## 2026 bets
1) Package “90-day sprint” deliverables more cleanly
2) Reduce template sprawl (one canonical set)
3) Improve pipeline discipline (definitions + close dates)

## Risks
- Partners bypass process.
- Client numbers are inconsistent (we need reconciliation artifacts by default).
""",
    )

    w(
        "service-lines.md",
        """# Service Lines

## Ops Diagnostic
- 2–4 weeks
- Output: bottleneck map + prioritized backlog + quick wins

## Finance Transformation
- Close process, reporting definitions, variance analysis

## Cost Model + Pricing
- Build a driver-based cost model, identify margin leaks

## Training Workshops
- Operator playbooks (KPIs, cadences, meeting design)
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) **Clarity under pressure** — reconcile the numbers, name the constraints.
2) **Operator-led transformation** — we build with your team, not around it.
3) **Artifacts that survive** — models, templates, and a handoff you can run.
""",
    )

    # Meeting notes + voice transcript
    w(
        "all-hands-notes-2025-09-03.md",
        """# All Hands Notes — 2025-09-03

Attendees: Elena, Ravi, Casey, Morgan, Priya, Jo + others

## Updates
- PINECONE (Evergreen) kickoff went well; data access is slower than promised.
- Template overhaul: please use the new deck skeleton (link… somewhere).

## Concerns
- Too many "final" decks floating around.
- Staffing plan says 18 consultants; roster reality differs.
""",
    )

    w(
        "partner-meeting-notes-2025-10-10.md",
        """# Partner Meeting Notes — 2025-10-10

Attendees: Elena, Ravi

## Pricing dispute (PINECONE)
- Elena: keep fixed fee; client wants certainty
- Ravi: push for change-order language; scope creep likely

## Decision
- Proposal goes out today (version unclear).
""",
    )

    w(
        "voice-note-transcript-2025-10-21.md",
        """# Voice Note Transcript — 2025-10-21 (Casey)

- Pipeline is lying to us. We have deals marked "at risk" that are basically already delivering.
- Finance booked something as recognized but HubSpot still shows it as stage 4.
- We need one client naming rule: legal name + codename on every artifact.
- Deck versioning is out of control.
""",
    )

    w(
        "engagement-playbook.md",
        """# Engagement Playbook (draft)

## Cadence (what the handbook says)
- Two-week sprints
- Weekly status note
- Monthly steering committee

## Cadence (what happens)
- Varies by partner + client urgency
- Status notes are inconsistent

## Required artifacts
- SOW
- KPI definitions (client + ours)
- Delivery plan
- Reconciliation note if numbers disagree
""",
    )

    # Sample engagement folder
    w(
        "sample-engagement-PINECONE/kickoff-notes.md",
        """# Kickoff Notes — PINECONE (Evergreen)

Client: Evergreen Transport (aka Evergreen Logistics LLC)
Codename: PINECONE

## Goals (client stated)
- Cut late shipments
- Reduce overtime
- Get a weekly operations dashboard that leadership trusts

## Data reality
- Dispatch data exists but definitions differ by region.
- Finance definitions differ from ops reports.
""",
    )

    w(
        "sample-engagement-PINECONE/weekly-status-2025-08-15.md",
        """# Weekly Status — 2025-08-15 (PINECONE)

## Wins
- Mapped current-state dispatch workflow

## Risks
- Client KPI definitions are inconsistent ("on-time" is not consistent)

## Next
- Draft reconciliation note
""",
    )

    w(
        "sample-engagement-PINECONE/weekly-status-2025-08-29.md",
        """# Weekly Status — 2025-08-29 (PINECONE)

## Wins
- Built v1 cost-to-serve model (rough)

## Risks
- Scope creep: client asked for pricing recommendations

## Next
- Deliverable deck v3 "FINAL" in progress
""",
    )

    w(
        "sample-engagement-PINECONE/closeout-summary.md",
        """# Closeout Summary — PINECONE

Status: draft (not widely circulated)

## What changed
- Standardized KPI definitions for on-time performance
- Implemented weekly ops cadence

## Open items
- Client data pipeline still fragile
""",
    )

    # Evaluation templates
    w(
        "engagement-retro-template.md",
        """# Engagement Retro Template

## What we expected

## What happened

## What we learned

## What we’ll change next time

## Artifacts to link
- SOW
- Deliverable deck
- Cost model
- KPI definitions
""",
    )

    w(
        "consultant-performance-rubric.md",
        """# Consultant Performance Rubric (internal)

Categories (1–5):
- Client trust
- Analytical rigor
- Artifact quality
- Ownership
- Communication

Notes:
- Promotions should reference concrete artifacts delivered.
""",
    )


def make_json():
    wjson(
        "lead-list.json",
        {
            "generated": str(date.today()),
            "owner": "Ravi Desai",
            "leads": [
                {
                    "account": "Evergreen Logistics LLC",
                    "aka": ["Evergreen Transport"],
                    "codename": "PINECONE",
                    "stage": "4 - Negotiation",
                    "est_close": "2025-07-12",
                    "amount": 180000,
                    "notes": "Delivery started before CRM close was updated.",
                },
                {
                    "account": "Harborview Medical Services",
                    "codename": "MARLIN",
                    "stage": "3 - Discovery",
                    "est_close": "2025-09-30",
                    "amount": 95000,
                },
                {
                    "account": "Cascadia Components",
                    "codename": "RIVET",
                    "stage": "2 - Qualified",
                    "est_close": "2025-10-28",
                    "amount": 120000,
                },
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2025-08-20",
            "kpis": [
                {
                    "name": "Booked revenue",
                    "definition": "Signed SOW value",
                    "source": "CRM + signed PDFs",
                },
                {
                    "name": "Recognized revenue",
                    "definition": "Revenue recognized by month per finance close",
                    "source": "Finance workbook",
                },
                {
                    "name": "Utilization %",
                    "definition": "Billable hours / total available hours",
                    "source": "Timesheets + staffing plan",
                },
            ],
        },
    )

    wjson(
        "client-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Client Intake",
            "fields": [
                {"id": "legal_name", "type": "string", "required": True},
                {"id": "preferred_name", "type": "string", "required": False},
                {"id": "codename", "type": "string", "required": True},
                {"id": "sponsor", "type": "string", "required": True},
                {"id": "problem_statement", "type": "string", "required": True},
                {"id": "success_metrics", "type": "array", "required": False, "items": {"type": "string"}},
                {"id": "data_sources", "type": "array", "required": False, "items": {"type": "string"}},
            ],
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def build_pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    build_pdf(
        BASE / "case-study-one-pager.pdf",
        "Case Study (sanitized) — PINECONE (Evergreen)",
        [
            "Client: Evergreen Transport (aka Evergreen Logistics LLC)",
            "Problem: late shipments + inconsistent reporting",
            "Approach: KPI reconciliation + operating cadence + cost-to-serve model",
            "Outcome: leadership-aligned dashboard definitions; early operational improvements",
            "Note: numbers differ by source; reconciliation memo attached in the engagement folder",
        ],
    )

    build_pdf(
        BASE / "sample-engagement-PINECONE/statement-of-work.pdf",
        "Statement of Work — PINECONE (Evergreen)",
        [
            "Term: 8 weeks",
            "Fee: $180,000 fixed (change orders allowed)",
            "Deliverables: KPI definitions, cost model, operating cadence, deliverable deck",
            "Client sponsor: VP Operations",
        ],
    )


def make_xlsx_finance_and_ops():
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill

    # Quarterly revenue history
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Quarterly Revenue"
    ws.append(["Year", "Quarter", "Booked Revenue", "Recognized Revenue", "Notes", "Variance", "Variance %"])

    rows = [
        (2024, "Q4", 420000, 390000, "Strong close", None, None),
        (2025, "Q1", 380000, 410000, "Recognition lag", None, None),
        (2025, "Q2", 460000, 455000, "Stable", None, None),
        # Intentional contradiction anchor: Q3 recognizes PINECONE while pipeline shows at-risk
        (2025, "Q3", 510000, 535000, "Includes PINECONE recognition", None, None),
        (2025, "Q4", 600000, 580000, "Year-end push", None, None),
    ]
    for r in rows:
        ws.append(list(r))

    # Style
    fill = PatternFill("solid", fgColor="0B1F3A")
    for c in ws[1]:
        c.fill = fill
        c.font = Font(color="FFFFFF", bold=True)
        c.alignment = Alignment(horizontal="center", wrap_text=True)

    # Formulas for variance
    for r in range(2, ws.max_row + 1):
        booked = ws.cell(r, 3).coordinate
        rec = ws.cell(r, 4).coordinate
        ws.cell(r, 6).value = f"={rec}-{booked}"
        ws.cell(r, 7).value = f"=IF({booked}=0,0,({rec}-{booked})/{booked})"
        ws.cell(r, 3).number_format = "$#,##0"
        ws.cell(r, 4).number_format = "$#,##0"
        ws.cell(r, 6).number_format = "$#,##0"
        ws.cell(r, 7).number_format = "0.0%"

    ws.freeze_panes = "A2"
    wb.save(BASE / "quarterly-revenue-history-2024-2025.xlsx")

    # Projection
    wb2 = openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "2026 Projection"
    ws2.append(["Month", "Retainers", "Projects", "Workshops", "Total", "Assumed GM%", "COGS (est)", "Gross Profit (est)", "Notes"])

    months = [
        ("2026-01", 120000, 210000, 15000, 0.46, "GTM refresh draft"),
        ("2026-02", 125000, 190000, 20000, 0.46, "Pipeline discipline push"),
        ("2026-03", 130000, 240000, 22000, 0.46, "Hiring contractor analysts"),
    ]
    for m, ret, proj, wsh, gm, notes in months:
        ws2.append([m, ret, proj, wsh, None, gm, None, None, notes])

    for cell in ws2[1]:
        cell.fill = PatternFill("solid", fgColor="2AA6A1")
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)

    for r in range(2, ws2.max_row + 1):
        ret = ws2.cell(r, 2).coordinate
        proj = ws2.cell(r, 3).coordinate
        wsh = ws2.cell(r, 4).coordinate
        total = ws2.cell(r, 5)
        gm = ws2.cell(r, 6).coordinate
        total.value = f"={ret}+{proj}+{wsh}"
        cogs = ws2.cell(r, 7)
        gp = ws2.cell(r, 8)
        cogs.value = f"={total.coordinate}*(1-{gm})"
        gp.value = f"={total.coordinate}-{cogs.coordinate}"
        for c in [2, 3, 4, 5, 7, 8]:
            ws2.cell(r, c).number_format = "$#,##0"
        ws2.cell(r, 6).number_format = "0.0%"

    wb2.save(BASE / "projected-revenue-2026.xlsx")

    # Pipeline report (xlsx)
    wb3 = openpyxl.Workbook()
    ws3 = wb3.active
    ws3.title = "Pipeline Q3 2025"
    ws3.append(["Account", "Codename", "Stage", "Probability", "Est Close", "Amount", "Weighted", "Notes"])

    deals = [
        ("Evergreen Logistics LLC", "PINECONE", "At risk", 0.45, "2025-07-12", 180000, None, "Delivery started; CRM not updated"),
        ("Harborview Medical Services", "MARLIN", "Discovery", 0.35, "2025-09-30", 95000, None, "Waiting on data access"),
        ("Cascadia Components", "RIVET", "Qualified", 0.25, "2025-10-28", 120000, None, "Pricing sensitivity"),
    ]
    for d in deals:
        ws3.append(list(d))

    for cell in ws3[1]:
        cell.fill = PatternFill("solid", fgColor="0B1F3A")
        cell.font = Font(color="FFFFFF", bold=True)

    for r in range(2, ws3.max_row + 1):
        prob = ws3.cell(r, 4).coordinate
        amt = ws3.cell(r, 6).coordinate
        ws3.cell(r, 7).value = f"={prob}*{amt}"
        ws3.cell(r, 6).number_format = "$#,##0"
        ws3.cell(r, 7).number_format = "$#,##0"
        ws3.cell(r, 4).number_format = "0%"

    ws3.cell(ws3.max_row + 2, 5).value = "Total weighted"
    ws3.cell(ws3.max_row, 7).value = f"=SUM(G2:G{ws3.max_row-2})"

    wb3.save(BASE / "pipeline-report-2025-q3.xlsx")

    # Utilization tracker
    wb4 = openpyxl.Workbook()
    ws4 = wb4.active
    ws4.title = "Utilization H2 2025"
    ws4.append(["Name", "Role", "FTE?", "Available Hours", "Billable Hours", "Utilization", "Notes"])

    people = [
        ("Morgan Lee", "Sr Consultant", "Y", 520, 455, None, ""),
        ("Priya Shah", "Analyst", "Y", 520, 410, None, ""),
        ("(Contractor) Dana Kim", "Analyst", "N", 300, 240, None, "Not counted in staffing plan"),
        ("(Contractor) Luis Moreno", "Consultant", "N", 320, 280, None, ""),
    ]
    for p in people:
        ws4.append(list(p))

    for cell in ws4[1]:
        cell.fill = PatternFill("solid", fgColor="6B7280")
        cell.font = Font(color="FFFFFF", bold=True)

    for r in range(2, ws4.max_row + 1):
        avail = ws4.cell(r, 4).coordinate
        bill = ws4.cell(r, 5).coordinate
        ws4.cell(r, 6).value = f"=IF({avail}=0,0,{bill}/{avail})"
        ws4.cell(r, 6).number_format = "0.0%"

    wb4.save(BASE / "utilization-tracker-2025-h2.xlsx")


def make_pptx_decks():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    NAVY = RGBColor(0x0B, 0x1F, 0x3A)
    OFFWHITE = RGBColor(0xF7, 0xF5, 0xF0)
    TEAL = RGBColor(0x2A, 0xA6, 0xA1)
    GRAY = RGBColor(0x55, 0x55, 0x55)

    def build(path: Path, *, tagline: str, pinecone_status_line: str):
        prs = Presentation()

        def bg(slide):
            f = slide.background.fill
            f.solid()
            f.fore_color.rgb = OFFWHITE

        def header(slide):
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
            bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
            bar.line.fill.background()
            acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
            acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
            acc.line.fill.background()

        def title_slide():
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); header(s)
            tx = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.2))
            p = tx.text_frame.paragraphs[0]
            p.text = "Northlake Advisory Group"
            p.font.name = "Georgia"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = NAVY
            sub = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11), Inches(0.8))
            p2 = sub.text_frame.paragraphs[0]
            p2.text = tagline
            p2.font.name = "Arial"; p2.font.size = Pt(20); p2.font.color.rgb = GRAY

        def bullets(title, items):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); header(s)
            t = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
            p = t.text_frame.paragraphs[0]
            p.text = title
            p.font.name = "Georgia"; p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY
            box = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(4.2))
            tf = box.text_frame
            tf.clear()
            for i, it in enumerate(items):
                pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pp.text = it
                pp.font.name = "Arial"; pp.font.size = Pt(20); pp.font.color.rgb = GRAY

        title_slide()
        bullets("What we do", [
            "Ops diagnostics + finance reporting definition",
            "90-day transformation sprints (practical implementation)",
            "Retainers for ongoing advisory",
        ])
        bullets("Selected engagement (sanitized)", [
            "Client: Evergreen Transport / Evergreen Logistics LLC (codename PINECONE)",
            pinecone_status_line,
            "Core work: KPI reconciliation + cadence + cost-to-serve model",
        ])
        bullets("How we work", [
            "Handbook says two-week sprints",
            "Reality: cadence varies; we document decisions and assumptions",
            "Deliverables: models + definitions + handoff",
        ])
        bullets("Next steps", [
            "Identify 1–2 pilot clients",
            "Agree KPI definitions up front",
            "Set a weekly decision cadence",
        ])

        prs.save(path)

    # Deck variants include a contradiction about PINECONE status
    build(BASE / "marketing-pitch-deck.pptx", tagline="Operations + finance transformation for mid-market operators", pinecone_status_line="Status: at risk in CRM (but delivery already started)")
    build(BASE / "sample-engagement-PINECONE/deliverable-deck-v1.pptx", tagline="Deliverable deck v1", pinecone_status_line="Status: closed/won (per delivery kickoff)")
    build(BASE / "sample-engagement-PINECONE/deliverable-deck-v3-FINAL.pptx", tagline="Deliverable deck v3 FINAL", pinecone_status_line="Status: recognized in finance close (Q3 2025)")


def main():
    make_markdown()
    make_json()
    make_pdfs()
    make_xlsx_finance_and_ops()
    make_pptx_decks()


if __name__ == "__main__":
    main()
