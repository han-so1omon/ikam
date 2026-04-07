from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Harborview Family Clinics Network — Mission, Vision, Values

## Mission
Provide accessible, high-quality care with clear communication and reliable operations.

## Vision
A clinic network patients trust for consistent care and predictable experiences.

## Values
- Patient dignity
- Definitions matter
- Privacy by default
- Fix the system
""",
    )

    w(
        "brand-guide.md",
        """# HFCN — Brand Guide (v0.1)

## Palette
- Clinic Blue: #1D4ED8
- Teal: #14B8A6
- White: #F8FAFC
- Slate: #64748B

## Voice
- Calm, plain-language, reassuring.
- Avoid jargon where possible.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Marco Diaz (COO)
Last updated: 2026-02-02

## Priorities
1) Standardize no-show definitions + interventions.
2) Improve collections with clearer patient communications.
3) Roll out privacy policy v4 + close audit findings.
4) Decide future of membership pilot.
""",
    )

    w(
        "clinic-ops-handbook-v2.md",
        """# Clinic Ops Handbook (v2)

Effective: 2024-08-01

- Check-in workflow
- No-show handling (late cancels excluded)
- Staffing ratios guidance
""",
    )

    w(
        "clinic-ops-handbook-v3.md",
        """# Clinic Ops Handbook (v3 draft)

Drafted: 2026-01-10

Changes vs v2:
- No-show includes late cancels within 24h
- Standard reminder cadence

Known issue:
- Sites still using v2 for training.
""",
    )

    w(
        "billing-process-map.md",
        """# Billing Process Map (high-level)

1) Eligibility verification
2) Encounter coding
3) Claim submission
4) Denials management
5) Patient statements

Known risk: denials spike after 2024-09 system change.
""",
    )

    w(
        "privacy-policy-v3.md",
        """# Privacy Policy (v3)

Effective: 2024-05-01

- Annual training required
- Minimum necessary access
""",
    )

    w(
        "privacy-policy-v4.md",
        """# Privacy Policy (v4)

Drafted: 2025-12-15

Changes:
- Clarified access review cadence
- Training deadline: Jan 31

Known issue:
- Rollout incomplete; some sites still reference v3.
""",
    )

    w(
        "site-director-call-2026-01-11.md",
        """# Site Director Call — 2026-01-11

- No-show debate: ops says 7.5% (no late cancels), sites say 10–12%.
- Staffing spreadsheet differs from HR roster.
""",
    )

    w(
        "site-director-call-2026-02-08.md",
        """# Site Director Call — 2026-02-08

- Collections definition disagreement continues.
- Compliance says training is "on track"; HR tracker shows many incomplete.
""",
    )

    w(
        "voice-note-transcript-2026-02-14.md",
        """# Voice Note Transcript — 2026-02-14 (Tara)

- Denials tracker says collections 89% if you exclude catch-up payments.
- Finance reports 94% including late payments.
- Sites count late cancels as no-shows; ops doesn’t.
""",
    )

    w(
        "patient-complaint-template.md",
        """# Patient Complaint Template

Date:
Site:
Complaint:
Resolution:
Owner:
""",
    )

    w(
        "incident-report-template.md",
        """# Incident Report Template

Date:
Site:
Incident:
Severity:
Actions taken:
Owner:
""",
    )

    w(
        "new-hire-onboarding-checklist.md",
        """# New Hire Onboarding Checklist

- HIPAA training
- Systems access
- Role shadowing
- Policy acknowledgement
""",
    )

    w(
        "performance-review-template.md",
        """# Performance Review Template

Role:
Period:

## Outcomes

## Strengths

## Growth

## Next goals
""",
    )


def make_json():
    wjson(
        "org-structure.json",
        {
            "generated": str(date.today()),
            "leaders": {
                "CMO": "Dr. Renee Lawson",
                "COO": "Marco Diaz",
                "RevenueCycle": "Tara Nguyen",
                "Compliance": "Omar Shah",
                "HR": "Jenna Cole",
                "Finance": "Vinh Tran",
            },
            "sites": 9,
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "No-show (ops)", "definition": "no-shows excluding late cancels", "notes": "lower"},
                {"name": "No-show (sites)", "definition": "no-shows including late cancels <24h", "notes": "higher"},
                {"name": "Collections (finance)", "definition": "collections incl catch-up payments", "notes": "optimistic"},
                {"name": "Collections (revcycle)", "definition": "collections excluding catch-up", "notes": "strict"},
            ],
        },
    )

    wjson(
        "patient-intake-form.json",
        {
            "schema_version": "0.2",
            "title": "Patient Intake",
            "fields": [
                {"id": "full_name", "type": "string", "required": True},
                {"id": "dob", "type": "date", "required": True},
                {"id": "insurance", "type": "string"},
                {"id": "primary_reason", "type": "string"},
                {"id": "consents", "type": "array", "items": {"type": "string"}},
            ],
            "notes": "Use synthetic data only in benchmarks; avoid PHI in sample fills.",
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "audit-memo-2026-02.pdf",
        "Privacy Audit Memo — Feb 2026 (Summary)",
        [
            "Finding: training completion tracking inconsistent",
            "Finding: policy v4 not fully rolled out",
            "Action: centralize acknowledgements",
        ],
    )

    pdf(
        BASE / "board-update-2026-02.pdf",
        "Board Update — Feb 2026 (Summary)",
        [
            "Collections: 94% (finance definition)",
            "No-show: 7.5% (ops definition)",
            "Privacy: rollout in progress",
        ],
    )

    pdf(
        BASE / "employee-handbook-excerpt.pdf",
        "Employee Handbook Excerpt (summary)",
        [
            "Privacy expectations",
            "Attendance policy",
            "Training requirements",
        ],
    )

    pdf(
        BASE / "service-line-one-pager.pdf",
        "Service Lines — Overview",
        [
            "Primary care",
            "Urgent care",
            "Occupational health",
            "Labs/referrals",
        ],
    )


def make_xlsx():
    import openpyxl

    # Revenue history
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Quarterly"
    ws.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q3",18200000,""),(2024,"Q4",19500000,"urgent care integration"),(2025,"Q1",18800000,""),(2025,"Q2",20100000,""),(2025,"Q3",21000000,"denials spike aftermath"),(2025,"Q4",22500000,"")]:
        ws.append(list(r))
    for i in range(2, ws.max_row+1): ws.cell(i,3).number_format="$#,##0"
    wb.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    # Projection
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="2026"
    ws2.append(["Month","Revenue","Notes"])
    for r in [("2026-01",7200000,""),("2026-02",6900000,""),("2026-03",7400000,"")]:
        ws2.append(list(r))
    for i in range(2, ws2.max_row+1): ws2.cell(i,2).number_format="$#,##0"
    wb2.save(BASE/"projected-revenue-2026.xlsx")

    # Staffing model
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Model"
    ws3.append(["Role","Per clinic FTE","Clinics","Total FTE"])
    rows=[("Front desk",2.2,9,None),("MA",3.0,9,None),("Provider",2.5,9,None)]
    for r in rows: ws3.append(list(r))
    for i in range(2, ws3.max_row+1):
        ws3.cell(i,4).value=f"=B{i}*C{i}"
    wb3.save(BASE/"staffing-model.xlsx")

    # Site director scorecard (no-show contradiction)
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-01"
    ws4.append(["Site","No-show (ops)","No-show (site)","Notes"])
    for r in [("Clinic 1",0.072,0.11,"late cancels"),("Clinic 2",0.078,0.10,"")]: ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="0.0%"; ws4.cell(i,3).number_format="0.0%"
    wb4.save(BASE/"site-director-scorecard-2026-01.xlsx")

    # Denials tracker
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="2026-01"
    ws5.append(["Payer","Denial rate","Top reason","Notes"])
    for r in [("Payer A",0.082,"eligibility","system change fallout"),("Payer B",0.065,"coding","")]: ws5.append(list(r))
    for i in range(2, ws5.max_row+1): ws5.cell(i,2).number_format="0.0%"
    wb5.save(BASE/"denials-tracker-2026-01.xlsx")

    # Payer mix
    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="2025"
    ws6.append(["Payer type","Share"])
    for r in [("Commercial",0.55),("Medicare",0.25),("Medicaid",0.15),("Self-pay",0.05)]: ws6.append(list(r))
    for i in range(2, ws6.max_row+1): ws6.cell(i,2).number_format="0.0%"
    wb6.save(BASE/"payer-mix-2025.xlsx")

    # Collections summary contradiction
    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="2026-01"
    ws7.append(["Metric","Value","Definition","Notes"])
    ws7.append(["Collections (finance)",0.94,"incl catch-up",""])
    ws7.append(["Collections (revcycle)",0.89,"excl catch-up",""])
    ws7.cell(2,2).number_format="0.0%"; ws7.cell(3,2).number_format="0.0%"
    wb7.save(BASE/"collections-summary-2026-01.xlsx")

    # HIPAA training tracker
    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="2026"
    ws8.append(["Dept","Required","Completed","Completion %","Notes"])
    rows=[("Clinical",180,142,None,"behind"),("Front desk",90,71,None,""),("Billing",55,52,None,"")]
    for r in rows: ws8.append(list(r))
    for i in range(2, ws8.max_row+1):
        ws8.cell(i,4).value=f"=C{i}/B{i}"
        ws8.cell(i,4).number_format="0.0%"
    wb8.save(BASE/"hipaa-training-tracker-2026.xlsx")

    # Hiring plan
    wb9=openpyxl.Workbook(); ws9=wb9.active; ws9.title="2026"
    ws9.append(["Dept","Role","Count","Priority","Notes"])
    for r in [("Clinics","MA",12,"High",""),("Revenue cycle","Denials specialist",3,"Med",""),("Compliance","Trainer",1,"Low","")]: ws9.append(list(r))
    wb9.save(BASE/"hiring-plan-2026.xlsx")

    # Incident log
    wb10=openpyxl.Workbook(); ws10=wb10.active; ws10.title="Q1 2026"
    ws10.append(["Date","Site","Incident","Severity","Notes"])
    for r in [("2026-01-19","Clinic 4","No-show reminder system failed","Med",""),("2026-02-03","Billing","Denials backlog spike","High","payer response lag")]: ws10.append(list(r))
    wb10.save(BASE/"incident-log-2026-q1.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    BLUE=RGBColor(0x1D,0x4E,0xD8)
    TEAL=RGBColor(0x14,0xB8,0xA6)
    OFF=RGBColor(0xF8,0xFA,0xFC)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=BLUE; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=TEAL; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Harborview Family Clinics Network"; p.font.name="Georgia"; p.font.size=Pt(40); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
