from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# LatticeOps — Mission, Vision, Values

Mission: help ops teams run workflows with clarity.
Vision: a mid-market ops platform customers trust.
Values: definitions, reliability, pragmatic delivery.
""",
    )

    w(
        "brand-guide.md",
        """# LatticeOps — Brand Guide (v0.1)

Palette:
- Navy: #0B1220
- Teal: #14B8A6
- Off-white: #F6F7FB

Voice: calm, precise.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Alex Chen

- Ship Garnet reliably
- Align KPI definitions
- Reduce backlog risk
""",
    )

    w(
        "product-overview.md",
        """# Product Overview

LatticeOps ingests ops signals and powers workflows + dashboards.

Feature Garnet: KPI definition registry + reconciliation UI.
""",
    )

    w(
        "prd-garnet-v1.md",
        """# PRD — Garnet (v1)

Date: 2026-01-10
Status: draft
""",
    )

    w(
        "prd-garnet-v3-FINAL.md",
        """# PRD — Garnet (v3 FINAL)

Date: 2026-02-12
Status: FINAL

Ship target: April (roadmap) but sprint log suggests May.
""",
    )

    w(
        "release-notes-2025-q4.md",
        """# Release Notes — 2025 Q4

- New workflow templates
- Partial backlog improvements
""",
    )

    w(
        "runbook-backlog.md",
        """# Runbook — Backlog Spike

Symptoms: delayed workflows.
Actions: scale workers, pause imports, drain backlog.
""",
    )

    w(
        "incident-2026-02-09-postmortem.md",
        """# Incident Postmortem — 2026-02-09

Narrative: vendor slowness caused backlog.
Action: add queue alerts.
""",
    )

    w(
        "exec-notes-2026-02-03.md",
        """# Exec Notes — 2026-02-03

- Garnet scope is growing.
- KPI definitions are inconsistent across teams.
""",
    )

    w(
        "product-review-notes-2026-02-17.md",
        """# Product Review — 2026-02-17

- Roadmap deck says Garnet ships April.
- Sam says May is realistic.
""",
    )

    w(
        "voice-note-transcript-2026-02-21.md",
        """# Voice Note Transcript — 2026-02-21 (Morgan)

- CS prefers active accounts.
- Product prefers weekly active users.
- Finance wants billable accounts.
""",
    )

    w(
        "incident-postmortem-template.md",
        """# Incident Postmortem Template

Summary:
Impact:
Timeline:
Root cause:
Actions:
""",
    )

    w(
        "prd-template.md",
        """# PRD Template

Problem:
Users:
Goals:
Risks:
Metrics:
""",
    )


def make_json():
    wjson(
        "org-chart.json",
        {
            "generated": str(date.today()),
            "leaders": {
                "CEO": "Jamie Ross",
                "CTO": "Priya Mehta",
                "Product": "Alex Chen",
                "Eng": "Sam Rivera",
                "CS": "Taylor Kim",
                "Finance": "Morgan Shah",
            },
        },
    )

    wjson(
        "kanban-snapshot-2026-02.json",
        {
            "snapshot_date": "2026-02-18",
            "board": "Platform",
            "columns": {
                "Backlog": ["GARNET-101", "GARNET-120"],
                "In Progress": ["QUEUE-33"],
                "Review": ["GARNET-132"],
                "Done": ["REL-2025Q4-4"],
            },
        },
    )

    wjson(
        "incident-2026-02-09-metrics.json",
        {
            "incident_id": "INC-2026-02-09",
            "signals": {
                "vendor_latency_ms_p95": 160,
                "queue_depth_max": 310000,
                "worker_lag_seconds_max": 2400,
            },
            "notes": "Metrics suggest internal backlog/worker lag was dominant; vendor latency secondary.",
        },
    )

    wjson(
        "top-customers.json",
        {
            "generated": str(date.today()),
            "customers": [
                {"account": "Bayline Logistics", "segment": "logistics"},
                {"account": "JuniperPay", "segment": "fintech"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-21",
            "kpis": [
                {"name": "Active accounts (CS)", "definition": "accounts with weekly engagement"},
                {"name": "Weekly active users (Product)", "definition": "users with >=1 action in 7 days"},
                {"name": "Billable accounts (Finance)", "definition": "accounts with active subscription"},
            ],
        },
    )

    wjson(
        "customer-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Customer Intake",
            "fields": [
                {"id": "legal_name", "type": "string", "required": True},
                {"id": "industry", "type": "string", "required": True},
                {"id": "kpis", "type": "array", "items": {"type": "string"}},
            ],
        },
    )


def make_xlsx():
    import openpyxl

    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Quarterly"
    ws.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q4",3800000,""),(2025,"Q1",4100000,""),(2025,"Q2",4400000,""),(2025,"Q3",4700000,""),(2025,"Q4",5200000,"")]: ws.append(list(r))
    for i in range(2, ws.max_row+1): ws.cell(i,3).number_format="$#,##0"
    wb.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="ARR"
    ws2.append(["Customer","MRR","ARR"])
    for r in [("JuniperPay",180000,None),("Bayline Logistics",150000,None)]: ws2.append(list(r))
    for i in range(2, ws2.max_row+1):
        ws2.cell(i,3).value=f"=B{i}*12"
        ws2.cell(i,2).number_format="$#,##0"; ws2.cell(i,3).number_format="$#,##0"
    wb2.save(BASE/"arr-model-2026.xlsx")

    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Sprints Q1 2026"
    ws3.append(["Sprint","Start","End","Committed","Completed","Carryover","Notes"])
    rows=[("S-01","2026-01-06","2026-01-17",90,80,None,""),("S-02","2026-01-20","2026-01-31",95,72,None,"queue work"),("S-03","2026-02-03","2026-02-14",100,65,None,"Garnet slip")]
    for r in rows: ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,6).value=f"=D{i}-E{i}"
    wb3.save(BASE/"sprint-log-2026-q1.xlsx")

    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="Pipeline Q2 2026"
    ws4.append(["Account","Stage","Close","Amount","Probability","Weighted"])
    for a,st,cl,amt,pr in [("CedarOps","Discovery","2026-05-30",950000,0.35),("Lakefront Labs","Negotiation","2026-06-20",1400000,0.55)]:
        ws4.append([a,st,cl,amt,pr,None])
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,6).value=f"=D{i}*E{i}"
        ws4.cell(i,4).number_format="$#,##0"; ws4.cell(i,6).number_format="$#,##0"; ws4.cell(i,5).number_format="0%"
    wb4.save(BASE/"pipeline-2026-q2.xlsx")

    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Renewals"
    ws5.append(["Customer","Renewal","ARR","Health","Notes"])
    ws5.append(["JuniperPay","2026-09-01",2160000,"Yellow","metric debate" ])
    ws5.cell(2,3).number_format="$#,##0"
    wb5.save(BASE/"renewal-risk.xlsx")


def make_pptx():
    from pptx import Presentation
    prs=Presentation(); prs.save(BASE/"roadmap-2026-h1.pptx")
    prs=Presentation(); prs.save(BASE/"sales-deck.pptx")
    prs=Presentation(); prs.save(BASE/"cs-qbr-template.pptx")


def main():
    make_markdown(); make_json(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
