from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Driftwood & Oak — Mission, Vision, Values

## Mission
Curate warm, functional home goods that make gifting easy and personal.

## Vision
Two great stores + a reliable online shop, without promos breaking operations.

## Values
- Taste over volume
- Keep it human
- Make the numbers make sense
""",
    )

    w(
        "brand-guide.md",
        """# Driftwood & Oak — Brand Guide (v0.1)

## Palette
- Oak brown: #7A5C3A
- Cream: #F4EFE6
- Muted blue: #4A6FA5
- Ink: #1A1A1A

## Voice
- Warm, specific, minimal hype.

## Known drift
- Promo names vary between calendar and social drafts.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Chris Nguyen
Last updated: 2026-02-01

## Priorities
1) Reduce inventory adjustment noise
2) Keep promo naming consistent
3) Reconcile weekly sales vs bookkeeping net
""",
    )

    w(
        "campaign-brief-spring.md",
        """# Campaign Brief — Spring Promo

Owner: Dani Park

Promo names in the wild:
- Spring Reset (calendar)
- New Season Drop (IG drafts)

Goal: drive foot traffic to East store and improve online conversion.
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Warm modern home goods
2) Giftable staples
3) Seasonal drops
""",
    )

    w(
        "product-style-guide.md",
        """# Product Style Guide

- Natural light photos
- Minimal props
- Name + material + size
""",
    )

    w(
        "store-call-2026-01-12.md",
        """# Store Call — 2026-01-12

- East store says candles are overstocked.
- Online says backorder.
- Need one inventory truth.
""",
    )

    w(
        "store-call-2026-02-09.md",
        """# Store Call — 2026-02-09

- Weekly sales report differs from bookkeeping by ~3%.
- Returns timing likely.
""",
    )

    w(
        "voice-note-transcript-2026-02-15.md",
        """# Voice Note Transcript — 2026-02-15

- Ops sends gross weekly sales.
- Bookkeeping reports net after returns.
- Promo naming drift is confusing attribution.
""",
    )

    w(
        "vendor-scorecard-template.md",
        """# Vendor Scorecard Template

Vendor:

Scores (1–5):
- Quality
- Lead time
- Packaging
- Responsiveness
""",
    )

    w(
        "promo-postmortem-template.md",
        """# Promo Postmortem Template

Promo:
Dates:

Results:
Issues:
Fixes:
""",
    )


def make_json():
    wjson(
        "vendor-list.json",
        {
            "generated": str(date.today()),
            "vendors": [
                {"id": "V001", "name": "Oakline Ceramics", "category": "ceramics"},
                {"id": "V002", "name": "BlueHarbor Textiles", "category": "textiles"},
                {"id": "V003", "name": "Candlecraft Co.", "category": "candles"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "Weekly sales (ops)", "definition": "gross sales"},
                {"name": "Weekly sales (bookkeeping)", "definition": "net sales after returns"},
                {"name": "Sell-through", "definition": "units sold / units received"},
            ],
        },
    )

    wjson(
        "online-order-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Online Order Intake",
            "fields": [
                {"id": "order_id", "type": "string", "required": True},
                {"id": "customer_email", "type": "string"},
                {"id": "items", "type": "array", "items": {"type": "string"}},
                {"id": "ship_by", "type": "date"},
            ],
        },
    )


def make_xlsx():
    import openpyxl
    from openpyxl.styles import Font, PatternFill

    # promo calendar
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="H1 2026"
    ws.append(["Week","Promo name","Channels","Notes"])
    for r in [("2026-W03","Spring Reset","All","aka New Season Drop"),("2026-W08","Home Refresh","Stores","" )]:
        ws.append(list(r))
    for c in ws[1]:
        c.fill=PatternFill("solid", fgColor="7A5C3A"); c.font=Font(color="FFFFFF", bold=True)
    wb.save(BASE/"promo-calendar-2026-h1.xlsx")

    # revenue history
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="Quarterly"
    ws2.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q4",820000,""),(2025,"Q1",860000,""),(2025,"Q2",910000,""),(2025,"Q3",940000,""),(2025,"Q4",1120000,"returns spike")]:
        ws2.append(list(r))
    for i in range(2, ws2.max_row+1): ws2.cell(i,3).number_format="$#,##0"
    wb2.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    # projection
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="2026"
    ws3.append(["Month","Revenue","Notes"])
    for r in [("2026-01",260000,""),("2026-02",245000,""),("2026-03",280000,"")]:
        ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,2).number_format="$#,##0"
    wb3.save(BASE/"projected-revenue-2026.xlsx")

    # weekly performance gross vs net
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-02"
    ws4.append(["Week","Ops Gross","Bookkeeping Net","Notes"])
    for r in [("2026-W05",61000,59200,"returns timing"),("2026-W06",64500,62500,"" )]:
        ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="$#,##0"; ws4.cell(i,3).number_format="$#,##0"
    wb4.save(BASE/"store-performance-weekly-2026-02.xlsx")

    # replenishment plan
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Replenishment"
    ws5.append(["SKU","Item","Vendor","Reorder point","Target stock","Notes"])
    for r in [("CND-SEA-08","Sea salt candle 8oz","Candlecraft",30,90,"stores claim overstock"),("MUG-CRM-01","Cream ceramic mug","Oakline",20,60,"")]:
        ws5.append(list(r))
    wb5.save(BASE/"inventory-replenishment-plan.xlsx")

    # inventory adjustments
    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="2026-01"
    ws6.append(["Date","Location","SKU","Adjustment","Reason","Notes"])
    for r in [("2026-01-10","East","CND-SEA-08",+18,"count correction",""),("2026-01-21","Online","CND-SEA-08",-24,"backorder","" )]:
        ws6.append(list(r))
    wb6.save(BASE/"inventory-adjustments-2026-01.xlsx")

    # labor scheduling template
    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="Template"
    ws7.append(["Store","Role","Mon","Tue","Wed","Thu","Fri","Sat","Sun","Total hrs"])
    ws7.append(["South","Retail associate",8,8,8,8,8,10,10,None])
    ws7.cell(2,10).value="=SUM(C2:I2)"
    wb7.save(BASE/"labor-scheduling-template.xlsx")

    # incident log
    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="Q1 2026"
    ws8.append(["Date","Location","Incident","Severity","Notes"])
    for r in [("2026-01-18","Online","Fulfillment backlog","Med",""),("2026-02-04","East","POS outage 20m","Low","")]:
        ws8.append(list(r))
    wb8.save(BASE/"incident-log-2026-q1.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    OAK=RGBColor(0x7A,0x5C,0x3A)
    CRM=RGBColor(0xF4,0xEF,0xE6)
    BLUE=RGBColor(0x4A,0x6F,0xA5)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=CRM
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=OAK; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=BLUE; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Driftwood & Oak"; p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
