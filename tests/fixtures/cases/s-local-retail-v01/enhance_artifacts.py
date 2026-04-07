from __future__ import annotations

from pathlib import Path

BASE = Path(__file__).resolve().parent


def enhance_brand_guide():
    p = BASE / "brand-guide.md"
    p.write_text(
        """# Bramble & Bitters — Brand Guide (v0.2)

Status: practical / small-team friendly
Owner: Maya
Last updated: 2026-02-05

## Brand feel (in plain English)
- **Cozy + curated** (not “luxury”, not “rustic cosplay”)
- **Neighborhood-smart**: specific recommendations, no buzzwords
- **Seasonal**: rotating moments, small drops, tasting-note energy

## Voice & tone
**Do:**
- Use sensory specificity (texture, aroma, ingredient origin)
- Be helpful and direct (“If they like X, get Y”)
- Keep it warm; mild wit is okay

**Don’t:**
- Overpromise (“life-changing”, “the best ever”)
- Trend-chase jargon (“elevated”, “authentic”, “curated” every sentence)

### Sample lines
- “A bright, peppery jam that loves a sharp cheddar.”
- “Bring this to dinner and you’ll look like you planned ahead.”

## Naming conventions (and known mess)
- Subscription program name is inconsistent in legacy posts:
  - **Bramble Box**
  - **Seasonal Pantry Box**
  - **Quarterly Crate**
- Initiatives get nicknames:
  - “Box Revamp” = subscription refresh

## Visual identity
### Color palette
Primary:
- **Bramble Green** — #1F3D2B
- **Cream** — #F4EFE6

Accent:
- **Copper** — #B87333
- **Ink** — #1A1A1A
- **Fog** — #D7D2C8

Usage rule of thumb:
- Backgrounds: Cream / Fog
- Headers + anchors: Bramble Green
- Accents (sparingly): Copper
- Body text: Ink

### Typography
If using what we already have:
- **Headers:** Georgia (or “Libre Baskerville” if available)
- **Body:** Inter (or Arial if needed)

Rules:
- Keep line length short; lots of white space.
- Use bold for emphasis; avoid all-caps except for tiny labels.

### Photography
- Natural light, slight grain is fine.
- Hands holding items, shelves, wrapping station.
- Avoid overly staged studio shots.

### Logo / mark (guidance, not strict)
- Wordmark should feel classic (serif) with a small “&”.
- Don’t put the logo on busy photos; use a cream block behind it.

## Layout patterns
### Instagram post pattern
- 1 hook line (benefit)
- 2–3 bullet-ish lines (what’s in it / why it’s good)
- 1 clear next step (come by / DM for corporate)

### One-pager pattern
- Who we are
- Bundles/offer
- Lead times
- Contact

## Accessibility
- Minimum contrast: Ink on Cream.
- Avoid Copper text on Cream for small type.

## File hygiene (lightweight)
- If it’s public-facing, include date in footer.
- Don’t overwrite old decks; save a new version (yes, even if it’s called “final-final”).
""",
        encoding="utf-8",
    )


def style_pptx(path: Path, *, variant: str):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Palette
    BRAMBLE = RGBColor(0x1F, 0x3D, 0x2B)
    CREAM = RGBColor(0xF4, 0xEF, 0xE6)
    COPPER = RGBColor(0xB8, 0x73, 0x33)
    INK = RGBColor(0x1A, 0x1A, 0x1A)

    prs = Presentation(path)

    def set_bg(slide, rgb):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def add_header_bar(slide):
        # top bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = BRAMBLE
        bar.line.fill.background()

        # copper accent line
        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06)
        )
        acc.fill.solid()
        acc.fill.fore_color.rgb = COPPER
        acc.line.fill.background()

    def restyle_text(shape):
        if not hasattr(shape, "text_frame"):
            return
        tf = shape.text_frame
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = "Georgia"
                run.font.color.rgb = INK

    for i, slide in enumerate(prs.slides):
        set_bg(slide, CREAM)
        add_header_bar(slide)

        # restyle title
        if slide.shapes.title:
            t = slide.shapes.title
            t.text_frame.paragraphs[0].font.name = "Georgia"
            t.text_frame.paragraphs[0].font.bold = True
            t.text_frame.paragraphs[0].font.size = Pt(40 if i == 0 else 34)
            t.text_frame.paragraphs[0].font.color.rgb = BRAMBLE

        # body placeholder
        for shape in slide.shapes:
            restyle_text(shape)

        # Add a small footer tag on all slides except title
        if i != 0:
            box = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.5), Inches(6), Inches(0.3))
            p = box.text_frame.paragraphs[0]
            p.text = f"Bramble & Bitters • {variant}"
            p.font.size = Pt(12)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        # Add a simple “sticker” circle on the title slide
        if i == 0:
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                prs.slide_width - Inches(2.2),
                Inches(1.1),
                Inches(1.6),
                Inches(1.6),
            )
            circ.fill.solid()
            circ.fill.fore_color.rgb = COPPER
            circ.line.color.rgb = COPPER
            tx = circ.text_frame
            tx.clear()
            p = tx.paragraphs[0]
            p.text = "Oakland\nCA"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = CREAM

    prs.save(path)


def enhance_xlsx_quarterly_history(path: Path):
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill

    wb = openpyxl.load_workbook(path)
    ws = wb["Quarterly Revenue"]

    # Add columns if not already present
    headers = [c.value for c in ws[1]]
    extra = [
        "Variance (POS - Book)",
        "Variance % (vs Book)",
        "YoY Growth (Book)",
    ]
    for h in extra:
        if h not in headers:
            ws.cell(row=1, column=len(headers) + 1, value=h)
            headers.append(h)

    # find indices
    col = {h: i + 1 for i, h in enumerate(headers)}

    # style header row (keep existing theme but ensure all headers styled)
    header_fill = PatternFill("solid", fgColor="2F4F4F")
    header_font = Font(color="FFFFFF", bold=True)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # formulas
    for r in range(2, ws.max_row + 1):
        book = ws.cell(r, col["Revenue (Bookkeeping)"]).coordinate
        pos = ws.cell(r, col["Revenue (POS)"]).coordinate

        ws.cell(r, col["Variance (POS - Book)"]).value = f"={pos}-{book}"
        ws.cell(r, col["Variance % (vs Book)"]).value = f"=IF({book}=0,0,({pos}-{book})/{book})"

        # YoY by matching quarter 4 rows above (since we have sequential quarters)
        # If r-4 exists, compute (this - prev)/prev else blank
        if r - 4 >= 2:
            prev_book = ws.cell(r - 4, col["Revenue (Bookkeeping)"]).coordinate
            ws.cell(r, col["YoY Growth (Book)"]).value = f"=IF({prev_book}=0,"",({book}-{prev_book})/{prev_book})"
        else:
            ws.cell(r, col["YoY Growth (Book)"]).value = ""

    # Formatting
    money_cols = ["Revenue (Bookkeeping)", "Revenue (POS)", "Variance (POS - Book)"]
    pct_cols = ["Variance % (vs Book)", "YoY Growth (Book)"]
    for r in range(2, ws.max_row + 1):
        for h in money_cols:
            ws.cell(r, col[h]).number_format = "$#,##0"
        for h in pct_cols:
            ws.cell(r, col[h]).number_format = "0.0%"

    ws.freeze_panes = "A2"

    # Column widths
    for h, w in {
        "Variance (POS - Book)": 18,
        "Variance % (vs Book)": 18,
        "YoY Growth (Book)": 16,
    }.items():
        ws.column_dimensions[openpyxl.utils.get_column_letter(col[h])].width = w

    wb.save(path)


def enhance_xlsx_projection(path: Path):
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill

    wb = openpyxl.load_workbook(path)
    ws = wb["2026 Projection"]

    # Ensure additional columns
    headers = [c.value for c in ws[1]]
    extras = [
        "COGS (est)",
        "Gross Profit (est)",
        "Marketing Spend (est)",
        "Operating Profit (est)",
    ]
    for h in extras:
        if h not in headers:
            ws.cell(row=1, column=len(headers) + 1, value=h)
            headers.append(h)

    col = {h: i + 1 for i, h in enumerate(headers)}

    # Restyle header row
    header_fill = PatternFill("solid", fgColor="1F4E79")
    header_font = Font(color="FFFFFF", bold=True)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Add formulas per row
    for r in range(2, ws.max_row + 1):
        if not ws.cell(r, col["Month"]).value:
            continue
        instore = ws.cell(r, col["In-store"]).coordinate
        online = ws.cell(r, col["Online"]).coordinate
        subs = ws.cell(r, col["Subscriptions"]).coordinate
        corp = ws.cell(r, col["Corporate"]).coordinate
        total = ws.cell(r, col["Total"]).coordinate
        gm = ws.cell(r, col["Assumed GM%"]).coordinate

        # Total formula (instead of value)
        ws.cell(r, col["Total"]).value = f"={instore}+{online}+{subs}+{corp}"

        # COGS = Total * (1 - GM)
        ws.cell(r, col["COGS (est)"]).value = f"={total}*(1-{gm})"
        ws.cell(r, col["Gross Profit (est)"]).value = f"={total}-{ws.cell(r, col['COGS (est)']).coordinate}"

        # Marketing spend: simple heuristic for small biz (2.5% of total, bump in Oct-Dec)
        month_cell = ws.cell(r, col["Month"]).coordinate
        # If month ends with -10/-11/-12 then 4% else 2.5%
        ws.cell(r, col["Marketing Spend (est)"]).value = (
            f"=IF(OR(RIGHT({month_cell},2)=\"10\",RIGHT({month_cell},2)=\"11\",RIGHT({month_cell},2)=\"12\"),{total}*0.04,{total}*0.025)"
        )

        # Operating profit: GP - marketing - fixed ops estimate (rent+utilities+base labor)
        # Fixed ops estimate: 38k/month baseline, 55k in Nov/Dec
        ws.cell(r, col["Operating Profit (est)"]).value = (
            f"={ws.cell(r, col['Gross Profit (est)']).coordinate}-{ws.cell(r, col['Marketing Spend (est)']).coordinate}-"
            f"IF(OR(RIGHT({month_cell},2)=\"11\",RIGHT({month_cell},2)=\"12\"),55000,38000)"
        )

    # number formats
    money = [
        "In-store",
        "Online",
        "Subscriptions",
        "Corporate",
        "Total",
        "COGS (est)",
        "Gross Profit (est)",
        "Marketing Spend (est)",
        "Operating Profit (est)",
    ]
    for r in range(2, ws.max_row + 1):
        for h in money:
            ws.cell(r, col[h]).number_format = "$#,##0"
        ws.cell(r, col["Assumed GM%"]).number_format = "0.0%"

    ws.freeze_panes = "A2"

    # Add a small summary block at bottom
    end_row = ws.max_row + 2
    ws.cell(end_row, 1, "Year totals")
    ws.cell(end_row, col["Total"], f"=SUM({ws.cell(2,col['Total']).coordinate}:{ws.cell(ws.max_row,col['Total']).coordinate})")
    ws.cell(end_row, col["Gross Profit (est)"], f"=SUM({ws.cell(2,col['Gross Profit (est)']).coordinate}:{ws.cell(ws.max_row,col['Gross Profit (est)']).coordinate})")
    ws.cell(end_row, col["Operating Profit (est)"], f"=SUM({ws.cell(2,col['Operating Profit (est)']).coordinate}:{ws.cell(ws.max_row,col['Operating Profit (est)']).coordinate})")
    for h in ["Total", "Gross Profit (est)"]:
        ws.cell(end_row, col[h]).number_format = "$#,##0"
    ws.cell(end_row, col["Operating Profit (est)"]).number_format = "$#,##0"

    wb.save(path)


def main():
    enhance_brand_guide()

    style_pptx(BASE / "marketing-pitch-deck.pptx", variant="Pitch Deck (v1)")
    style_pptx(BASE / "marketing-pitch-deck-final-final.pptx", variant="Pitch Deck (v2)")

    enhance_xlsx_quarterly_history(BASE / "quarterly-revenue-history-2024-2025.xlsx")
    enhance_xlsx_projection(BASE / "projected-revenue-2026.xlsx")


if __name__ == "__main__":
    main()
