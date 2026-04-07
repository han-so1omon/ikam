import json
import os
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_xlsx_quarterly_history():
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Quarterly Revenue"

    headers = [
        "Year",
        "Quarter",
        "Revenue (Bookkeeping)",
        "Revenue (POS)",
        "Notes",
    ]
    ws.append(headers)

    rows = [
        (2024, "Q2", 62000, 64000, "Soft launch mid-quarter"),
        (2024, "Q3", 88000, 90500, "Back-to-school gifting"),
        (2024, "Q4", 155000, 160200, "First big holiday season"),
        (2025, "Q1", 93000, 95050, "Rainy season foot traffic dip"),
        (2025, "Q2", 112000, 116300, "Local delivery pilot (limited)"),
        (2025, "Q3", 141000, 146900, "Subscription launch late quarter"),
        # Intentional mismatch for Q4 2025 (idea calls out contradictions)
        (2025, "Q4", 286000, 295000, "POS total excludes a week of manual returns"),
    ]
    for r in rows:
        ws.append(list(r))

    # Styling
    header_fill = PatternFill("solid", fgColor="2F4F4F")
    header_font = Font(color="FFFFFF", bold=True)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    ws.column_dimensions["A"].width = 10
    ws.column_dimensions["B"].width = 10
    ws.column_dimensions["C"].width = 22
    ws.column_dimensions["D"].width = 14
    ws.column_dimensions["E"].width = 48

    wb.save(BASE / "quarterly-revenue-history-2024-2025.xlsx")


def make_xlsx_projected_revenue_2026():
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "2026 Projection"

    ws.append(["Month", "In-store", "Online", "Subscriptions", "Corporate", "Total", "Assumed GM%", "Notes"])

    # Note: projection uses 48% GM (contradiction vs strategy target 52%)
    months = [
        ("2026-01", 62000, 11000, 8000, 6000, 0.48, "Post-holiday lull; new marketing contractor trial"),
        ("2026-02", 58000, 9000, 9000, 6500, 0.48, "Shipping footprint reduced mid-month"),
        ("2026-03", 65000, 9500, 10500, 8000, 0.48, "Spring refresh; vendor onboarding"),
        ("2026-04", 70000, 10000, 11500, 9000, 0.48, "Local corporate gifting push"),
        ("2026-05", 72000, 10500, 12000, 8500, 0.48, "Mother’s Day bundles"),
        ("2026-06", 76000, 11000, 12500, 9500, 0.48, "Summer traffic"),
        ("2026-07", 74000, 12000, 13000, 10000, 0.48, "Subscription refresh rollout"),
        ("2026-08", 78000, 12500, 13500, 11000, 0.48, "Back-to-school gifting"),
        ("2026-09", 80000, 13000, 14000, 12000, 0.48, "Prep for holiday ramp"),
        ("2026-10", 90000, 15000, 15000, 16000, 0.48, "Holiday pre-orders"),
        ("2026-11", 120000, 18000, 17000, 28000, 0.48, "Holiday staffing required"),
        ("2026-12", 170000, 22000, 18000, 45000, 0.48, "Peak season"),
    ]

    for m, instore, online, subs, corp, gm, notes in months:
        total = instore + online + subs + corp
        ws.append([m, instore, online, subs, corp, total, gm, notes])

    # Styling
    header_fill = PatternFill("solid", fgColor="1F4E79")
    header_font = Font(color="FFFFFF", bold=True)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    widths = [12, 12, 12, 14, 12, 12, 12, 44]
    for i, w_ in enumerate(widths, start=1):
        ws.column_dimensions[chr(64 + i)].width = w_

    wb.save(BASE / "projected-revenue-2026.xlsx")


def make_xlsx_inventory_plan_spring_2026():
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Spring 2026"

    ws.append(["SKU", "Item name", "Vendor", "Category", "Reorder point", "Target stock", "Unit cost", "Retail price", "Notes"])
    items = [
        ("JAM-ALM-08", "Almond Apricot Jam (8oz)", "Sunfield Preserves", "Pantry", 12, 36, 5.10, 12.00, "Top seller; label changed (old: 'Apricot-Almond')"),
        ("SCE-HOT-05", "Smoked Chili Hot Sauce", "Oak Ember Foods", "Pantry", 10, 24, 4.40, 11.00, "Seasonal spike"),
        ("CHO-SEA-01", "Sea Salt Dark Chocolate", "Nightjar Chocolate", "Sweets", 20, 60, 2.20, 6.50, "Bundle-friendly"),
        ("TEA-JSM-10", "Jasmine Pearl Tea (tin)", "Juniper & Co.", "Beverage", 6, 18, 7.80, 18.00, "Slow mover but high margin"),
        ("GFT-CORP-01", "Corporate Mini Gift Pack", "(Assembled)", "Gifts", 8, 25, 14.50, 35.00, "COGS varies by contents"),
    ]
    for row in items:
        ws.append(list(row))

    header_fill = PatternFill("solid", fgColor="3C7D22")
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = Font(color="FFFFFF", bold=True)

    wb.save(BASE / "inventory-plan-spring-2026.xlsx")


def make_pptx_deck(path: str, *, q4_2025_revenue: str, subs_name: str, gm_target: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(20)

    add_title_slide(
        "Bramble & Bitters",
        "Local specialty retail • curated gifts • corporate bundles",
    )

    add_bullets(
        "What we are",
        [
            "Neighborhood shop in Oakland + limited ship-to",
            f"Subscription program: {subs_name}",
            "Corporate gifting bundles for local companies",
        ],
    )

    add_bullets(
        "Traction (selected)",
        [
            "Opened storefront: Mar 2024",
            "First corporate order: Aug 2024",
            f"Q4 2025 revenue: {q4_2025_revenue}",
        ],
    )

    add_bullets(
        "Strategy (2026)",
        [
            "Increase repeat purchase via seasonal curation",
            "Improve corporate gifting lead time + packaging consistency",
            f"Target gross margin: {gm_target}",
        ],
    )

    add_bullets(
        "Asks",
        [
            "Local corporate introductions",
            "Press / partnerships with neighborhood events",
            "Packaging vendor recommendations",
        ],
    )

    prs.save(BASE / path)


def make_markdown_docs():
    w(
        "mission-vision-values.md",
        """# Bramble & Bitters — Mission, Vision, Values

## Mission
Make gifting (and everyday treats) feel personal again by curating small-batch goods and telling the makers’ stories.

## Vision
Be the East Bay’s go-to neighborhood shop for thoughtful gifts and seasonal pantry discoveries — with a small, sustainable shipping footprint.

## Values
- **Curate with taste:** fewer items, better picks.
- **Make it human:** warm service, real notes, remember regulars.
- **Support makers:** pay on time, spotlight local talent.
- **Keep it practical:** margins matter; seasonal planning beats constant firefighting.
- **Leave it better:** low-waste packaging when feasible.
""",
    )

    w(
        "brand-guide.md",
        """# Bramble & Bitters — Brand Guide (light)

## Voice + tone
- Warm, specific, a little poetic — but never precious.
- Write like a helpful shopkeeper, not a lifestyle brand.
- Avoid buzzwords. Prefer sensory details.

## Messaging themes
- Seasonal curation (what’s new, why it’s good)
- Maker stories (short, factual, respectful)
- Giftability (make it easy to choose)

## Visual notes (intentionally non-precise)
- Primary colors: deep green, cream, copper
- Type feel: classic serif for headers, clean sans for body
- Photography: natural light; hands holding products; shelves, not staged studio

## Naming conventions
- Subscription program name: varies (Bramble Box / Seasonal Pantry Box / Quarterly Crate)
- Initiatives often have nicknames (e.g., “Box Revamp”)
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Maya Chen
Last updated: 2026-02-03

## Situation
We grew quickly in holiday 2025, but fulfillment (shipping + returns) introduced noise into revenue reporting and strained the team.

## Strategy pillars
1) **Repeat purchase over constant acquisition**
   - Seasonal curation that brings regulars back.
2) **Corporate gifting as a stabilizer**
   - Standard bundles, clear lead times, smoother intake.
3) **Subscription refresh**
   - Simplify the box, improve naming + benefits.
4) **Operational calm**
   - Better vendor lead times, inventory reorder points, and a real holiday staffing plan.

## Targets (directional)
- Target gross margin: **52%** (mix shift + packaging improvements)
- Reduce shipping footprint to protect service quality.

## Risks
- Inventory cash tied up in slow movers.
- Over-customizing corporate orders.
- Marketing cadence inconsistent if contractor churns.
""",
    )

    w(
        "annual-goals-okrs-2026.md",
        """# Annual Goals + OKRs — 2026

## Objective 1: Make corporate gifting reliable and repeatable
- KR1: Launch a standardized intake flow (form + confirmation email)
- KR2: Reduce average corporate lead time from 10 days to 6 days
- KR3: 20 repeat corporate customers by end of year

## Objective 2: Improve subscription retention
- KR1: Rename + repackage subscription (one canonical name, keep aliases in old posts)
- KR2: Reduce churn from ~35% to < 25% by Q4
- KR3: 150 active subscribers by Q4

## Objective 3: Calm operations
- KR1: Inventory plan in place for top 30 SKUs
- KR2: Holiday staffing plan written by Sep 1
- KR3: Monthly close process documented (even if simple)
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

## Pillar 1: Seasonal curation
We do the browsing so you don’t have to.

## Pillar 2: Local makers, real stories
Short maker bios; highlight ingredients and process.

## Pillar 3: Gifts that don’t feel generic
Bundles that look intentional — not like a last-minute scramble.

## Proof points
- In-store tasting notes
- Rotating shelves + small drops
- Corporate gifting with clear lead times
""",
    )

    w(
        "customer-personas.md",
        """# Customer Personas

## Persona A: The Neighbor Regular
- Wants: a small treat, something new, friendly chat
- Trigger: seasonal shelf refresh

## Persona B: The Gift Buyer
- Wants: “I need a gift that looks thoughtful, today.”
- Trigger: birthdays, dinners, holidays

## Persona C: The Office Hero
- Wants: employee/client gifts that won’t cause drama
- Needs: invoice, shipping clarity, lead time, consistency
""",
    )

    w(
        "roadmap-2026-h1.md",
        """# Roadmap — 2026 H1 (draft)

Owner: Maya
Status: draft / partially updated after shipping pivot

## Feb–Mar
- Corporate intake form + standardized bundles (Jordan + Maya)
- Subscription refresh: packaging + naming decision (Maya)
- Vendor list cleanup (Talia supports copy)

## Apr–May
- Spring inventory plan execution
- Corporate outreach: top 30 local companies
- Improve signage + in-store “gift helper” corner

## Jun
- Prep holiday 2026 playbook outline (start early)
- Local events partnerships (2–3)

## Notes
- Shipping footprint reduced Feb 2026 (not fully reflected in older marketing docs).
""",
    )

    w(
        "prior-development-cycles.md",
        """# Prior “Development Cycles” (how we work)

We don’t run formal engineering sprints. But we do run **seasonal cycles**.

## Holiday 2025 cycle (example)
- **Plan (early Oct):** decide bundles, packaging, vendors
- **Build (late Oct–Nov):** stock up, write product copy, assemble bundles
- **Run (Nov–Dec):** daily adjustments, restocks, signage, corporate orders
- **Retro (Jan):** what sold, what didn’t, what broke (shipping/returns)

## Artifacts created during cycles
- A single “bundle list” spreadsheet
- Vendor email threads
- A few meeting notes
- Instagram draft captions
""",
    )

    w(
        "meeting-notes-2025-10-15.md",
        """# Meeting Notes — 2025-10-15 (Vendor planning)

Attendees: Maya, Jordan

## Decisions
- Prioritize 3 bundle tiers (mini / classic / deluxe)
- Keep local maker ratio high (feels true to the shop)

## Open questions
- Packaging vendor lead times
- Whether to do nationwide shipping for corporate (tentative yes)

## Notes
- Jordan flagged that returns policy is confusing at checkout.
""",
    )

    w(
        "meeting-notes-2025-11-20.md",
        """# Meeting Notes — 2025-11-20 (Holiday staffing)

Attendees: Maya, Jordan

## Staffing
- Hire 2 seasonal associates (names in notes vary: “Sam”/“Samuel”, “Nina”)
- Extend store hours weekends in December

## Risks
- Corporate orders stacking up with custom requests
- Shipping days conflicting with in-store rush

## Action items
- Draft a simple corporate intake email
- Pre-build 40 mini bundles by Dec 1
""",
    )

    w(
        "voice-note-transcript-2026-01-07.md",
        """# Voice Note Transcript — 2026-01-07 (Maya)

> Okay so… Q4 was huge and also kind of a mess. The shop did great but the shipping thing was… painful.

- I think the *real* number for Q4 is somewhere around **$300k** but it depends on what you count.
- Square shows one total, but we had that week where returns were handled manually because the workflow was broken and we were slammed.
- Corporate is the stabilizer. Less chaos if we standardize bundles.
- Subscription needs a single name. People keep calling it three different things.
- Jordan is basically running the floor now. I should make that official.

Also: we can’t ship everywhere. It makes us worse at everything.
""",
    )

    w(
        "vendor-evaluation-template.md",
        """# Vendor Evaluation Template

Vendor name:

## Scores (1–5)
- Product quality:
- Reliability / lead time:
- Margin friendliness:
- Packaging + labeling consistency:
- Responsiveness:

## Notes
- What went well:
- What went wrong:
- Any label/name variants to track:

## Decision
- [ ] Keep
- [ ] Keep (reduce volume)
- [ ] Trial
- [ ] Replace
""",
    )

    w(
        "campaign-postmortem-template.md",
        """# Campaign Postmortem Template

Campaign:
Dates:
Owner:

## Goals

## What happened (timeline)

## Results
- Reach:
- Clicks:
- Sales attributed (if any):

## What worked

## What didn’t

## Follow-ups

## Files / links
""",
    )


def make_json_artifacts():
    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "currency": "USD",
            "kpis": [
                {
                    "name": "Revenue (POS)",
                    "description": "Gross sales captured by Square POS exports.",
                    "formula": "sum(transactions.total_amount) - sum(transactions.refunds_processed_in_pos)",
                    "notes": "May exclude manual returns/adjustments during peak weeks.",
                },
                {
                    "name": "Revenue (Bookkeeping)",
                    "description": "Revenue recognized in quarterly close by North Ledger LLC.",
                    "formula": "revenue_accounts_total_for_period",
                    "notes": "May include manual adjustments and accrual timing.",
                },
                {
                    "name": "Gross Margin %",
                    "description": "Blended gross margin across channels.",
                    "formula": "(revenue - cogs) / revenue",
                    "notes": "Target stated as 52% in strategy doc; projections use 48% assumption.",
                },
                {
                    "name": "AOV",
                    "description": "Average order value (POS).",
                    "formula": "Revenue (POS) / count(orders)",
                },
            ],
        },
    )

    wjson(
        "vendor-list.json",
        {
            "generated": str(date.today()),
            "vendors": [
                {
                    "vendor_id": "V001",
                    "name": "Sunfield Preserves",
                    "aka": ["Sunfield", "Sunfield Preserves LLC"],
                    "category": "jams",
                    "risk_notes": "Seasonal supply constraints; labels changed in 2025.",
                },
                {
                    "vendor_id": "V002",
                    "name": "Oak Ember Foods",
                    "aka": ["OakEmber"],
                    "category": "hot sauce",
                    "risk_notes": "Lead times unpredictable near holidays.",
                },
                {
                    "vendor_id": "V003",
                    "name": "Nightjar Chocolate",
                    "aka": ["Nightjar", "Nightjar Chocolate Co."],
                    "category": "chocolate",
                    "risk_notes": "Reliable; packaging occasionally arrives scuffed.",
                },
                {
                    "vendor_id": "V004",
                    "name": "Juniper & Co.",
                    "aka": ["Juniper and Co", "Juniper"],
                    "category": "tea",
                    "risk_notes": "Slow mover but high margin; watch cash tied up.",
                },
                {
                    "vendor_id": "V005",
                    "name": "PaperMoss Packaging",
                    "aka": ["Paper Moss"],
                    "category": "packaging",
                    "risk_notes": "Switched to this vendor mid-2025 after cost increase with prior vendor.",
                },
            ],
        },
    )

    wjson(
        "corporate-order-intake-form.json",
        {
            "schema_version": "0.2",
            "title": "Corporate Gift Order Intake",
            "fields": [
                {"id": "company_name", "type": "string", "required": True},
                {"id": "contact_name", "type": "string", "required": True},
                {"id": "contact_email", "type": "string", "required": True},
                {"id": "invoice_required", "type": "boolean", "required": False, "default": True},
                {"id": "desired_delivery_date", "type": "date", "required": True},
                {
                    "id": "bundle_tier",
                    "type": "enum",
                    "required": True,
                    "values": ["mini", "classic", "deluxe"],
                },
                {"id": "quantity", "type": "integer", "required": True, "min": 1},
                {"id": "shipping", "type": "enum", "required": True, "values": ["pickup", "local_delivery", "ship"], "notes": "Shipping currently limited (CA/nearby)."},
                {"id": "addresses", "type": "array", "required": False, "items": {"type": "object", "fields": [{"id": "name"}, {"id": "address"}]}},
                {"id": "custom_note", "type": "string", "required": False},
                {"id": "budget_per_bundle", "type": "number", "required": False},
            ],
        },
    )

    # POS export, intentionally "raw-ish" and imperfect
    wjson(
        "pos-export-2025-q4.json",
        {
            "source": "Square",
            "exported_at": "2026-01-03T09:12:00-08:00",
            "period": {"start": "2025-10-01", "end": "2025-12-31"},
            "notes": [
                "Export excludes some manual returns processed off-register during 2025-12-14 to 2025-12-20.",
                "Tips not included.",
            ],
            "summary": {
                "gross_sales": 312450,
                "discounts": 7800,
                "refunds_processed_in_pos": 9650,
                "net_sales": 295000,
            },
            "transactions_sample": [
                {"id": "T-10001", "date": "2025-10-03", "channel": "in_store", "total_amount": 64.50, "items": ["Sea Salt Dark Chocolate", "Smoked Chili Hot Sauce"]},
                {"id": "T-10888", "date": "2025-11-18", "channel": "in_store", "total_amount": 42.00, "items": ["Almond Apricot Jam (8oz)"]},
                {"id": "T-12002", "date": "2025-12-05", "channel": "online", "total_amount": 118.00, "items": ["Classic Gift Bundle"], "ship_to": "San Jose, CA"},
                {"id": "T-12990", "date": "2025-12-19", "channel": "in_store", "total_amount": 0.00, "items": [], "note": "Peak week: some returns handled manually, not in POS"},
            ],
        },
    )


def make_pdfs_with_reportlab():
    """Generate lightweight PDFs without requiring LaTeX."""
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def build_pdf(out_path: Path, title: str, blocks: list):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = []
        story.append(Paragraph(f"<b>{title}</b>", styles["Title"]))
        story.append(Spacer(1, 12))
        for b in blocks:
            kind = b.get("kind")
            if kind == "p":
                story.append(Paragraph(b["text"], styles["BodyText"]))
                story.append(Spacer(1, 10))
            elif kind == "h":
                story.append(Paragraph(f"<b>{b['text']}</b>", styles["Heading2"]))
                story.append(Spacer(1, 6))
            elif kind == "ul":
                items = [ListItem(Paragraph(t, styles["BodyText"])) for t in b["items"]]
                story.append(ListFlowable(items, bulletType="bullet"))
                story.append(Spacer(1, 10))
        doc.build(story)

    # Corporate one-pager
    build_pdf(
        BASE / "corporate-gifting-one-pager.pdf",
        "Bramble & Bitters — Corporate Gifting (One Pager)",
        [
            {"kind": "p", "text": "<b>Who we are</b><br/>A neighborhood specialty shop in Oakland curating small-batch pantry goods and giftable items."},
            {"kind": "h", "text": "Bundles"},
            {"kind": "ul", "items": [
                "<b>Mini</b> ($35) — small thank-you gift",
                "<b>Classic</b> ($55) — best for teams",
                "<b>Deluxe</b> ($85) — client-facing gifts",
            ]},
            {"kind": "h", "text": "Lead times"},
            {"kind": "p", "text": "Typical: <b>6–10 business days</b> (depends on quantity + customization)"},
            {"kind": "h", "text": "Logistics"},
            {"kind": "ul", "items": [
                "Pickup in Oakland",
                "Local delivery (East Bay)",
                "Shipping: <b>limited</b> (CA + nearby)",
            ]},
            {"kind": "h", "text": "Contact"},
            {"kind": "p", "text": "Maya Chen — Owner/GM<br/>hello@brambleandbitters.example"},
        ],
    )

    # Bookkeeping summary
    build_pdf(
        BASE / "bookkeeping-q4-2025-summary.pdf",
        "North Ledger LLC — Q4 2025 Summary (Bramble & Bitters)",
        [
            {"kind": "p", "text": "Prepared for: Maya Chen<br/>Prepared by: North Ledger LLC<br/>Period: 2025-10-01 to 2025-12-31"},
            {"kind": "h", "text": "Highlights"},
            {"kind": "ul", "items": [
                "Strong holiday performance driven by bundles + corporate gifting.",
                "Noted operational strain: fulfillment + returns workflow issues in mid-December.",
            ]},
            {"kind": "h", "text": "Revenue (bookkeeping basis)"},
            {"kind": "p", "text": "Q4 2025 revenue: <b>$286,000</b>"},
            {"kind": "h", "text": "Notes on reconciliation"},
            {"kind": "ul", "items": [
                "POS exports may show different totals due to timing and manual return adjustments.",
                "Recommend documenting a simple returns log for peak weeks.",
            ]},
            {"kind": "h", "text": "Recommendations"},
            {"kind": "ul", "items": [
                "Standardize corporate bundle SKUs.",
                "Document a month-end close checklist (even if minimal).",
            ]},
        ],
    )


def main():
    make_markdown_docs()
    make_json_artifacts()
    make_xlsx_quarterly_history()
    make_xlsx_projected_revenue_2026()
    make_xlsx_inventory_plan_spring_2026()

    # PPTX decks: v1 and 'final-final' with intentional contradictions
    make_pptx_deck(
        "marketing-pitch-deck.pptx",
        q4_2025_revenue="$295k (POS net)",
        subs_name="Seasonal Pantry Box",
        gm_target="52%",
    )
    make_pptx_deck(
        "marketing-pitch-deck-final-final.pptx",
        q4_2025_revenue="$310k (reported)",
        subs_name="Bramble Box",
        gm_target="50%+",
    )

    make_pdfs_with_reportlab()

    # Cleanliness: leave sources for auditability


if __name__ == "__main__":
    main()
