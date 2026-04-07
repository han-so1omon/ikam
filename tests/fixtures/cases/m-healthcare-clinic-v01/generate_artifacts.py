from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Mesa Ridge Urgent & Primary Care — Mission, Vision, Values

## Mission
Provide accessible urgent and primary care with predictable operations and clear communication.

## Vision
A small clinic group known for consistency and patient trust.

## Values
- Plain-language care
- Privacy by default
- Fix the system
""",
    )

    w(
        "brand-guide.md",
        """# Mesa Ridge — Brand Guide (v0.1)

Palette:
- Blue: #1D4ED8
- Teal: #14B8A6
- White: #F8FAFC

Voice:
- Calm, reassuring, direct.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Ryan Patel
Last updated: 2026-02-01

Priorities:
- Align no-show definition
- Reduce denials
- Improve training completion tracking
""",
    )

    w(
        "clinic-ops-handbook-v1.md",
        """# Clinic Ops Handbook (v1)

- Check-in process
- No-show handling (late cancels excluded)
- Staffing ratios guidance
""",
    )

    w(
        "billing-process-map.md",
        """# Billing Process Map

1) Eligibility
2) Coding
3) Claim submission
4) Denials
5) Patient statements
""",
    )

    w(
        "privacy-policy.md",
        """# Privacy Policy

- Annual HIPAA training required
- Minimum necessary access
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Same-day urgent care
2) Friendly primary care
3) Clear billing expectations
""",
    )

    w(
        "ops-call-2026-01-15.md",
        """# Ops Call — 2026-01-15

- No-show definition differs by site.
- Staffing sheet is maintained separately by Ryan.
""",
    )

    w(
        "ops-call-2026-02-12.md",
        """# Ops Call — 2026-02-12

- Collections numbers differ between billing tracker and finance summary.
- Training completion is behind for front desk.
""",
    )

    w(
        "voice-note-transcript-2026-02-16.md",
        """# Voice Note Transcript — 2026-02-16

- Ops no-show excludes late cancels; sites include them.
- Denials mostly eligibility and coding.
- Finance includes catch-up payments; billing doesn't.
""",
    )

    w(
        "patient-complaint-template.md",
        """# Patient Complaint Template

Date:
Site:
Complaint:
Resolution:
Owner:
""",
    )

    w(
        "incident-report-template.md",
        """# Incident Report Template

Date:
Site:
Incident:
Severity:
Actions:
Owner:
""",
    )

    w(
        "new-hire-onboarding-checklist.md",
        """# New Hire Onboarding Checklist

- HIPAA training
- Systems access
- Shadowing
- Policy acknowledgement
""",
    )

    w(
        "performance-review-template.md",
        """# Performance Review Template

Role:
Period:

Outcomes:
Strengths:
Growth:
""",
    )


def make_json():
    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "No-show (ops)", "definition": "excludes late cancels"},
                {"name": "No-show (sites)", "definition": "includes late cancels <24h"},
                {"name": "Collections (finance)", "definition": "incl catch-up"},
                {"name": "Collections (billing)", "definition": "excl catch-up"},
            ],
        },
    )

    wjson(
        "patient-intake-form.json",
        {
            "schema_version": "0.2",
            "title": "Patient Intake",
            "fields": [
                {"id": "full_name", "type": "string", "required": True},
                {"id": "dob", "type": "date", "required": True},
                {"id": "insurance", "type": "string"},
                {"id": "reason", "type": "string"},
            ],
            "notes": "Synthetic only; avoid PHI in benchmark fills.",
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "employee-handbook-excerpt.pdf",
        "Employee Handbook Excerpt (summary)",
        ["Privacy expectations", "Attendance", "Training requirements"],
    )

    pdf(
        BASE / "service-line-one-pager.pdf",
        "Service Lines — Overview",
        ["Urgent care", "Primary care", "Employer contracts"],
    )


def make_xlsx():
    import openpyxl

    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Quarterly"
    ws.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q4",2400000,""),(2025,"Q1",2550000,""),(2025,"Q2",2620000,""),(2025,"Q3",2700000,"denials spike"),(2025,"Q4",2850000,"")]: ws.append(list(r))
    for i in range(2, ws.max_row+1): ws.cell(i,3).number_format="$#,##0"
    wb.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="2026"
    ws2.append(["Month","Revenue","Notes"])
    for r in [("2026-01",920000,""),("2026-02",880000,""),("2026-03",940000,"")]: ws2.append(list(r))
    for i in range(2, ws2.max_row+1): ws2.cell(i,2).number_format="$#,##0"
    wb2.save(BASE/"projected-revenue-2026.xlsx")

    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Model"
    ws3.append(["Role","Per clinic FTE","Clinics","Total FTE"])
    for r in [("Front desk",2.0,3,None),("MA",2.5,3,None),("Provider",2.0,3,None)]: ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,4).value=f"=B{i}*C{i}"
    wb3.save(BASE/"staffing-model.xlsx")

    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="2026-01"
    ws4.append(["Site","No-show (ops)","No-show (site)","Notes"])
    for r in [("Clinic A",0.071,0.105,"late cancels"),("Clinic B",0.079,0.11,"")]: ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="0.0%"; ws4.cell(i,3).number_format="0.0%"
    wb4.save(BASE/"site-scorecard-2026-01.xlsx")

    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Q1 2026"
    ws5.append(["Date","Site","Incident","Severity","Notes"])
    for r in [("2026-01-22","Clinic C","Reminder calls missed","Med",""),("2026-02-05","Billing","Eligibility denials spike","High","")]: ws5.append(list(r))
    wb5.save(BASE/"incident-log-2026-q1.xlsx")

    wb6=openpyxl.Workbook(); ws6=wb6.active; ws6.title="2026-01"
    ws6.append(["Payer","Denial rate","Top reason","Notes"])
    for r in [("Payer A",0.075,"eligibility",""),("Payer B",0.062,"coding","")]: ws6.append(list(r))
    for i in range(2, ws6.max_row+1): ws6.cell(i,2).number_format="0.0%"
    wb6.save(BASE/"denials-tracker-2026-01.xlsx")

    wb7=openpyxl.Workbook(); ws7=wb7.active; ws7.title="2026-01"
    ws7.append(["Metric","Value","Definition"])
    ws7.append(["Collections (finance)",0.92,"incl catch-up"])
    ws7.append(["Collections (billing)",0.89,"excl catch-up"])
    ws7.cell(2,2).number_format="0.0%"; ws7.cell(3,2).number_format="0.0%"
    wb7.save(BASE/"collections-summary-2026-01.xlsx")

    wb8=openpyxl.Workbook(); ws8=wb8.active; ws8.title="2026"
    ws8.append(["Dept","Required","Completed","Completion %","Notes"])
    for r in [("Clinical",60,52,None,""),("Front desk",30,21,None,"behind")]: ws8.append(list(r))
    for i in range(2, ws8.max_row+1):
        ws8.cell(i,4).value=f"=C{i}/B{i}"
        ws8.cell(i,4).number_format="0.0%"
    wb8.save(BASE/"hipaa-training-tracker-2026.xlsx")

    wb9=openpyxl.Workbook(); ws9=wb9.active; ws9.title="2026"
    ws9.append(["Dept","Role","Count","Priority","Notes"])
    for r in [("Clinics","MA",3,"High",""),("Billing","Denials specialist",1,"Med","")]: ws9.append(list(r))
    wb9.save(BASE/"hiring-plan-2026.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    BLUE=RGBColor(0x1D,0x4E,0xD8)
    TEAL=RGBColor(0x14,0xB8,0xA6)
    OFF=RGBColor(0xF8,0xFA,0xFC)

    prs=Presentation()
    def header(slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=BLUE; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=TEAL; acc.line.fill.background()

    s=prs.slides.add_slide(prs.slide_layouts[6]); header(s)
    tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
    p=tb.text_frame.paragraphs[0]; p.text="Mesa Ridge Care"; p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
    prs.save(BASE/"marketing-pitch-deck.pptx")


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
