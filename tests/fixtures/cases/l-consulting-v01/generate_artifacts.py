from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# Greybridge Partners — Mission, Vision, Values

## Mission
Help leadership teams make hard decisions quickly: stabilize operations, fix costs, and execute.

## Vision
Be the first-call firm for mid-market turnarounds and operator-led transformations.

## Values
- Work from facts
- Don’t hide tradeoffs
- Build artifacts that survive
- Own outcomes
""",
    )

    w(
        "brand-guide.md",
        """# Greybridge Partners — Brand Guide (v0.1)

## Palette
- Charcoal: #1F2937
- Off-white: #F7F5F0
- Muted gold: #B89B5E
- Slate: #64748B

## Voice
- Calm, analytical, direct.
- Prefer short sentences and quantified claims.

## Slide style
- Charcoal header bar, gold accent line
- Simple diagrams; minimal stock photos
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Marisol Vega
Last updated: 2026-02-01

## Focus
- Turnaround practice growth
- Reduce internal template sprawl
- Align utilization definitions across practices
""",
    )

    w(
        "practices-overview.md",
        """# Practices Overview

- Turnaround / restructuring
- Ops excellence
- Pricing + cost transformation
- PMO

Note: practices overlap; naming varies in internal materials.
""",
    )

    w(
        "proposal-template.md",
        """# Proposal Template

## Client context

## Hypothesis

## Scope

## Approach

## Team

## Fees

## Timeline
""",
    )

    w(
        "messaging-pillars.md",
        """# Messaging Pillars

1) Stabilize first
2) Fix the cost structure
3) Build the operating cadence
4) Deliver artifacts the client can run
""",
    )

    # Sample engagement
    w(
        "sample-engagement-IRONCLAD/kickoff-notes.md",
        """# Kickoff Notes — IRONCLAD

Client: (alias) Ironclad Industries

## Goals
- Stabilize cash
- Reduce expedited freight
- Build a 13-week cash forecast
""",
    )

    w(
        "sample-engagement-IRONCLAD/weekly-status-2026-01-12.md",
        """# Weekly Status — 2026-01-12 (IRONCLAD)

- Data access improving
- Benefits tracker v1 started
""",
    )

    w(
        "sample-engagement-IRONCLAD/weekly-status-2026-01-26.md",
        """# Weekly Status — 2026-01-26 (IRONCLAD)

- Proposal scope expanded (quietly)
- Partner says deal is verbally won; CRM still shows Proposal stage
""",
    )

    w(
        "sample-engagement-IRONCLAD/closeout-summary.md",
        """# Closeout Summary — IRONCLAD

Status: partial

- Cost actions identified; execution in progress.
- Success fee discussion unresolved.
""",
    )

    w(
        "partner-meeting-notes-2026-02-03.md",
        """# Partner Meeting Notes — 2026-02-03

- Utilization debate: practice dashboard says 78%, finance says 71%.
- Pipeline truth: some deals are delivering before "Closed Won".
""",
    )

    w(
        "practice-leads-notes-2026-02-10.md",
        """# Practice Leads Notes — 2026-02-10

- Turnaround practice wants success-fee stories.
- Finance is pushing back on booked amounts.
""",
    )

    w(
        "voice-note-transcript-2026-02-18.md",
        """# Voice Note Transcript — 2026-02-18 (Nina)

- Marketing keeps claiming $1.2M success fee. Finance can only book $0.9M right now.
- Utilization depends on what we count as billable.
""",
    )

    w(
        "engagement-retro-template.md",
        """# Engagement Retro Template

## Outcomes

## What worked

## What didn't

## Artifacts
- SOW
- Deliverable deck
- Benefits tracker
""",
    )

    w(
        "consultant-performance-rubric.md",
        """# Consultant Performance Rubric

Categories (1–5):
- Client trust
- Analytical rigor
- Artifact quality
- Ownership
- Communication
""",
    )


def make_json():
    wjson(
        "lead-list.json",
        {
            "generated": str(date.today()),
            "leads": [
                {"account": "Ironclad Industries", "codename": "IRONCLAD", "stage": "Proposal", "amount": 650000, "notes": "Partner says verbally won"},
                {"account": "Lakefront Packaging", "codename": "FOIL", "stage": "Discovery", "amount": 420000},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-01",
            "kpis": [
                {"name": "Utilization (practice)", "definition": "billable + approved internal / capacity"},
                {"name": "Utilization (finance)", "definition": "billable only / capacity"},
                {"name": "Success fee (marketing)", "definition": "claimed total benefit capture"},
                {"name": "Success fee (finance)", "definition": "booked amount net of contingencies"},
            ],
        },
    )

    wjson(
        "client-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Client Intake",
            "fields": [
                {"id": "legal_name", "type": "string", "required": True},
                {"id": "alias", "type": "string", "required": False},
                {"id": "codename", "type": "string", "required": True},
                {"id": "problem", "type": "string", "required": True},
                {"id": "success_metrics", "type": "array", "items": {"type": "string"}},
            ],
        },
    )


def make_pdfs():
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    styles = getSampleStyleSheet()

    def pdf(out_path: Path, title: str, bullets: list[str]):
        doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, title=title)
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        items = [ListItem(Paragraph(b, styles["BodyText"])) for b in bullets]
        story.append(ListFlowable(items, bulletType="bullet"))
        doc.build(story)

    pdf(
        BASE / "case-studies-portfolio.pdf",
        "Case Studies Portfolio (selected)",
        [
            "IRONCLAD — turnaround plan + cash forecast",
            "Claimed success fee: $1.2M (marketing)",
            "Booked success fee: $0.9M (finance)",
        ],
    )

    pdf(
        BASE / "sample-engagement-IRONCLAD/statement-of-work.pdf",
        "Statement of Work — IRONCLAD",
        [
            "Term: 10 weeks",
            "Fees: $650k + success fee contingent",
            "Deliverables: turnaround plan, benefits tracker, 13-week cash forecast",
        ],
    )


def make_xlsx():
    import openpyxl

    # Revenue history
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Quarterly"
    ws.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q3",4800000,""),(2024,"Q4",5200000,""),(2025,"Q1",5000000,""),(2025,"Q2",5400000,""),(2025,"Q3",5600000,""),(2025,"Q4",6100000,"Success-fee case")]:
        ws.append(list(r))
    for i in range(2, ws.max_row+1):
        ws.cell(i,3).number_format = "$#,##0"
    wb.save(BASE / "quarterly-revenue-history-2024-2025.xlsx")

    # Projection
    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="2026"
    ws2.append(["Month","Revenue","Notes"])
    for r in [("2026-01",1600000,""),("2026-02",1550000,""),("2026-03",1750000,"" )]:
        ws2.append(list(r))
    for i in range(2, ws2.max_row+1):
        ws2.cell(i,2).number_format="$#,##0"
    wb2.save(BASE / "projected-revenue-2026.xlsx")

    # Pipeline
    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Pipeline Q1 2026"
    ws3.append(["Account","Codename","Stage","Amount","Probability","Weighted"])
    rows=[("Ironclad Industries","IRONCLAD","Proposal",650000,0.6,None),("Lakefront Packaging","FOIL","Discovery",420000,0.35,None)]
    for r in rows:
        ws3.append(list(r))
    for i in range(2, ws3.max_row+1):
        ws3.cell(i,6).value=f"=D{i}*E{i}"
        ws3.cell(i,4).number_format="$#,##0"; ws3.cell(i,6).number_format="$#,##0"; ws3.cell(i,5).number_format="0%"
    wb3.save(BASE / "pipeline-report-2026-q1.xlsx")

    # Utilization tracker contradiction
    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="Utilization 2025"
    ws4.append(["Month","Utilization (practice)","Utilization (finance)","Notes"])
    for r in [("2025-10",0.77,0.70,""),("2025-11",0.78,0.71,"definition dispute"),("2025-12",0.76,0.69,"")]:
        ws4.append(list(r))
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,2).number_format="0%"; ws4.cell(i,3).number_format="0%"
    wb4.save(BASE / "utilization-tracker-2025.xlsx")

    # Benefits tracker
    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Benefits"
    ws5.append(["Initiative","Owner","Estimated Annual Benefit","Status","Notes"])
    rows=[("Reduce expedited freight","Ops","$450,000","In progress",""),("Inventory cleanup","Finance","$300,000","Planned",""),("Vendor renegotiation","Procurement","$450,000","In progress","Marketing claims full benefit")]
    for r in rows:
        ws5.append(list(r))
    wb5.save(BASE / "sample-engagement-IRONCLAD/benefits-tracker.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    CHAR = RGBColor(0x1F,0x29,0x37)
    OFF = RGBColor(0xF7,0xF5,0xF0)
    GOLD = RGBColor(0xB8,0x9B,0x5E)

    def header(prs, slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb = OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = CHAR; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb = GOLD; acc.line.fill.background()

    def deck(path: Path, title: str, subtitle: str, slides: list[tuple[str,list[str]]]):
        prs=Presentation()
        s=prs.slides.add_slide(prs.slide_layouts[6]); header(prs,s)
        tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
        p=tb.text_frame.paragraphs[0]; p.text=title
        p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True; p.font.color.rgb=CHAR
        sb=s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11), Inches(0.7))
        p2=sb.text_frame.paragraphs[0]; p2.text=subtitle
        p2.font.name="Arial"; p2.font.size=Pt(20); p2.font.color.rgb=CHAR

        for t, items in slides:
            s2=prs.slides.add_slide(prs.slide_layouts[6]); header(prs,s2)
            tt=s2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
            pp=tt.text_frame.paragraphs[0]; pp.text=t
            pp.font.name="Georgia"; pp.font.size=Pt(34); pp.font.bold=True; pp.font.color.rgb=CHAR
            box=s2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(4.2))
            tf=box.text_frame; tf.clear()
            for i,it in enumerate(items):
                pr=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                pr.text=it
                pr.font.name="Arial"; pr.font.size=Pt(20); pr.font.color.rgb=CHAR

        prs.save(path)

    deck(
        BASE/"marketing-pitch-deck.pptx",
        "Greybridge Partners",
        "Turnaround + ops transformation",
        [("What we do",["Stabilize", "Fix cost structure", "Build cadence"]), ("Proof",["IRONCLAD (case)", "Benefits tracker", "Turnaround plan"])],
    )

    deck(
        BASE/"sample-engagement-IRONCLAD/deliverable-turnaround-plan-v1.pptx",
        "IRONCLAD Turnaround Plan",
        "v1",
        [("Situation",["Cash pressure", "Expedite freight"]), ("Plan",["13-week cash forecast", "Cost actions"])],
    )

    deck(
        BASE/"sample-engagement-IRONCLAD/deliverable-turnaround-plan-v4-FINAL.pptx",
        "IRONCLAD Turnaround Plan",
        "v4 FINAL",
        [("Updated plan",["Expanded scope", "Success fee narrative diverges"]), ("Actions",["Owner list", "Cadence"])],
    )


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
