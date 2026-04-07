from __future__ import annotations

import json
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent


def w(path: str, content: str):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.rstrip() + "\n", encoding="utf-8")


def wjson(path: str, obj):
    p = BASE / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def make_markdown():
    w(
        "mission-vision-values.md",
        """# QuillStack Systems — Mission, Vision, Values

## Mission
Help operators run workflows with clarity and accountability.

## Vision
A reporting + workflow platform customers trust when the data is messy.

## Values
- Definitions matter
- Own reliability
- Build artifacts that survive
""",
    )

    w(
        "brand-guide.md",
        """# QuillStack Systems — Brand Guide (v0.1)

Palette:
- Navy: #0B1220
- Violet: #7C3AED
- Off-white: #F6F7FB

Voice:
- Calm, precise, no hype.
""",
    )

    w(
        "high-level-strategy-2026.md",
        """# High-Level Strategy — 2026

Owner: Chloe Reyes
Last updated: 2026-02-01

Priorities:
- Ship Fable with stable queue behavior
- Make KPI definitions explicit
- Reduce PRD sprawl
""",
    )

    w(
        "product-overview.md",
        """# Product Overview

QuillStack ingests events + docs, normalizes entities, and powers workflows + dashboards.

Feature Fable: definition registry + reconciliation UI.
""",
    )

    w(
        "prd-fable-v1.md",
        """# PRD — Fable (v1)

Date: 2025-12-02
Status: draft

Goal: Create a definition registry and a reconciliation workflow.
""",
    )

    w(
        "prd-fable-v4-FINAL.md",
        """# PRD — Fable (v4 FINAL)

Date: 2026-02-10
Status: FINAL

Ship target: Roadmap deck says May; sprint log suggests June.
""",
    )

    w(
        "release-notes-2025-q4.md",
        """# Release Notes — 2025 Q4

- New dashboard widgets
- Partial queue improvements
- Early Fable groundwork
""",
    )

    w(
        "runbook-queue-saturation.md",
        """# Runbook — Queue Saturation

Symptoms:
- Ingestion latency increases
- Workflows delayed

Actions:
- Scale consumers
- Apply rate limits
- Drain backlog
""",
    )

    w(
        "incident-2026-02-11-postmortem.md",
        """# Incident Postmortem — 2026-02-11

Narrative: third-party API latency caused backlog.

Action: add queue-depth alerts.
""",
    )

    w(
        "exec-notes-2026-02-01.md",
        """# Exec Notes — 2026-02-01

- Fable scope creep
- Roadmap vs sprint reality diverging
""",
    )

    w(
        "product-review-notes-2026-02-18.md",
        """# Product Review — 2026-02-18

- Roadmap deck says Fable May.
- Devin says earliest June.
""",
    )

    w(
        "voice-note-transcript-2026-02-20.md",
        """# Voice Note Transcript — 2026-02-20 (Morgan)

- CS says "active accounts" is engagement.
- Finance wants "billable accounts".
- We need one KPI definition sheet.
""",
    )

    w(
        "incident-postmortem-template.md",
        """# Incident Postmortem Template

Summary:
Impact:
Timeline:
Root cause:
Actions:
""",
    )

    w(
        "prd-template.md",
        """# PRD Template

Problem:
Users:
Goals:
Non-goals:
Risks:
Metrics:
""",
    )


def make_json():
    wjson(
        "org-chart.json",
        {
            "generated": str(date.today()),
            "leaders": {
                "CEO": "Paige Turner",
                "CTO": "Arman Liu",
                "Product": "Chloe Reyes",
                "Eng": "Devin Shaw",
                "Security": "Nia Brooks",
                "Finance": "Morgan Patel",
            },
        },
    )

    wjson(
        "kanban-snapshot-2026-02.json",
        {
            "snapshot_date": "2026-02-19",
            "board": "Platform",
            "columns": {
                "Backlog": ["FABLE-101", "FABLE-120"],
                "In Progress": ["QUEUE-55"],
                "Review": ["FABLE-132"],
                "Done": ["REL-2025Q4-8"],
            },
        },
    )

    wjson(
        "incident-2026-02-11-metrics.json",
        {
            "incident_id": "INC-2026-02-11",
            "signals": {
                "third_party_latency_ms_p95": 180,
                "queue_depth_max": 420000,
                "consumer_lag_seconds_max": 3600,
            },
            "notes": "Metrics show queue depth/consumer lag as dominant driver; third-party latency secondary.",
        },
    )

    wjson(
        "top-customers.json",
        {
            "generated": str(date.today()),
            "customers": [
                {"account": "Harborline Ops", "segment": "logistics"},
                {"account": "JuniperPay", "segment": "fintech"},
            ],
        },
    )

    wjson(
        "kpi-definitions.json",
        {
            "version": "2026-02-20",
            "kpis": [
                {"name": "Active accounts (CS)", "definition": "accounts with >N weekly actions"},
                {"name": "Billable accounts (Finance)", "definition": "accounts with active subscription"},
            ],
        },
    )

    wjson(
        "customer-intake-form.json",
        {
            "schema_version": "0.1",
            "title": "Customer Intake",
            "fields": [
                {"id": "legal_name", "type": "string", "required": True},
                {"id": "industry", "type": "string", "required": True},
                {"id": "kpis", "type": "array", "items": {"type": "string"}},
            ],
        },
    )


def make_pdfs():
    # none required for this case
    return


def make_xlsx():
    import openpyxl

    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Quarterly"
    ws.append(["Year","Quarter","Revenue","Notes"])
    for r in [(2024,"Q4",7200000,""),(2025,"Q1",7800000,""),(2025,"Q2",8300000,""),(2025,"Q3",8600000,""),(2025,"Q4",9100000,"")]: ws.append(list(r))
    for i in range(2, ws.max_row+1): ws.cell(i,3).number_format="$#,##0"
    wb.save(BASE/"quarterly-revenue-history-2024-2025.xlsx")

    wb2=openpyxl.Workbook(); ws2=wb2.active; ws2.title="ARR"
    ws2.append(["Customer","MRR","ARR"])
    for r in [("JuniperPay",310000,None),("Harborline Ops",260000,None)]: ws2.append(list(r))
    for i in range(2, ws2.max_row+1):
        ws2.cell(i,3).value=f"=B{i}*12"
        ws2.cell(i,2).number_format="$#,##0"; ws2.cell(i,3).number_format="$#,##0"
    wb2.save(BASE/"arr-model-2026.xlsx")

    wb3=openpyxl.Workbook(); ws3=wb3.active; ws3.title="Sprints Q1 2026"
    ws3.append(["Sprint","Start","End","Committed","Completed","Carryover","Notes"])
    rows=[("S-01","2026-01-05","2026-01-16",120,102,None,""),("S-02","2026-01-19","2026-01-30",130,95,None,"queue work"),("S-03","2026-02-02","2026-02-13",140,88,None,"Fable slip")]
    for r in rows: ws3.append(list(r))
    for i in range(2, ws3.max_row+1): ws3.cell(i,6).value=f"=D{i}-E{i}"
    wb3.save(BASE/"sprint-log-2026-q1.xlsx")

    wb4=openpyxl.Workbook(); ws4=wb4.active; ws4.title="Pipeline Q2 2026"
    ws4.append(["Account","Stage","Close","Amount","Probability","Weighted"])
    for a,st,cl,amt,pr in [("CedarOps","Discovery","2026-05-30",1800000,0.35),("Lakefront Labs","Negotiation","2026-06-20",2400000,0.55)]:
        ws4.append([a,st,cl,amt,pr,None])
    for i in range(2, ws4.max_row+1):
        ws4.cell(i,6).value=f"=D{i}*E{i}"
        ws4.cell(i,4).number_format="$#,##0"; ws4.cell(i,6).number_format="$#,##0"; ws4.cell(i,5).number_format="0%"
    wb4.save(BASE/"pipeline-2026-q2.xlsx")

    wb5=openpyxl.Workbook(); ws5=wb5.active; ws5.title="Renewals"
    ws5.append(["Customer","Renewal","ARR","Health","Notes"])
    ws5.append(["JuniperPay","2026-09-01",3720000,"Yellow","definition debate" ])
    ws5.cell(2,3).number_format="$#,##0"
    wb5.save(BASE/"renewal-risk.xlsx")


def make_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    NAVY=RGBColor(0x0B,0x12,0x20)
    VIO=RGBColor(0x7C,0x3A,0xED)
    OFF=RGBColor(0xF6,0xF7,0xFB)

    def header(prs, slide):
        f=slide.background.fill; f.solid(); f.fore_color.rgb=OFF
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
        acc=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), prs.slide_width, Inches(0.06))
        acc.fill.solid(); acc.fill.fore_color.rgb=VIO; acc.line.fill.background()

    def simple_deck(path: Path, title: str):
        prs=Presentation()
        s=prs.slides.add_slide(prs.slide_layouts[6]); header(prs,s)
        tb=s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(1.0))
        p=tb.text_frame.paragraphs[0]; p.text=title
        p.font.name="Georgia"; p.font.size=Pt(44); p.font.bold=True
        prs.save(path)

    simple_deck(BASE/"roadmap-2026-h1.pptx", "Roadmap — H1 2026")
    simple_deck(BASE/"sales-deck.pptx", "QuillStack Systems")
    simple_deck(BASE/"cs-qbr-template.pptx", "QBR Template")


def main():
    make_markdown(); make_json(); make_pdfs(); make_xlsx(); make_pptx()


if __name__ == "__main__":
    main()
